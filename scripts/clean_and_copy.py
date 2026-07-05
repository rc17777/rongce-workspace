import sys, os, shutil
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = r'D:\openclaw-workspace\output\v5_expanded.pptx'
DST = r'D:\openclaw-workspace\output\v6_clean.pptx'
DESK = r'C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v6_模板.pptx'

prs = Presentation(SRC)
removed = 0
for si, slide in enumerate(prs.slides):
    shapes_to_remove = []
    for shape in slide.shapes:
        is_ph = hasattr(shape, 'is_placeholder') and shape.is_placeholder
        has_text = shape.has_text_frame
        if has_text:
            text = shape.text_frame.text.strip()
            is_empty = not text or text in [
                '点击此处添加标题', '点击此处添加文本', 
                '单击此处添加标题', '单击此处添加文本',
                'Click to add title', 'Click to add text'
            ]
        else:
            is_empty = False
        if is_ph and has_text and is_empty:
            shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        removed += 1

print(f'Removed {removed} placeholder shapes')
prs.save(DST)

# Copy to desktop
for i in range(1, 20):
    if i == 1:
        target = DESK
    else:
        base, ext = os.path.splitext(DESK)
        target = f'{base}_{i}{ext}'
    try:
        if os.path.exists(target):
            os.remove(target)
        shutil.copy2(DST, target)
        print(f'Saved to: {target}')
        break
    except Exception as e:
        continue
