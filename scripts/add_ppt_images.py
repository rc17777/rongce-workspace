#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""给第24/25/27/28/29页右侧嵌入配图"""
import sys,os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Cm, Emu

prs = Presentation(r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx')

img_dir = r'D:\openclaw-workspace\temp\ppt_images'

# 配置：{页码索引(0-based): ('图片文件', 左上角位置, 宽, 高)}
config = {
    23: ('slide_24.png', Cm(23), Cm(4.5), Cm(9), Cm(12)),
    24: ('slide_25.png', Cm(23), Cm(4.5), Cm(9), Cm(12)),
    26: ('slide_27.png', Cm(23), Cm(4.5), Cm(9), Cm(12)),
    27: ('slide_28.png', Cm(23), Cm(4.5), Cm(9), Cm(12)),
    28: ('slide_29.png', Cm(23), Cm(4.5), Cm(9), Cm(12)),
}

for slide_idx, (img_name, left, top, width, height) in config.items():
    slide = prs.slides[slide_idx]
    img_path = os.path.join(img_dir, img_name)
    if os.path.exists(img_path):
        # 在右侧添加图片
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        print(f'Slide {slide_idx+1}: added {img_name} ({os.path.getsize(img_path)//1024}KB)')
    else:
        print(f'Slide {slide_idx+1}: img not found: {img_path}')

# 保存
prs.save(r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx')
print(f'\nDone. Size: {os.path.getsize(r"D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx")//1024}KB')
