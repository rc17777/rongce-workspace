"""
模板迁移脚本 v4：逐形状重建，保留原版式结构
"""

import os, copy, io, time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

TEMPLATE = r"C:\Users\scrccpa\Desktop\物资管理制度培训-2025.1.20.pptx"
CONTENT  = r"C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v2.pptx"

tpl = Presentation(TEMPLATE)
src = Presentation(CONTENT)
print("Template: %d slides, Content: %d slides" % (len(tpl.slides), len(src.slides)))

# ── 查找layout ──
def find_layout(prs, name):
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name == name:
            return i
    return 0

out = Presentation(TEMPLATE)
while len(out.slides) > 0:
    rId = out.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    out.part.drop_rel(rId)
    out.slides._sldIdLst.remove(out.slides._sldIdLst[0])

L_TITLE   = find_layout(out, '标题幻灯片')
L_SECTION = find_layout(out, '节标题')
L_CONTENT = find_layout(out, '标题和内容')

# ── 颜色映射：v2深色 → 模板浅色 ──
CLR_MAP = {
    'FFFFFF': '0B75B4',  # 白→蓝
    '1E2761': '0B75B4',  # v2深蓝→模板蓝
    '333333': '333333',  # 保持
    '666666': '666666',  # 保持
    '8899BB': '666666',  # 浅蓝灰→灰色
    'CADCFC': '0B75B4',  # 浅蓝→模板蓝
    'F5F5F5': 'F0F0F0',  # 浅灰→稍浅灰
    'F96167': 'C0392B',  # 红→暗红
    'B85042': 'C0392B',  # 棕红→暗红
    'D4A843': '0B75B4',  # 金→蓝
    '2C5F2D': '0B75B4',  # 绿→蓝
    '1E2761': '0B75B4',  # 重复
}

def map_color(hex_color):
    """映射颜色：v2的深色系→模板的浅色蓝色系"""
    if hex_color is None:
        return RGBColor(0x33, 0x33, 0x33)
    # 处理RGBColor对象
    if hasattr(hex_color, '__str__'):
        s = str(hex_color)
    else:
        s = str(hex_color)
    s = s.upper().replace('#', '')
    mapped = CLR_MAP.get(s, None)
    if mapped:
        return RGBColor(int(mapped[0:2], 16), int(mapped[2:4], 16), int(mapped[4:6], 16))
    # 如果原色不够深（太浅），加深它
    try:
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
        brightness = (r+g+b)/3
        if brightness < 128:  # 暗色→映射到模板蓝
            return RGBColor(0x0B, 0x75, 0xB4)
        elif brightness > 230:  # 很浅→白色
            return RGBColor(0x33, 0x33, 0x33)
        return RGBColor(0x33, 0x33, 0x33)
    except:
        return RGBColor(0x33, 0x33, 0x33)

def is_empty_shape(shape):
    """判断形状是否为空白装饰"""
    if not shape.has_text_frame:
        return False
    txt = shape.text_frame.text.strip()
    return len(txt) == 0

def get_shape_fill_color(shape):
    """获取形状填充色（如果有）"""
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc.type is not None:
                return str(fc.rgb)
    except:
        pass
    return None

def add_bg_images(slide):
    """从模板第3页复制背景装饰图"""
    proto = tpl.slides[2]
    for shape in proto.shapes:
        if shape.shape_type == 13:  # Picture
            try:
                image_blob = shape.image.blob
                pic = slide.shapes.add_picture(
                    io.BytesIO(image_blob), shape.left, shape.top,
                    shape.width, shape.height)
                if shape.rotation != 0:
                    pic.rotation = shape.rotation
            except:
                pass

def add_text_shape(slide, text, left, top, width, height,
                   font_name='微软雅黑', font_size=14, bold=False,
                   color=RGBColor(0x33,0x33,0x33), align=PP_ALIGN.LEFT,
                   valign='center', bg_color=None):
    """添加一个匹配位置的文本框"""
    # 如果有背景色，先加一个矩形
    if bg_color:
        rect = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE = 1
            left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_color
        rect.line.fill.background()
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    # 垂直对齐
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(3)
    
    return txBox

def clone_shape(src_shape, dst_slide, scale_x=1.0, scale_y=1.0, offset_x=0, offset_y=0):
    """将源形状的内容复制到目标幻灯片上，确保浅底可见"""
    if not src_shape.has_text_frame:
        return None
    
    text = src_shape.text_frame.text.strip()
    if not text:
        return None
    
    # 位置
    left = int(src_shape.left * scale_x) + offset_x
    top = int(src_shape.top * scale_y) + offset_y
    width = int(src_shape.width * scale_x)
    height = int(src_shape.height * scale_y)
    
    # 字体：统一使用微软雅黑
    font_name = '微软雅黑'
    font_size = 14
    bold = False
    align = PP_ALIGN.LEFT
    
    # 从源形状提取字体大小和加粗
    try:
        for para in src_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_size = run.font.size / 12700
                if run.font.bold is not None:
                    bold = run.font.bold
                break
            break
    except:
        pass
    
    # 对齐
    try:
        a = src_shape.text_frame.paragraphs[0].alignment
        if a is not None:
            align = a
    except:
        pass
    
    # 颜色：确保浅底可读
    orig_color = None
    try:
        for para in src_shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    c = run.font.color.rgb
                    if c:
                        orig_color = str(c)
                except:
                    pass
                break
            break
    except:
        pass
    
    # 决策文本颜色：浅色→#333，深色→保持或#0B75B4
    if orig_color:
        try:
            r = int(orig_color[0:2], 16)
            g = int(orig_color[2:4], 16)
            b = int(orig_color[4:6], 16)
            brightness = (r + g + b) / 3.0
            if brightness > 180:
                text_color = RGBColor(0x33, 0x33, 0x33)  # 浅色→深灰
            elif brightness < 80:
                text_color = RGBColor(r, g, b)  # 深色保持
            else:
                text_color = RGBColor(0x33, 0x33, 0x33)  # 中间色→深灰
        except:
            text_color = RGBColor(0x33, 0x33, 0x33)
    else:
        text_color = RGBColor(0x33, 0x33, 0x33)
    
    # 背景色：模板不用色块背景
    bg_color = None
    
    return add_text_shape(dst_slide, text, left, top, width, height,
                         font_name, font_size, bold, text_color, align,
                         bg_color=None)

# ── 页类型映射 ──
PAGE_MAP = {
    1: 'title', 2: 'agenda', 3: 'section',
    4: 'content', 5: 'content', 6: 'content', 7: 'content',
    8: 'content', 9: 'content', 10: 'content',
    11: 'section', 12: 'content', 13: 'content', 14: 'content',
    15: 'content', 16: 'content', 17: 'content',
    18: 'section', 19: 'content', 20: 'content',
    21: 'content', 22: 'content', 23: 'content', 24: 'end',
}

# ── 逐页逐形状重建 ──
for src_idx in range(len(src.slides)):
    page_num = src_idx + 1
    ptype = PAGE_MAP.get(page_num, 'content')
    src_slide = src.slides[src_idx]
    
    # 选layout
    if ptype in ('title', 'end'):
        layout_idx = L_TITLE
    elif ptype == 'section':
        layout_idx = L_SECTION
    else:
        layout_idx = L_CONTENT
    new_slide = out.slides.add_slide(out.slide_layouts[layout_idx])
    
    # 背景图
    if ptype not in ('title', 'end'):
        add_bg_images(new_slide)
    
    if ptype in ('title', 'section', 'end'):
        # 这几类特殊处理
        if ptype == 'title':
            # 封面：覆盖原内容，用自定义排版
            for shape in src_slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    pass  # 跳过，使用固定封面
            # 自定义封面文字
            from pptx.util import Cm
            add_text_shape(new_slide, "提升公司全员\n审计风险意识",
                          Cm(4), Cm(5), Cm(26), Cm(5),
                          '楷体', 40, True, RGBColor(0x0B,0x75,0xB4), PP_ALIGN.CENTER)
            add_text_shape(new_slide, "四川轨道公司专题培训",
                          Cm(4), Cm(11), Cm(26), Cm(2),
                          '楷体', 22, False, RGBColor(0x66,0x66,0x66), PP_ALIGN.CENTER)
            add_text_shape(new_slide, "四川融策会计师事务所    2026年6月",
                          Cm(4), Cm(13.5), Cm(26), Cm(1.5),
                          '楷体', 16, False, RGBColor(0x66,0x66,0x66), PP_ALIGN.CENTER)
        
        elif ptype == 'section':
            all_text = []
            for shape in src_slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        all_text.append(t)
            main_text = all_text[0] if all_text else ""
            sub_text = all_text[1] if len(all_text) > 1 else ""
            from pptx.util import Cm
            add_text_shape(new_slide, main_text,
                          Cm(3), Cm(6), Cm(28), Cm(5),
                          '方正小标宋简体', 48, True, RGBColor(0x0B,0x75,0xB4), PP_ALIGN.CENTER)
            if sub_text:
                add_text_shape(new_slide, sub_text,
                              Cm(4), Cm(11.5), Cm(26), Cm(3),
                              '微软雅黑', 20, False, RGBColor(0x66,0x66,0x66), PP_ALIGN.CENTER)
        
        elif ptype == 'end':
            from pptx.util import Cm
            add_text_shape(new_slide, "谢谢！",
                          Cm(4), Cm(6), Cm(26), Cm(4),
                          '楷体', 48, True, RGBColor(0x0B,0x75,0xB4), PP_ALIGN.CENTER)
            add_text_shape(new_slide, "欢迎会后交流提问",
                          Cm(4), Cm(11), Cm(26), Cm(2),
                          '楷体', 20, False, RGBColor(0x66,0x66,0x66), PP_ALIGN.CENTER)
    
    else:
        # 内容页：逐形状克隆
        for shape in src_slide.shapes:
            if is_empty_shape(shape):
                continue
            if not shape.has_text_frame:
                continue
            
            clone_shape(shape, new_slide)

# ── 保存 ──
ts = time.strftime("%H%M%S")
out_file = r"D:\openclaw-workspace\output\sc-v4-" + ts + ".pptx"
out.save(out_file)
print("\nDone: " + out_file)
print("Slides: " + str(len(out.slides)))
