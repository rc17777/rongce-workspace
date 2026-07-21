# -*- coding: utf-8 -*-
"""
融策·券商风整报告布局母版 PPTX 生成器
====================================
生成一套完整可编辑的研报版式体系，而不是单张图表：
1. 封面
2. 核心观点摘要
3. 目录
4. 章节开篇页
5. 正文图文分析页
6. 数据表格页
7. 方法论/框架页
8. 结论与行动建议页
9. 免责声明

用法：
    python -X utf8 scripts/generate_securities_report_layout_pptx.py output/融策_券商风整报告布局母版.pptx
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

NAVY = RGBColor(6, 26, 51)
NAVY2 = RGBColor(10, 42, 74)
TEAL = RGBColor(26, 111, 120)
GOLD = RGBColor(184, 138, 68)
GOLD2 = RGBColor(214, 176, 113)
INK = RGBColor(31, 41, 51)
GRAY = RGBColor(83, 97, 109)
MUTED = RGBColor(135, 147, 160)
PANEL = RGBColor(244, 246, 248)
WHITE = RGBColor(255, 255, 255)
GRID = RGBColor(217, 222, 229)
RED = RGBColor(184, 64, 58)
GREEN = RGBColor(43, 122, 85)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def font(run, size=10, color=INK, bold=False):
    run.font.name = 'Microsoft YaHei'
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def textbox(slide, x, y, w, h, text='', size=10, color=INK, bold=False,
            align=PP_ALIGN.LEFT, fill=None, line=None, margin=0.08):
    shape = slide.shapes.add_textbox(x, y, w, h)
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line; shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    font(r, size, color, bold)
    return shape


def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def round_rect(slide, x, y, w, h, fill, line=GRID):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(0.75)
    return s


def header(slide, section_title, page_title, right='融策·审盾研究'):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.52), NAVY)
    rect(slide, Inches(0), Inches(0.52), SLIDE_W, Inches(0.06), GOLD)
    textbox(slide, Inches(0.38), Inches(0.08), Inches(2.6), Inches(0.16), section_title, 8, GOLD2, True)
    textbox(slide, Inches(0.38), Inches(0.25), Inches(8.2), Inches(0.22), page_title, 14, WHITE, True)
    textbox(slide, Inches(9.5), Inches(0.23), Inches(3.4), Inches(0.22), right, 8.5, RGBColor(201,210,220), False, PP_ALIGN.RIGHT)


def footer(slide, page_no, source='融策AI审计中台'):
    textbox(slide, Inches(0.38), Inches(7.08), Inches(5.8), Inches(0.16), f'资料来源：{source}', 7.2, MUTED)
    textbox(slide, Inches(10.8), Inches(7.08), Inches(2.1), Inches(0.16), f'{page_no:02d}', 7.2, MUTED, False, PP_ALIGN.RIGHT)


def kpi(slide, x, y, label, value, color=NAVY):
    round_rect(slide, x, y, Inches(2.35), Inches(0.65), PANEL)
    textbox(slide, x + Inches(0.12), y + Inches(0.08), Inches(2.05), Inches(0.16), label, 7.5, GRAY)
    textbox(slide, x + Inches(0.12), y + Inches(0.30), Inches(2.05), Inches(0.24), value, 15, color, True)


def bullet_box(slide, x, y, w, h, title, bullets, title_fill=NAVY2):
    round_rect(slide, x, y, w, h, RGBColor(250,251,252), GRID)
    rect(slide, x, y, w, Inches(0.38), title_fill)
    textbox(slide, x + Inches(0.14), y + Inches(0.08), w - Inches(0.25), Inches(0.18), title, 9.5, WHITE, True)
    tf = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.55), w - Inches(0.34), h - Inches(0.68)).text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, b in enumerate(bullets, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.space_after = Pt(11)
        r1 = p.add_run(); r1.text = f'{i}. '; font(r1, 8.8, GOLD, True)
        r2 = p.add_run(); r2.text = b; font(r2, 8.8, INK, False)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    rect(slide, Inches(0), Inches(0), Inches(2.0), SLIDE_H, NAVY)
    rect(slide, Inches(2.0), Inches(0), Inches(0.08), SLIDE_H, GOLD)
    textbox(slide, Inches(0.35), Inches(0.45), Inches(1.3), Inches(0.25), '融策', 18, GOLD2, True, PP_ALIGN.CENTER)
    textbox(slide, Inches(0.33), Inches(0.78), Inches(1.35), Inches(0.18), '审盾研究', 9, RGBColor(200,210,220), False, PP_ALIGN.CENTER)
    textbox(slide, Inches(2.65), Inches(0.95), Inches(8.4), Inches(0.3), '行业深度报告 | AI审计中台系列', 12, GOLD, True)
    textbox(slide, Inches(2.62), Inches(1.55), Inches(8.9), Inches(1.15), '融策AI审计中台：\n重塑政府审计的数据穿透力', 31, NAVY, True)
    textbox(slide, Inches(2.68), Inches(2.85), Inches(8.4), Inches(0.38), '从信息化核查到智能交叉验证的0到1跃迁', 17, TEAL)
    rect(slide, Inches(2.65), Inches(3.55), Inches(1.05), Inches(0.05), GOLD)
    textbox(slide, Inches(2.65), Inches(3.9), Inches(7.4), Inches(1.55),
            '本报告复刻券商深度研报的结构逻辑，以“观点先行、图表作证、判断落地”为核心，形成适用于绩效评价、经责审计、专项资金、工程决算等业务线的标准化研报母版。', 12, INK)
    textbox(slide, Inches(2.65), Inches(6.35), Inches(4.3), Inches(0.22), '分析师：融策右护卫 | 报告日期：2026年7月21日', 9, GRAY)
    textbox(slide, Inches(9.6), Inches(6.35), Inches(2.7), Inches(0.22), '内部研报母版 v1.0', 9, GRAY, False, PP_ALIGN.RIGHT)
    return slide


def executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, '核心摘要', '三句话讲清本报告的结论')
    cards = [
        ('01', '不是做“漂亮图表”，而是做证据链页面', '每页必须有一个明确判断，并用KPI、图表、右侧解释和来源共同支撑。'),
        ('02', '研报模板应服务融策业务线，而非照搬券商金融话术', '绩效评价、经责审计、工程决算、专项资金各自需要不同的数据口径与图表母版。'),
        ('03', '可编辑PPTX用于汇报打磨，Word用于正式归档出具', 'PPTX承载图文协同和客户沟通，Word承载正式报告文本、附件和签批流程。'),
    ]
    for i, (no, title, body) in enumerate(cards):
        y = Inches(1.05 + i * 1.55)
        round_rect(slide, Inches(0.72), y, Inches(11.9), Inches(1.22), PANEL)
        textbox(slide, Inches(0.96), y + Inches(0.20), Inches(0.58), Inches(0.35), no, 20, GOLD, True)
        textbox(slide, Inches(1.75), y + Inches(0.18), Inches(5.8), Inches(0.28), title, 13, NAVY, True)
        textbox(slide, Inches(1.75), y + Inches(0.58), Inches(9.8), Inches(0.32), body, 10, INK)
    footer(slide, 2)


def toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, '目录', '报告结构总览')
    items = [
        ('01', '行业痛点：传统审计报告为何不够像研报'),
        ('02', '模板逻辑：观点-证据-判断-行动闭环'),
        ('03', '图表体系：趋势、结构、对比、矩阵四类母版'),
        ('04', '业务适配：绩效评价/经责/工程/专项资金'),
        ('05', '实施路径：从样板报告到公司级模板库'),
    ]
    for i, (num, text) in enumerate(items):
        y = Inches(1.2 + i * 0.82)
        textbox(slide, Inches(1.25), y, Inches(0.65), Inches(0.32), num, 17, GOLD, True)
        rect(slide, Inches(2.05), y + Inches(0.18), Inches(0.45), Inches(0.02), GRID)
        textbox(slide, Inches(2.65), y, Inches(7.5), Inches(0.32), text, 14, INK, True)
        textbox(slide, Inches(11.25), y, Inches(0.4), Inches(0.28), f'{i+4}', 10, MUTED, False, PP_ALIGN.RIGHT)
    footer(slide, 3)


def section_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.12), GOLD)
    textbox(slide, Inches(0.8), Inches(1.05), Inches(1.2), Inches(0.35), '01', 24, GOLD2, True)
    textbox(slide, Inches(0.78), Inches(1.75), Inches(8.8), Inches(0.72), '行业痛点', 34, WHITE, True)
    textbox(slide, Inches(0.82), Inches(2.65), Inches(8.6), Inches(0.42), '传统审计报告为何“不像研报”', 20, RGBColor(200,210,220))
    textbox(slide, Inches(0.85), Inches(4.15), Inches(7.8), Inches(0.9), '报告不是资料堆砌，而是围绕判断组织证据。\n真正的研报版式，先让读者相信结论，再让读者看到证据。', 14, RGBColor(220,226,232))
    textbox(slide, Inches(11.1), Inches(6.7), Inches(1.4), Inches(0.22), '04', 9, RGBColor(190,200,210), False, PP_ALIGN.RIGHT)


def chart_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, '图文分析', '财政收入修复斜率放缓，税收收入弹性仍是关键变量')
    kpi(slide, Inches(0.55), Inches(0.9), '2025E预算收入', '22.8万亿', NAVY)
    kpi(slide, Inches(3.02), Inches(0.9), '2024-2025E增量', '+0.8万亿', GOLD)
    kpi(slide, Inches(5.49), Inches(0.9), '税收占比', '83.3%', TEAL)
    textbox(slide, Inches(0.55), Inches(1.82), Inches(6.8), Inches(0.26), '图1：全国一般公共预算收入与税收收入走势（万亿元）', 11.5, INK, True)
    data = CategoryChartData()
    data.categories = ['2020', '2021', '2022', '2023', '2024', '2025E']
    data.add_series('一般公共预算收入', (18.29, 20.25, 20.37, 21.68, 22.00, 22.80))
    data.add_series('税收收入', (15.43, 17.27, 16.66, 18.11, 18.20, 19.00))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.55), Inches(2.22), Inches(7.35), Inches(4.02), data).chart
    chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.TOP; chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = GRID
    chart.font.name = 'Microsoft YaHei'; chart.font.size = Pt(8)
    chart.series[0].format.line.color.rgb = NAVY; chart.series[0].format.line.width = Pt(2.25)
    chart.series[1].format.line.color.rgb = GOLD; chart.series[1].format.line.width = Pt(2.25)
    bullet_box(slide, Inches(8.35), Inches(1.82), Inches(4.25), Inches(4.55), '核心判断', [
        '财政收入修复并非线性扩张，税基质量和房地产链条回暖仍决定后续弹性。',
        '税收收入占比维持高位，说明非税收入拉动空间有限，应关注税源真实性。',
        '若预算收入增速明显高于经济增速，应回查一次性收入、非税缴库和跨期调节。',
    ])
    footer(slide, 5, '财政部，Wind，融策会计师事务所整理')


def table_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, '数据表格', '项目结构决定审计资源配置')
    textbox(slide, Inches(0.55), Inches(0.95), Inches(7.4), Inches(0.38), '表1：不同业务线的模板化程度与AI复核适配度', 14, NAVY, True)
    table = slide.shapes.add_table(6, 5, Inches(0.55), Inches(1.55), Inches(7.75), Inches(3.45)).table
    headers = ['业务线', '数据结构化', '报告标准化', 'AI适配度', '优先级']
    rows = [
        ['绩效评价', '中', '高', '高', 'P0'],
        ['工程决算', '高', '中', '高', 'P0'],
        ['经责审计', '中', '中', '中高', 'P1'],
        ['专项资金', '中高', '中', '高', 'P1'],
        ['资产清查', '高', '高', '中', 'P2'],
    ]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs: font(r, 8.5, WHITE, True)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j); cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = PANEL if i % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs: font(r, 8.5, INK, False)
    bullet_box(slide, Inches(8.65), Inches(1.55), Inches(3.85), Inches(3.45), '解读', [
        '优先从标准化程度高、复核规则清晰的业务线切入。',
        '模板不是装饰，而是把数据口径和报告口径固化。',
        '每类业务线应沉淀一套图表库和一套风险判断句库。',
    ])
    footer(slide, 6, '融策项目台账，内部复盘')


def methodology_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, '方法论', '一页研报的标准生产逻辑')
    steps = [
        ('观点', '一句话判断', '先给结论，不绕弯'),
        ('指标', '3个关键数字', '用数字压住场'),
        ('图表', '一张主证据图', '趋势/结构/对比/矩阵'),
        ('解释', '3条核心判断', '说明为什么重要'),
        ('行动', '下一步建议', '落到项目管理动作'),
    ]
    for i, (title, subtitle, desc) in enumerate(steps):
        x = Inches(0.65 + i * 2.45)
        round_rect(slide, x, Inches(1.65), Inches(1.75), Inches(3.6), PANEL)
        textbox(slide, x + Inches(0.18), Inches(1.95), Inches(1.38), Inches(0.32), title, 18, NAVY, True, PP_ALIGN.CENTER)
        rect(slide, x + Inches(0.42), Inches(2.55), Inches(0.92), Inches(0.04), GOLD)
        textbox(slide, x + Inches(0.15), Inches(2.85), Inches(1.45), Inches(0.28), subtitle, 10, TEAL, True, PP_ALIGN.CENTER)
        textbox(slide, x + Inches(0.18), Inches(3.32), Inches(1.38), Inches(0.7), desc, 9, GRAY, False, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            textbox(slide, x + Inches(1.9), Inches(3.05), Inches(0.35), Inches(0.3), '→', 20, GOLD, True, PP_ALIGN.CENTER)
    footer(slide, 7, '融策研报模板方法论')


def conclusion_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, '结论', '从样板报告到公司级研报模板库')
    bullet_box(slide, Inches(0.7), Inches(1.1), Inches(3.7), Inches(4.8), '短期动作', [
        '先选绩效评价报告作为样板，做1套完整研报版式。',
        '固定4类图表：趋势、结构、对比、风险矩阵。',
        '每张图绑定一条正文判断和一条资料来源。',
    ])
    bullet_box(slide, Inches(4.8), Inches(1.1), Inches(3.7), Inches(4.8), '中期建设', [
        '建立业务线模板库：绩效/经责/工程/专项资金。',
        '形成数据口径、图表样式、判断句库三件套。',
        '把AI复核嵌入模板交付流程。',
    ], TEAL)
    bullet_box(slide, Inches(8.9), Inches(1.1), Inches(3.7), Inches(4.8), '长期价值', [
        '形成融策自己的“研究型审计报告”品牌。',
        '让客户感知从审计服务升级到决策支持。',
        '为审盾产品化和投标展示提供样板资产。',
    ], GOLD)
    footer(slide, 8, '融策AI审计中台建设路线')


def disclaimer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, '免责声明', '报告使用边界')
    textbox(slide, Inches(0.75), Inches(1.25), Inches(11.6), Inches(1.15),
            '本报告为四川融策会计师事务所及四川融策工程咨询公司内部研报模板示例，所列数据与案例仅用于展示图文协同版式，不构成正式审计结论、投资建议或对外承诺。正式项目报告应以经复核的底稿、取证材料、法规依据和签批流程为准。',
            11, INK)
    textbox(slide, Inches(0.75), Inches(3.05), Inches(11.6), Inches(0.32), '使用要求', 15, NAVY, True)
    textbox(slide, Inches(0.75), Inches(3.55), Inches(10.9), Inches(1.35),
            '1. 图表数据必须注明来源；\n2. 重大金额、比率和结论必须与底稿交叉核验；\n3. AI生成内容必须经项目经理和质控复核后方可用于正式报告；\n4. 对外版本应删除内部方法、模型路由和未公开数据来源。',
            10.5, GRAY)
    footer(slide, 9, '融策内部模板')


def build(output):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    cover(prs)
    executive_summary(prs)
    toc(prs)
    section_page(prs)
    chart_page(prs)
    table_page(prs)
    methodology_page(prs)
    conclusion_page(prs)
    disclaimer(prs)
    Path(output).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f'✅ 整报告布局母版已生成: {output}')


if __name__ == '__main__':
    output = sys.argv[1] if len(sys.argv) > 1 else 'output/融策_券商风整报告布局母版_可编辑.pptx'
    build(output)
