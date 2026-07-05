#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
融策标书终极组合生成器
整合：DeepSeek内容 + deepseek-charting图表 + 火山引擎图片 + draw.io流程图 + PPT整合
作者：融策右护卫
版本：1.0.0
"""

import os
import sys

# Windows 编码修复
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

import json
import time
import argparse
from datetime import datetime
from pathlib import Path

# ==================== 配置区 ====================
CONFIG = {
    "project_name": "",
    "project_type": "绩效评价",  # 绩效评价/工程审计/经责审计/ etc.
    "output_dir": "./output",
    "template_pptx": "./templates/融策标书模板.pptx",
    "brand_colors": {
        "dark_blue": "#0A1F3F",
        "teal": "#1A5C6E",
        "gold": "#C5955C",
        "warm_gray": "#F5F2EC",
        "white": "#FFFFFF"
    },
    # API Keys (从环境变量读取)
    "deepseek_api_key": os.environ.get("DEEPSEEK_API_KEY", ""),
    "seedream_api_key": os.environ.get("SEEDREAM_API_KEY", ""),
    "seedream_base_url": os.environ.get("SEEDREAM_BASE_URL", "https://ark.cn-beijing.volces.com/api/v3/images/generations"),
}

# ==================== 步骤1：DeepSeek生成内容 ====================
def step1_generate_content(project_name, project_type):
    """
    调用 DeepSeek API 生成标书内容大纲
    实际使用时替换为真实 API 调用
    """
    print(f"\n{'='*60}")
    print("步骤1：DeepSeek 生成内容大纲")
    print(f"{'='*60}")
    
    prompt = f"""请为"{project_name}"项目生成一份完整的投标方案大纲。
项目类型：{project_type}

要求包含以下章节：
1. 项目理解与服务方案
2. 技术方案与方法论
3. 项目组织架构与人员配置
4. 质量保证措施
5. 进度计划与里程碑
6. 风险应对措施

请输出结构化的 Markdown 格式，每节包含3-5个要点。"""
    
    # 模拟输出（实际使用时调用 DeepSeek API）
    content = f"""# {project_name} 投标方案

## 1. 项目理解与服务方案
- 深入理解项目背景与需求
- 明确审计目标与范围
- 制定针对性服务策略

## 2. 技术方案与方法论
- 采用风险导向审计方法
- 运用数据分析技术
- 建立多维度评价指标体系

## 3. 项目组织架构
- 项目经理负责制
- 专业分工协作机制
- 质量控制三级复核

## 4. 质量保证措施
- 全过程质量监控
- 底稿三级复核制度
- 重大问题会商机制

## 5. 进度计划
- 准备阶段（2周）
- 实施阶段（6周）
- 报告阶段（2周）

## 6. 风险应对
- 数据获取风险
- 人员变动风险
- 进度延误风险
"""
    
    output_file = f"{CONFIG['output_dir']}/01_内容大纲.md"
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(content)
    
    print(f"✅ 内容大纲已生成: {output_file}")
    return content

# ==================== 步骤2：deepseek-charting 生成图表 ====================
def step2_generate_charts():
    """
    生成 Mermaid 流程图和 ECharts 数据图表
    """
    print(f"\n{'='*60}")
    print("步骤2：deepseek-charting 生成图表")
    print(f"{'='*60}")
    
    # 2.1 Mermaid 流程图
    mermaid_code = """graph TD
    A[项目立项] --> B[组建审计组]
    B --> C[制定方案]
    C --> D[进场实施]
    D --> E[数据分析]
    E --> F[编制报告]
    F --> G[复核提交]
    
    style A fill:#0A1F3F,color:#fff
    style B fill:#1A5C6E,color:#fff
    style C fill:#1A5C6E,color:#fff
    style D fill:#C5955C,color:#fff
    style E fill:#C5955C,color:#fff
    style F fill:#1A5C6E,color:#fff
    style G fill:#0A1F3F,color:#fff
"""
    
    mermaid_file = f"{CONFIG['output_dir']}/02_流程图.mmd"
    with open(mermaid_file, 'w', encoding='utf-8') as f:
        f.write(mermaid_code)
    print(f"✅ Mermaid 流程图已生成: {mermaid_file}")
    print("   提示：复制代码到 https://mermaid.live 渲染，或导出 PNG")
    
    # 2.2 ECharts 甘特图 HTML
    echarts_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>项目进度甘特图</title>
    <script src="https://cdn.jsdelivr.net/npm/echarts@5.5.0/dist/echarts.min.js"></script>
    <style>
        body {{ margin: 0; padding: 20px; font-family: "微软雅黑", sans-serif; background: #f5f5f5; }}
        .container {{ max-width: 1200px; margin: 0 auto; background: white; padding: 20px; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }}
        h2 {{ color: {CONFIG['brand_colors']['dark_blue']}; border-bottom: 2px solid {CONFIG['brand_colors']['gold']}; padding-bottom: 10px; }}
        #chart {{ width: 100%; height: 500px; }}
    </style>
</head>
<body>
    <div class="container">
        <h2>项目进度计划</h2>
        <div id="chart"></div>
    </div>
    <script>
        var chart = echarts.init(document.getElementById('chart'));
        var data = [
            {{ name: '项目进场', value: [0, 7], itemStyle: {{ color: '{CONFIG['brand_colors']['dark_blue']}' }} }},
            {{ name: '资料收集', value: [7, 21], itemStyle: {{ color: '{CONFIG['brand_colors']['teal']}' }} }},
            {{ name: '现场调研', value: [21, 35], itemStyle: {{ color: '{CONFIG['brand_colors']['gold']}' }} }},
            {{ name: '数据分析', value: [28, 42], itemStyle: {{ color: '{CONFIG['brand_colors']['dark_blue']}' }} }},
            {{ name: '报告撰写', value: [42, 56], itemStyle: {{ color: '{CONFIG['brand_colors']['teal']}' }} }},
            {{ name: '专家复核', value: [56, 63], itemStyle: {{ color: '{CONFIG['brand_colors']['gold']}' }} }},
            {{ name: '成果提交', value: [63, 70], itemStyle: {{ color: '{CONFIG['brand_colors']['dark_blue']}' }} }}
        ];
        
        var option = {{
            tooltip: {{
                formatter: function(params) {{
                    return params.marker + params.name + '<br/>第 ' + (params.value[0]+1) + ' 天 - 第 ' + params.value[1] + ' 天';
                }}
            }},
            grid: {{ left: '15%', right: '10%', top: '10%', bottom: '10%' }},
            xAxis: {{
                type: 'value',
                name: '项目天数',
                min: 0, max: 80,
                axisLabel: {{ formatter: '第{{value}}天' }},
                splitLine: {{ lineStyle: {{ type: 'dashed' }} }}
            }},
            yAxis: {{
                type: 'category',
                data: data.map(item => item.name),
                axisLabel: {{ fontSize: 13, color: '#333' }}
            }},
            series: [{{
                type: 'custom',
                renderItem: function(params, api) {{
                    var categoryIndex = api.value(2);
                    var start = api.coord([api.value(0), categoryIndex]);
                    var end = api.coord([api.value(1), categoryIndex]);
                    var height = api.size([0, 1])[1] * 0.6;
                    return {{
                        type: 'rect',
                        shape: {{ x: start[0], y: start[1] - height / 2, width: end[0] - start[0], height: height, r: 4 }},
                        style: api.style()
                    }};
                }},
                encode: {{ x: [0, 1], y: 2 }},
                data: data.map((item, index) => ({{
                    value: [item.value[0], item.value[1], index],
                    itemStyle: item.itemStyle
                }}))
            }}]
        }};
        chart.setOption(option);
        window.addEventListener('resize', () => chart.resize());
    </script>
</body>
</html>"""
    
    echarts_file = f"{CONFIG['output_dir']}/03_甘特图.html"
    with open(echarts_file, 'w', encoding='utf-8') as f:
        f.write(echarts_html)
    print(f"✅ ECharts 甘特图已生成: {echarts_file}")
    print("   提示：用浏览器打开即可查看交互式图表")
    
    return {"mermaid": mermaid_file, "echarts": echarts_file}

# ==================== 步骤3：火山引擎生成封面背景 ====================
def step3_generate_cover_image(project_name):
    """
    调用火山引擎 Seedream API 生成封面背景图
    """
    print(f"\n{'='*60}")
    print("步骤3：火山引擎 Seedream 生成封面背景")
    print(f"{'='*60}")
    
    prompt = f"商务科技风格标书封面背景，深蓝色调，金色几何线条装饰，专业大气，留白区域适合叠加文字标题，现代简约，高品质，16:9比例"
    
    # 检查 API Key
    if not CONFIG['seedream_api_key']:
        print("⚠️ 警告：SEEDREAM_API_KEY 未设置，跳过图片生成")
        print("   请设置环境变量：SEEDREAM_API_KEY")
        return None
    
    # 调用 Seedream 脚本
    import subprocess
    seedream_script = "./skills/seedream-image-generation/scripts/seedream.py"
    
    cmd = [
        "python", seedream_script,
        "--prompt", prompt,
        "--size", "2K",
        "--download_dir", CONFIG['output_dir']
    ]
    
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode == 0:
            print("✅ 封面背景图生成成功")
            # 解析输出找到图片路径
            output = result.stdout
            if "local_path" in output:
                # 简单提取路径
                start = output.find('"local_path": "') + 15
                end = output.find('"', start)
                image_path = output[start:end]
                print(f"   图片路径: {image_path}")
                return image_path
        else:
            print(f"⚠️ 生成失败: {result.stderr[:200]}")
            return None
    except Exception as e:
        print(f"⚠️ 调用失败: {e}")
        return None

# ==================== 步骤4：draw.io 精确流程图 ====================
def step4_generate_drawio_diagram():
    """
    生成 draw.io 可编辑的流程图文件
    """
    print(f"\n{'='*60}")
    print("步骤4：draw.io 精确流程图")
    print(f"{'='*60}")
    
    # draw.io 的 XML 格式流程图
    drawio_xml = """<mxfile host="app.diagrams.net" modified="2024-01-01T00:00:00.000Z" 
    agent="Mozilla/5.0" version="21.0.0" etag="abc123" type="device">
    <diagram name="审计流程" id="audit-flow">
        <mxGraphModel dx="1422" dy="762" grid="1" gridSize="10" guides="1" 
            tooltips="1" connect="1" arrows="1" fold="1" page="1" 
            pageScale="1" pageWidth="1169" pageHeight="827" math="0" shadow="0">
            <root>
                <mxCell id="0" />
                <mxCell id="1" parent="0" />
                
                <!-- 节点定义 -->
                <mxCell id="2" value="项目立项" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#0A1F3F;fontColor=#FFFFFF;strokeColor=none;" vertex="1" parent="1">
                    <mxGeometry x="400" y="40" width="200" height="50" as="geometry" />
                </mxCell>
                
                <mxCell id="3" value="组建审计组" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#1A5C6E;fontColor=#FFFFFF;strokeColor=none;" vertex="1" parent="1">
                    <mxGeometry x="400" y="120" width="200" height="50" as="geometry" />
                </mxCell>
                
                <mxCell id="4" value="制定审计方案" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#1A5C6E;fontColor=#FFFFFF;strokeColor=none;" vertex="1" parent="1">
                    <mxGeometry x="400" y="200" width="200" height="50" as="geometry" />
                </mxCell>
                
                <mxCell id="5" value="进场实施" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#C5955C;fontColor=#FFFFFF;strokeColor=none;" vertex="1" parent="1">
                    <mxGeometry x="400" y="280" width="200" height="50" as="geometry" />
                </mxCell>
                
                <mxCell id="6" value="数据分析与报告" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#1A5C6E;fontColor=#FFFFFF;strokeColor=none;" vertex="1" parent="1">
                    <mxGeometry x="400" y="360" width="200" height="50" as="geometry" />
                </mxCell>
                
                <mxCell id="7" value="成果提交" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#0A1F3F;fontColor=#FFFFFF;strokeColor=none;" vertex="1" parent="1">
                    <mxGeometry x="400" y="440" width="200" height="50" as="geometry" />
                </mxCell>
                
                <!-- 连接线 -->
                <mxCell id="e1" style="edgeStyle=orthogonalEdgeStyle;rounded=1;orthogonalLoop=1;jettySize=auto;html=1;strokeColor=#C5955C;strokeWidth=2;" edge="1" parent="1" source="2" target="3">
                    <mxGeometry relative="1" as="geometry" />
                </mxCell>
                <mxCell id="e2" style="edgeStyle=orthogonalEdgeStyle;rounded=1;orthogonalLoop=1;jettySize=auto;html=1;strokeColor=#C5955C;strokeWidth=2;" edge="1" parent="1" source="3" target="4">
                    <mxGeometry relative="1" as="geometry" />
                </mxCell>
                <mxCell id="e3" style="edgeStyle=orthogonalEdgeStyle;rounded=1;orthogonalLoop=1;jettySize=auto;html=1;strokeColor=#C5955C;strokeWidth=2;" edge="1" parent="1" source="4" target="5">
                    <mxGeometry relative="1" as="geometry" />
                </mxCell>
                <mxCell id="e4" style="edgeStyle=orthogonalEdgeStyle;rounded=1;orthogonalLoop=1;jettySize=auto;html=1;strokeColor=#C5955C;strokeWidth=2;" edge="1" parent="1" source="5" target="6">
                    <mxGeometry relative="1" as="geometry" />
                </mxCell>
                <mxCell id="e5" style="edgeStyle=orthogonalEdgeStyle;rounded=1;orthogonalLoop=1;jettySize=auto;html=1;strokeColor=#C5955C;strokeWidth=2;" edge="1" parent="1" source="6" target="7">
                    <mxGeometry relative="1" as="geometry" />
                </mxCell>
            </root>
        </mxGraphModel>
    </diagram>
</mxfile>"""
    
    drawio_file = f"{CONFIG['output_dir']}/04_精确流程图.drawio"
    with open(drawio_file, 'w', encoding='utf-8') as f:
        f.write(drawio_xml)
    
    print(f"✅ draw.io 流程图已生成: {drawio_file}")
    print("   提示：用 draw.io 打开编辑，导出高清 PNG (scale 8)")
    
    return drawio_file

# ==================== 步骤5：powerpoint-pptx 整合 ====================
def step5_integrate_pptx(cover_image, content_md, charts):
    """
    使用 python-pptx 整合所有素材到 PPT 模板
    """
    print(f"\n{'='*60}")
    print("步骤5：powerpoint-pptx 整合标书")
    print(f"{'='*60}")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("⚠️ 缺少 python-pptx 库，请安装: pip install python-pptx")
        return None
    
    # 检查模板文件
    if not os.path.exists(CONFIG['template_pptx']):
        print(f"⚠️ 模板文件不存在: {CONFIG['template_pptx']}")
        print("   将创建空白演示文稿")
        prs = Presentation()
    else:
        prs = Presentation(CONFIG['template_pptx'])
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 获取空白布局
    blank_layout = prs.slide_layouts[6]  # 通常索引6是空白布局
    
    # 1. 添加封面页
    slide = prs.slides.add_slide(blank_layout)
    
    # 添加封面背景图
    if cover_image and os.path.exists(cover_image):
        slide.shapes.add_picture(cover_image, Inches(0), Inches(0), 
                                   width=prs.slide_width, height=prs.slide_height)
    
    # 添加标题文字
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = CONFIG['project_name']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "投标方案"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(197, 149, 92)  # 铜金色
    p2.alignment = PP_ALIGN.CENTER
    
    # 2. 添加目录页
    slide = prs.slides.add_slide(blank_layout)
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "目 录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 31, 63)  # 深蓝色
    
    # 添加目录项
    toc_items = [
        "一、项目理解与服务方案",
        "二、技术方案与方法论",
        "三、项目组织架构与人员配置",
        "四、质量保证措施",
        "五、进度计划与里程碑",
        "六、风险应对措施"
    ]
    
    for i, item in enumerate(toc_items):
        box = slide.shapes.add_textbox(Inches(2), Inches(1.5 + i*0.8), Inches(9), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(10, 31, 63)
    
    # 3. 添加甘特图页面（图片占位）
    slide = prs.slides.add_slide(blank_layout)
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "项目进度计划"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 31, 63)
    
    # 添加占位提示
    placeholder = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(2))
    tf = placeholder.text_frame
    p = tf.paragraphs[0]
    p.text = "[请插入甘特图截图]\n\n提示：打开 03_甘特图.html 用浏览器截图"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER
    
    # 4. 添加流程图页面（占位）
    slide = prs.slides.add_slide(blank_layout)
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "审计实施流程"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 31, 63)
    
    placeholder = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(2))
    tf = placeholder.text_frame
    p = tf.paragraphs[0]
    p.text = "[请插入流程图]\n\n提示：打开 04_精确流程图.drawio 导出 PNG"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER
    
    # 保存最终文件
    output_pptx = f"{CONFIG['output_dir']}/05_标书整合版.pptx"
    prs.save(output_pptx)
    
    print(f"✅ 标书 PPT 已整合: {output_pptx}")
    print(f"   包含: 封面页 + 目录页 + 甘特图占位 + 流程图占位")
    
    return output_pptx

# ==================== 主函数 ====================
def main():
    parser = argparse.ArgumentParser(description='融策标书终极组合生成器')
    parser.add_argument('--project', '-p', required=True, help='项目名称')
    parser.add_argument('--type', '-t', default='绩效评价', help='项目类型')
    parser.add_argument('--output', '-o', default='./output', help='输出目录')
    args = parser.parse_args()
    
    # 更新配置
    CONFIG['project_name'] = args.project
    CONFIG['project_type'] = args.type
    CONFIG['output_dir'] = args.output
    
    # 创建输出目录
    os.makedirs(CONFIG['output_dir'], exist_ok=True)
    
    print(f"\n{'#'*60}")
    print(f"# 融策标书终极组合生成器 v1.0")
    print(f"# 项目: {args.project}")
    print(f"# 类型: {args.type}")
    print(f"# 输出: {args.output}")
    print(f"{'#'*60}")
    
    # 执行五步骤
    start_time = time.time()
    
    # Step 1: 内容生成
    content = step1_generate_content(args.project, args.type)
    
    # Step 2: 图表生成
    charts = step2_generate_charts()
    
    # Step 3: 封面图片
    cover_image = step3_generate_cover_image(args.project)
    
    # Step 4: 精确流程图
    drawio_file = step4_generate_drawio_diagram()
    
    # Step 5: PPT整合
    pptx_file = step5_integrate_pptx(cover_image, content, charts)
    
    elapsed = time.time() - start_time
    
    # 输出汇总
    print(f"\n{'='*60}")
    print("生成完成！文件清单")
    print(f"{'='*60}")
    print(f"📄 内容大纲:    {CONFIG['output_dir']}/01_内容大纲.md")
    print(f"📊 Mermaid流程: {CONFIG['output_dir']}/02_流程图.mmd")
    print(f"📈 ECharts甘特:  {CONFIG['output_dir']}/03_甘特图.html")
    print(f"🎨 封面图片:    {cover_image or '（未生成）'}")
    print(f"🔧 draw.io流程: {CONFIG['output_dir']}/04_精确流程图.drawio")
    print(f"📑 PPT整合版:   {CONFIG['output_dir']}/05_标书整合版.pptx")
    print(f"\n⏱️  耗时: {elapsed:.1f} 秒")
    print(f"{'='*60}")
    
    print("\n💡 后续操作建议:")
    print("   1. 打开 03_甘特图.html → 浏览器截图 → 插入 PPT")
    print("   2. 打开 04_精确流程图.drawio → draw.io 导出 PNG → 插入 PPT")
    print("   3. 在 PPT 中微调文字和排版")
    print("   4. 导出为 PDF 用于投标")

if __name__ == '__main__':
    main()
