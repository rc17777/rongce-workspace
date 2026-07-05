"""Update PPT cover title and speech doc title to option A"""
import sys, os, shutil
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Update PPT ──
PPT_SRC = r'D:\openclaw-workspace\output\v6_clean.pptx'
PPT_DST = r'D:\openclaw-workspace\output\v6_final.pptx'
PPT_DESK = r'C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v6_模板.pptx'

prs = Presentation(PPT_SRC)
slide1 = prs.slides[0]

# Find and update the title text box (the one with 提升公司全员审计风险意识)
for shape in slide1.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if '提升公司全员审计风险意识' in text:
            # Clear and set new title
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '筑牢防线·守好底线'
            p.font.name = '微软雅黑'
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
            print(f'Updated PPT title: {p.text}')
            break

# Find and update subtitle
for shape in slide1.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if '四川融策会计师事务所' in text or '2026年6月' in text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '穿透式监管时代的企业审计风险防控'
            p.font.name = '微软雅黑'
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            p.alignment = PP_ALIGN.CENTER
            # Add second line
            p2 = tf.add_paragraph()
            p2.text = '四川融策会计师事务所  |  2026年6月  |  四川轨道公司专题培训'
            p2.font.name = '微软雅黑'
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
            p2.alignment = PP_ALIGN.CENTER
            print(f'Updated PPT subtitle: {p.text}')
            break

prs.save(PPT_DST)

# Copy to desktop
for i in range(1, 20):
    target = PPT_DESK if i == 1 else f'{os.path.splitext(PPT_DESK)[0]}_{i}.pptx'
    try:
        if os.path.exists(target):
            os.remove(target)
        shutil.copy2(PPT_DST, target)
        print(f'PPT saved to: {target}')
        break
    except:
        continue

# ── Update Speech Doc ──
from docx import Document
from docx.shared import Pt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn

DOC_SRC = r'D:\openclaw-workspace\output\四川轨道公司审计风险培训-演讲稿_v6.docx'
DOC_DST = r'D:\openclaw-workspace\output\四川轨道公司审计风险培训-演讲稿_v6_final.docx'
DOC_DESK = r'C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-演讲稿_v6.docx'

doc = Document(DOC_SRC)

# Update first paragraph (title)
for p in doc.paragraphs:
    if '提升公司全员审计风险意识' in p.text:
        for run in p.runs:
            run.text = run.text.replace('提升公司全员审计风险意识', '筑牢防线·守好底线')
        print(f'Updated doc title')
        break

doc.save(DOC_DST)

for i in range(1, 20):
    target = DOC_DESK if i == 1 else f'{os.path.splitext(DOC_DESK)[0]}_{i}.docx'
    try:
        if os.path.exists(target):
            os.remove(target)
        shutil.copy2(DOC_DST, target)
        print(f'Doc saved to: {target}')
        break
    except:
        continue

print('\nDone! Both PPT and speech doc updated with new title.')
