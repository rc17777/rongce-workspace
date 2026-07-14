"""
融策·政府审计AI赋能方法论 — 投资人演示PPT
来源：研究方法论_核心机密.docx → 面向投资人的精简版
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
NAVY = RGBColor(0x0B, 0x1D, 0x3A)      # 主色 深藏蓝
DARK_BLUE = RGBColor(0x1A, 0x2D, 0x5E)  # 深蓝
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1) # 中蓝
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)  # 浅蓝
GOLD = RGBColor(0xD4, 0xA0, 0x1E)        # 金色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x84, 0x49)

# ── Helpers ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=BLACK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=12, color=BLACK, spacing=Pt(6)):
    """lines: list of (text, bold, font_size_override, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line[0] if isinstance(line, tuple) else line
        p.text = text
        fs = font_size
        clr = color
        bld = False
        if isinstance(line, tuple):
            bld = line[1] if len(line) > 1 else False
            fs = line[2] if len(line) > 2 else font_size
            clr = line[3] if len(line) > 3 else color
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
    return txBox

def add_tag(slide, left, top, width, height, text, bg_color, font_color=WHITE):
    shape = add_rect(slide, left, top, width, height, bg_color)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(2)
    return shape

def add_card(slide, left, top, width, height, title, items, title_bg=ACCENT_BLUE, body_bg=WHITE):
    """Card with title bar and body items"""
    # title bar
    bar = add_rect(slide, left, top, width, Inches(0.45), title_bg)
    bar.text_frame.word_wrap = True
    bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = bar.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    # body
    body = add_rect(slide, left, top + Inches(0.45), width, height - Inches(0.45), body_bg, RGBColor(0xDD, 0xDD, 0xDD))
    body.text_frame.word_wrap = True
    tf = body.text_frame
    tf.paragraphs[0].text = ''
    for i, item in enumerate(items[:-1]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = BLACK
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    return body

def add_slide_number(slide, num):
    add_text_box(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                 str(num), font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

def new_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, NAVY)
# accent line
add_rect(s, Inches(1.2), Inches(2.5), Inches(0.08), Inches(2.5), GOLD)
# title
add_text_box(s, Inches(1.6), Inches(2.3), Inches(10), Inches(1.2),
             '融策', font_size=24, color=GOLD, bold=True)
add_text_box(s, Inches(1.6), Inches(2.9), Inches(10), Inches(1.2),
             '政府审计信息化AI赋能方法论', font_size=36, color=WHITE, bold=True)
add_text_box(s, Inches(1.6), Inches(3.8), Inches(10), Inches(0.6),
             '基于混合AI架构的审计智能化解决方案', font_size=18, color=LIGHT_BLUE)
add_text_box(s, Inches(1.6), Inches(5.0), Inches(10), Inches(0.5),
             '四川融策会计师事务所 · 四川融策工程咨询公司', font_size=14, color=GRAY)
add_text_box(s, Inches(1.6), Inches(5.5), Inches(10), Inches(0.5),
             '2026年5月  |  版本 1.0  |  机密', font_size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '行业痛点与市场机会', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
             '政府审计领域存在结构性矛盾，形成巨大的效率缺口', font_size=16, color=GRAY)

problems = [
    ('📈 业务量激增', '社保审计、招投标审计、绩效评价、专项债等需求快速增长\n传统人力模式在数据密集场景下效率见顶'),
    ('👥 人才断层', '资深审计师信息化能力不足\nIT人才不懂审计业务逻辑\n"懂审计又懂信息化"的团队极度稀缺'),
    ('⏳ 交付瓶颈', '数据量的指数增长 vs 审计人力的线性增长\n形成不可逆的效率剪刀差'),
    ('🏰 竞争壁垒', '"懂审计的不懂信息化，懂信息化的不懂审计"\n手握真实业务+深度领域知识的团队稀缺'),
]
for i, (title, desc) in enumerate(problems):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.1)
    card = add_rect(s, x, y, Inches(2.8), Inches(2.8), LIGHT_GRAY)
    # icon-ish top
    add_text_box(s, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(0.5),
                 title, font_size=16, color=NAVY, bold=True)
    add_text_box(s, x + Inches(0.15), y + Inches(0.8), Inches(2.5), Inches(1.8),
                 desc, font_size=12, color=BLACK)

# bottom highlight
add_rect(s, Inches(0), Inches(5.3), Inches(13.333), Inches(1.2), RGBColor(0xF8, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(5.4), Inches(11), Inches(1.0),
             '核心命题：如何将AI能力系统性地注入传统政府审计流程，实现提质增效，\n形成可复用的方法论与工具平台？',
             font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# data point callout
add_rect(s, Inches(9.5), Inches(5.8), Inches(3.2), Inches(0.55), GOLD)
add_text_box(s, Inches(9.6), Inches(5.85), Inches(3), Inches(0.45),
             '目标：审计效率提升10倍+', font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — UNIQUE POSITIONING
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '我们的独特定位', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '融策为什么能做这件事', font_size=16, color=GRAY)

advantages = [
    ('🏗️', '业务壁垒', '手握大量真实的政府审计业务\n覆盖10+种审计业务类型\n这是纯技术公司无法短期获取的'),
    ('🧠', '领域知识', '深度理解审计法规、业务流程、客户需求\n资深审计师团队+实战方法论\n"审计语言"是AI落地的关键翻译层'),
    ('⏰', '先发优势', '同行尚未系统布局AI审计\n率先完成方法论构建\n12-18个月黄金窗口期'),
]
for i, (icon, title, desc) in enumerate(advantages):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.0)
    # card bg
    add_rect(s, x, y, Inches(3.7), Inches(3.5), LIGHT_GRAY)
    # icon circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), y + Inches(0.3), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = icon
    circle.text_frame.paragraphs[0].font.size = Pt(28)
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # title
    add_text_box(s, x + Inches(0.2), y + Inches(1.5), Inches(3.3), Inches(0.5),
                 title, font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    # desc
    add_text_box(s, x + Inches(0.2), y + Inches(2.1), Inches(3.3), Inches(1.2),
                 desc, font_size=12, color=BLACK, alignment=PP_ALIGN.CENTER)

# bottom tagline
add_rect(s, Inches(0), Inches(5.9), Inches(13.333), Inches(0.8), NAVY)
add_text_box(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.6),
             '融策 = 真实业务 + 领域知识 + 先发技术布局      这是AI审计赛道最稀缺的组合',
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(s, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION: THREE-IN-ONE ARCHITECTURE
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '"三合一"混合架构', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '不走单一技术路线。规则引擎 + 机器学习 + 大语言模型 三层协同', font_size=16, color=GRAY)

layers = [
    ('确定性层', '规则引擎 + 知识图谱', '合规判定、标准比对', '高召回、低误判\n法规直接转化\n输出可追溯引用', NAVY),
    ('概率性层', '机器学习 / 异常检测', '围标串标 / 异常交易 / 风险预警', '模式识别\n可解释模型优先\n人工复核确认', DARK_BLUE),
    ('认知层', '大语言模型（LLM）', '文档理解 / 报告生成 / 交互分析', '私有化部署\n数据不出域\nAI不做最终判定', ACCENT_BLUE),
]
for i, (label, tech, scenario, features, bg_color) in enumerate(layers):
    y = Inches(2.0 + i * 1.65)
    # left label
    add_rect(s, Inches(1), y, Inches(1.8), Inches(1.35), bg_color)
    add_text_box(s, Inches(1.1), y + Inches(0.15), Inches(1.6), Inches(0.5),
                 label, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.1), y + Inches(0.65), Inches(1.6), Inches(0.5),
                 tech, font_size=11, color=RGBColor(0xCC, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
    # right content
    add_rect(s, Inches(3), y, Inches(5.5), Inches(1.35), LIGHT_GRAY)
    add_text_box(s, Inches(3.2), y + Inches(0.1), Inches(5.1), Inches(0.4),
                 f'适用：{scenario}', font_size=13, color=NAVY, bold=True)
    add_text_box(s, Inches(3.2), y + Inches(0.5), Inches(5.1), Inches(0.7),
                 features, font_size=11, color=BLACK)
    # arrow
    if i < 2:
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.8), y + Inches(1.35), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

# core principle
add_rect(s, Inches(9), Inches(2.0), Inches(3.8), Inches(4.65), NAVY)
add_text_box(s, Inches(9.3), Inches(2.3), Inches(3.2), Inches(0.5),
             '⚡ 硬原则', font_size=18, color=GOLD, bold=True)
principles_text = [
    '① 审计思维主导，技术服务于业务',
    '② 渐进式替代，不可一步到位',
    '③ AI不做最终审计结论',
    '④ 所有判定可追溯、可举证',
    '⑤ 数据安全优先，本地化部署',
    '',
    '━━━━━━━━━━━━━',
    '',
    '三引擎依次调度：',
    '规则引擎 → 合规初筛',
    'ML引擎 → 深度检测',
    'LLM引擎 → 解读生成',
    '',
    '人审最终结论。全程留痕。',
]
add_multi_text(s, Inches(9.3), Inches(3.0), Inches(3.2), Inches(3.5),
               principles_text, font_size=11, color=WHITE, spacing=Pt(3))

add_slide_number(s, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — TECHNICAL ARCHITECTURE
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.35), Inches(11), Inches(0.55),
             '四层技术架构', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(0.85), Inches(11), Inches(0.4),
             '从数据到应用的全链路AI基础设施', font_size=14, color=GRAY)

# Architectured layers: bottom-up stack with compact height
arch_layers = [
    # (label, bar_desc, detail_left, detail_right, color, y)
    ('应用层', '面向审计师交付',
     '审计报告  |  疑点管理  |  风险看板',
     '交互式审计助手  |  工作底稿生成', ACCENT_BLUE, Inches(5.8)),
    ('智能层', '三层引擎协同调度',
     '规则引擎：合规初筛（高召回）',
     'ML引擎：异常检测  |  LLM引擎：认知理解', DARK_BLUE, Inches(4.5)),
    ('知识层', '审计经验结构化',
     '审计法规库：法规→可执行规则',
     '案例知识图谱  |  行业指标库', NAVY, Inches(3.2)),
    ('数据层', '多源异构数据标准化',
     '多源接入适配 → 智能清洗',
     '语义标准化 → 数据湖构建', RGBColor(0x1A, 0x1A, 0x2E), Inches(1.9)),
]

for label, bar_desc, detail_left, detail_right, color, y in arch_layers:
    # Layer label strip
    lbl_strip = add_rect(s, Inches(1), y, Inches(1.6), Inches(0.95), color)
    add_text_box(s, Inches(1.05), y + Inches(0.1), Inches(1.5), Inches(0.45),
                 label, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.05), y + Inches(0.55), Inches(1.5), Inches(0.35),
                 bar_desc, font_size=9, color=RGBColor(0xCC, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
    # Left detail
    left_box = add_rect(s, Inches(2.8), y, Inches(5.0), Inches(0.95), LIGHT_GRAY)
    add_text_box(s, Inches(2.95), y + Inches(0.15), Inches(4.7), Inches(0.65),
                 detail_left, font_size=12, color=BLACK)
    # Right detail
    right_box = add_rect(s, Inches(7.95), y, Inches(4.6), Inches(0.95), RGBColor(0xE8, 0xEC, 0xF0))
    add_text_box(s, Inches(8.1), y + Inches(0.15), Inches(4.3), Inches(0.65),
                 detail_right, font_size=11, color=GRAY)

# Down arrows between layers
for i in range(3):
    arrow_y = Inches(1.9 + i * 1.3) + Inches(0.95)
    arrow = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6.1), arrow_y, Inches(0.6), Inches(0.28))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY
    arrow.line.fill.background()
    arrow.rotation = 180.0

# Bottom insight bar
add_rect(s, Inches(0), Inches(6.9), Inches(13.333), Inches(0.5), RGBColor(0xF8, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(6.93), Inches(11.3), Inches(0.45),
             '数据层汇聚 → 知识层结构化 → 智能层分析 → 应用层交付    硬原则：AI不做最终审计结论，所有判定须经审计师确认',
             font_size=11, color=NAVY, bold=False, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — DATA + KNOWLEDGE LAYERS
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '数据层 → 知识层：把原始数据变成机器可用的知识', font_size=26, color=NAVY, bold=True)

# Left: Data Pipeline
add_rect(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.6), NAVY)
add_text_box(s, Inches(0.7), Inches(1.65), Inches(5.4), Inches(0.5),
             '数据层 — 多源异构数据标准化', font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

data_steps = [
    ('① 多源接入', '政务系统 / 财务账套\n公开招投标 / API / 数据库'),
    ('② 智能清洗', '缺失值填充 / Benford异常检测\n去重 / 冲突解决'),
    ('③ 语义标准化', '编码映射 / 科目对齐\n单位名称消歧'),
    ('④ 数据湖', '按项目分区 / 增量更新\n元数据与血缘追踪'),
]
for i, (title, desc) in enumerate(data_steps):
    x = Inches(0.5 + i * 1.45)
    y = Inches(2.4)
    add_rect(s, x, y, Inches(1.35), Inches(0.55), ACCENT_BLUE)
    add_text_box(s, x + Inches(0.05), y + Inches(0.02), Inches(1.25), Inches(0.5),
                 title, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.05), y + Inches(0.65), Inches(1.35), Inches(1.0),
                 desc, font_size=9, color=BLACK, alignment=PP_ALIGN.CENTER)
    if i < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.35), y + Inches(0.12), Inches(0.1), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

# Right: Knowledge Layer
add_rect(s, Inches(7), Inches(1.6), Inches(5.8), Inches(0.6), DARK_BLUE)
add_text_box(s, Inches(7.2), Inches(1.65), Inches(5.4), Inches(0.5),
             '知识层 — 把审计经验结构化', font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

kbs = [
    ('📋 审计法规库', '法律/规章/地方性规定\n→ 转化为可执行规则\n→ 支持版本管理与热更新'),
    ('🕸️ 案例知识图谱', '历史审计问题→实体关系网络\n单位-人员-资金-项目-问题\n→ 关联穿透分析'),
    ('📊 行业指标库', '各行业标准比率/阈值\n正常区间库\n→ 异常检测基准参照'),
]
for i, (title, desc) in enumerate(kbs):
    x = Inches(7 + i * 2.05)
    y = Inches(2.5)
    add_rect(s, x, y, Inches(1.9), Inches(2.3), LIGHT_GRAY)
    add_text_box(s, x + Inches(0.1), y + Inches(0.1), Inches(1.7), Inches(0.5),
                 title, font_size=14, color=NAVY, bold=True)
    add_text_box(s, x + Inches(0.1), y + Inches(0.7), Inches(1.7), Inches(1.4),
                 desc, font_size=11, color=BLACK)

# Bottom insight
add_rect(s, Inches(0), Inches(5.4), Inches(13.333), Inches(0.8), RGBColor(0xF8, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
             '💡 知识层是融策的核心资产 — 把资深审计师脑子里的隐性经验变成机器可用的显性知识',
             font_size=15, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(s, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — KEY APPLICATION SCENARIOS
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '十一大应用场景覆盖', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '从社保审计到专项成本核算，AI赋能路径已清晰定义', font_size=15, color=GRAY)

scenarios = [
    ('社保资金审计', '人周→小时', '✅ 已有方法', NAVY),
    ('招投标审计', '围标串标检测', '✅ POC优先', ACCENT_BLUE),
    ('收支审计', '人天→分钟', '✅ 已有方法', NAVY),
    ('绩效评价', '指标+对标+报告', '✅ 已有方法', ACCENT_BLUE),
    ('监督检查', '交叉比对+线索发现', '✅ 已有方法', NAVY),
    ('预算执行审计', '编制→执行→决算', '✅ 已有方法', ACCENT_BLUE),
    ('单位内部审计', '持续监控+底稿', '✅ 已有方法', NAVY),
    ('内控体系建设', '流程建模+缺陷识别', '✅ 已有方法', ACCENT_BLUE),
    ('专项债一案两书', '模板+测算+合规', '🟡 待开发', NAVY),
    ('行政事业专项审计', '方案生成+多专项', '🟡 待开发', ACCENT_BLUE),
    ('专项成本核算', '归集+对标+分摊', '🟡 待开发', NAVY),
]

for i, (name, highlight, status, bg) in enumerate(scenarios):
    row = i // 4
    col = i % 4
    x = Inches(0.6 + col * 3.15)
    y = Inches(1.9 + row * 1.8)
    card = add_rect(s, x, y, Inches(2.95), Inches(1.55), LIGHT_GRAY)
    add_text_box(s, x + Inches(0.15), y + Inches(0.1), Inches(2.65), Inches(0.5),
                 name, font_size=15, color=NAVY, bold=True)
    add_text_box(s, x + Inches(0.15), y + Inches(0.65), Inches(2.65), Inches(0.4),
                 f'效率目标：{highlight}', font_size=11, color=ACCENT_BLUE, bold=True)
    tag_color = GREEN if '✅' in status else GOLD
    add_tag(s, x + Inches(0.15), y + Inches(1.1), Inches(1.2), Inches(0.3),
            status.replace('✅ ','').replace('🟡 ',''), tag_color, WHITE)

add_slide_number(s, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — TOP 3 KEY SCENARIOS (DEEP DIVE)
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '核心场景深潜', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '三个最具商业价值的场景', font_size=15, color=GRAY)

deep_dives = [
    ('🏥 社保资金审计', '百万级参保数据',
     ['多年度/多险种数据统一清洗对齐',
      '已故人员继续领取、重复参保等规则筛查',
      '聚类分析识别异常群体行为模式',
      '自动生成疑点清单 + 证据链描述'],
     '数据清洗与初筛：人周 → 小时'),
    ('📋 招投标审计', '围标串标检测',
     ['投标人-关联人-关联企业知识图谱构建',
      '报价规律性分析 + 文本相似度检测',
      '关联方共同投标模式识别',
      '多份投标文件关键信息自动比对'],
     '⭐ 推荐为首个POC场景'),
    ('💰 收支审计', '多系统数据融合',
     ['多年度收支数据自动采集清洗标准化',
      '收支匹配与勾稽关系自动校验',
      '大额无依据支出、截留收入等异常识别',
      '智能生成收支审计工作底稿'],
     '收支核对：人天 → 分钟'),
]
for i, (title, subtitle, points, result) in enumerate(deep_dives):
    x = Inches(0.5 + i * 4.2)
    # card header
    add_rect(s, x, Inches(1.9), Inches(3.9), Inches(0.7), NAVY)
    add_text_box(s, x + Inches(0.15), Inches(1.95), Inches(3.6), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(s, x + Inches(0.15), Inches(2.35), Inches(3.6), Inches(0.3),
                 subtitle, font_size=12, color=LIGHT_BLUE)
    # body
    add_rect(s, x, Inches(2.6), Inches(3.9), Inches(3.2), LIGHT_GRAY)
    add_multi_text(s, x + Inches(0.15), Inches(2.7), Inches(3.6), Inches(2.6),
                   points, font_size=12, spacing=Pt(8))
    # result bar
    add_rect(s, x, Inches(5.8), Inches(3.9), Inches(0.5), GOLD)
    add_text_box(s, x + Inches(0.15), Inches(5.85), Inches(3.6), Inches(0.4),
                 f'🎯 {result}', font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION ROADMAP
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '实施路线图', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '四阶段渐进式落地，从POC到平台化', font_size=15, color=GRAY)

phases = [
    ('Phase 1', '基础建设', '0-6个月',
     ['组建技术团队(2-3名AI工程师+1名架构师)',
      '完成招投标审计POC',
      '搭建数据清洗标准流水线',
      '构建基础法规库与规则库'],
     NAVY),
    ('Phase 2', '方法验证', '6-12个月',
     ['完成3个核心场景完整方案',
      '形成数据标准化SOP',
      '积累行业指标库',
      '工具平台V1.0内部上线'],
     DARK_BLUE),
    ('Phase 3', '平台化', '12-18个月',
     ['审计工具平台对外版本',
      '支持多租户、多项目并行',
      '知识库持续更新机制',
      '与主流财务/政务系统接口标准化'],
     ACCENT_BLUE),
    ('Phase 4', '规模化', '18个月+',
     ['覆盖主要政府审计业务类型',
      '形成行业解决方案包',
      '从"项目交付"→"平台+服务"',
      '建立持续复购的商业模型'],
     GOLD),
]

for i, (phase, title, timeline, items, color) in enumerate(phases):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.0)
    # phase connector line
    if i < 3:
        add_rect(s, x + Inches(3), y + Inches(2.0), Inches(0.4), Inches(0.06), GRAY)
    # phase badge
    add_rect(s, x, y, Inches(2.9), Inches(0.55), color)
    add_text_box(s, x + Inches(0.1), y + Inches(0.02), Inches(1.0), Inches(0.5),
                 phase, font_size=14, color=WHITE, bold=True)
    add_text_box(s, x + Inches(1.2), y + Inches(0.05), Inches(1.6), Inches(0.5),
                 f'{title} | {timeline}', font_size=11, color=WHITE, alignment=PP_ALIGN.RIGHT)
    # items
    add_rect(s, x, y + Inches(0.55), Inches(2.9), Inches(2.5), LIGHT_GRAY)
    add_multi_text(s, x + Inches(0.15), y + Inches(0.7), Inches(2.6), Inches(2.2),
                   items, font_size=11, spacing=Pt(6))
    # timeline
    add_text_box(s, x + Inches(0.1), y + Inches(3.2), Inches(2.7), Inches(0.5),
                 timeline, font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# bottom metric
add_rect(s, Inches(0), Inches(5.7), Inches(13.333), Inches(1.0), RGBColor(0xF8, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(5.8), Inches(11.3), Inches(0.8),
             '关键里程碑：POC（招投标审计）→ V1.0内部上线 → 对外版本 → 行业解决方案包',
             font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(s, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 — TECH STACK
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '技术选型', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '实用主义选型：可解释、可审计、政府场景适配优先', font_size=15, color=GRAY)

tech_items = [
    ('数据清洗', 'Python + Pandas/Polars', '灵活、生态丰富、审计团队易上手', NAVY),
    ('规则引擎', 'Drools / 自研轻量DSL', '可解释、可审计、支持规则版本管理', DARK_BLUE),
    ('知识图谱', 'Neo4j / 自研图存储', '关系网络分析核心能力、关联穿透', ACCENT_BLUE),
    ('ML框架', 'Scikit-learn + XGBoost', '可解释性强、政府场景适配度高', NAVY),
    ('LLM', '开源模型 + 领域微调', '私有化部署、数据不出域、成本可控', DARK_BLUE),
    ('可视化', 'draw.io + ECharts', '审计报告级图表、手绘/正式双风格', ACCENT_BLUE),
    ('后端', 'Python FastAPI', '轻量、AI生态原生、快速迭代', NAVY),
    ('前端', 'React / Vue', '成熟生态、组件丰富', DARK_BLUE),
]

for i, (module, tech, reason, item_color) in enumerate(tech_items):
    row = i // 2
    col_idx = i % 2
    x = Inches(0.6 + col_idx * 6.3)
    y = Inches(2.0 + row * 1.2)
    add_rect(s, x, y, Inches(5.9), Inches(1.0), LIGHT_GRAY)
    add_text_box(s, x + Inches(0.2), y + Inches(0.1), Inches(1.2), Inches(0.4),
                 module, font_size=14, color=item_color, bold=True)
    add_text_box(s, x + Inches(1.5), y + Inches(0.1), Inches(4.2), Inches(0.4),
                 tech, font_size=13, color=BLACK, bold=True)
    add_text_box(s, x + Inches(1.5), y + Inches(0.55), Inches(4.2), Inches(0.4),
                 reason, font_size=11, color=GRAY)

# bottom note
add_text_box(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
             '选型原则：可解释性 > 先进性  |  数据安全优先  |  开源/国产化偏好  |  审计团队可维护',
             font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(s, 10)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 — RISK & MITIGATION
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '风险与对策', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '已知风险均有明确应对方案', font_size=15, color=GRAY)

risks = [
    ('数据获取困难', '🟡 中', '提前与客户沟通数据范围，合同约定数据提供义务'),
    ('AI结果审计责任界定', '🔴 高', '明确AI为辅助工具，审计师对结论负责，系统全程留痕'),
    ('模型幻觉导致错误疑点', '🟡 中', 'LLM只用于解读和生成初稿，不做判定；输出可追溯'),
    ('法规变动导致规则失效', '🟡 中', '规则与模型版本化管理，支持热更新'),
    ('技术人才招聘困难', '🟡 中', '与高校/培训机构合作定向培养，逐步内部培养'),
    ('客户对AI审计接受度', '🟡 中', '渐进式导入，从辅助功能开始，用实际效果建立信任'),
]

for i, (risk, level, mitigation) in enumerate(risks):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.0 + row * 2.5)
    # risk card
    card = add_rect(s, x, y, Inches(3.9), Inches(2.2), LIGHT_GRAY)
    # level badge
    badge_color = RED if '🔴' in level else GOLD
    add_rect(s, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.35), badge_color)
    add_text_box(s, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.35),
                 level, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.95), y + Inches(0.15), Inches(2.8), Inches(0.4),
                 risk, font_size=15, color=NAVY, bold=True)
    # divider
    add_rect(s, x + Inches(0.15), y + Inches(0.65), Inches(3.6), Inches(0.02), GRAY)
    add_text_box(s, x + Inches(0.15), y + Inches(0.8), Inches(3.6), Inches(1.2),
                 f'对策：{mitigation}', font_size=12, color=BLACK)

add_slide_number(s, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 — INVESTMENT OPPORTUNITY
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
add_text_box(s, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             '投资机会', font_size=28, color=NAVY, bold=True)
add_text_box(s, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
             '为什么现在是投资融策的最佳时机', font_size=15, color=GRAY)

# Left: Why Now
add_rect(s, Inches(0.5), Inches(1.9), Inches(6), Inches(0.6), NAVY)
add_text_box(s, Inches(0.7), Inches(1.95), Inches(5.6), Inches(0.5),
             '🔥 为什么是现在', font_size=18, color=WHITE, bold=True)

why_now = [
    '政策东风：政府数字化转型加速，审计信息化预算持续增长',
    '市场空白：审计AI赛道尚无头部玩家，先发优势可转化为网络效应',
    '技术成熟：LLM+知识图谱+ML三件套已具备工程化落地条件',
    '数据飞轮：每个审计项目产生数据→反哺模型→服务更好→更多项目',
    '商业模式升级：从人力密集型项目交付 → 平台+服务，毛利率跃升',
]
add_multi_text(s, Inches(0.7), Inches(2.7), Inches(5.6), Inches(3.0),
               why_now, font_size=13, spacing=Pt(10))

# Right: What We Need
add_rect(s, Inches(7), Inches(1.9), Inches(5.8), Inches(0.6), GOLD)
add_text_box(s, Inches(7.2), Inches(1.95), Inches(5.4), Inches(0.5),
             '🎯 本轮融资用途', font_size=18, color=WHITE, bold=True)

use_of_funds = [
    ('技术团队建设', '2-3名AI/数据工程师 + 1名架构师'),
    ('POC开发与验证', '招投标审计POC + 2个场景扩展'),
    ('基础设施', '算力/数据存储/安全合规'),
    ('知识库建设', '法规库+案例库+行业指标库'),
    ('市场拓展', '标杆客户落地 + 行业方案输出'),
]
for i, (item, detail) in enumerate(use_of_funds):
    y = Inches(2.7 + i * 0.85)
    add_text_box(s, Inches(7.2), y, Inches(2.0), Inches(0.4),
                 f'▸ {item}', font_size=14, color=NAVY, bold=True)
    add_text_box(s, Inches(9.2), y, Inches(3.4), Inches(0.4),
                 detail, font_size=12, color=BLACK)

# Bottom: Key Metrics
add_rect(s, Inches(0), Inches(5.6), Inches(13.333), Inches(0.7), RGBColor(0xF8, 0xF0, 0xE0))
metrics = [
    ('11', '应用场景'),
    ('3层', '智能引擎'),
    ('10x', '效率提升目标'),
    ('18月', '到规模化'),
    ('0', '同类竞品'),
]
for i, (num, label) in enumerate(metrics):
    x = Inches(1.0 + i * 2.6)
    add_text_box(s, x, Inches(5.65), Inches(2.2), Inches(0.35),
                 num, font_size=28, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x, Inches(6.0), Inches(2.2), Inches(0.3),
                 label, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 12)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 — THANK YOU
# ═══════════════════════════════════════════════════════════
s = new_slide()
add_bg(s, NAVY)
add_rect(s, Inches(1.2), Inches(2.0), Inches(0.08), Inches(2.5), GOLD)
add_text_box(s, Inches(1.6), Inches(2.3), Inches(10), Inches(1.2),
             '谢谢', font_size=48, color=WHITE, bold=True)
add_text_box(s, Inches(1.6), Inches(3.5), Inches(10), Inches(0.6),
             '融策 · 政府审计信息化AI赋能', font_size=20, color=LIGHT_BLUE)
add_text_box(s, Inches(1.6), Inches(4.3), Inches(10), Inches(0.5),
             '四川融策会计师事务所  ·  四川融策工程咨询公司', font_size=14, color=GRAY)
add_text_box(s, Inches(1.6), Inches(4.9), Inches(10), Inches(0.5),
             '成都 · 2026', font_size=12, color=GRAY)

# ── Save ──
output_path = r'D:\openclaw-workspace\output\融策AI审计赋能方法论_投资人演示.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
