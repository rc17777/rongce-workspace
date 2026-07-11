#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""修改轨道培训PPT - 替换1号文/2号文为成都市两政策"""
import sys,os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import copy

src = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_616.pptx'
dst = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_更新版.pptx'

prs = Presentation(src)

def set_shape_text(shape, new_texts):
    """设置shape的文本, new_texts是字符串或字符串列表"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 清除现有文本
    for para in tf.paragraphs:
        para.clear()
    
    if isinstance(new_texts, str):
        new_texts = [new_texts]
    
    for i, text in enumerate(new_texts):
        if i >= len(tf.paragraphs):
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[i]
        
        # 设置文本
        run = p.add_run()
        run.text = text
        
        # 保持原有字体样式(复制第一个run的样式)
        # 默认用微软雅黑
        run.font.name = '微软雅黑'
        
        # 如果是标题(第一个段落), 设粗体
        if i == 0:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

def replace_slide_content(slide_idx, title, badge, content_lines):
    """替换幻灯片内容: 标题、标签、正文"""
    slide = prs.slides[slide_idx]
    shapes = slide.shapes
    shapes_list = list(shapes)  # 按顺序
    
    # 找到文本shape
    text_shapes = [s for s in shapes_list if s.has_text_frame]
    
    if len(text_shapes) >= 3:
        # Shape假设结构: [0]背景图 [1]logo [2]标题 [3]标签矩形 [4]正文
        title_shape = text_shapes[0]  # 第一个文本shape通常是标题
        badge_shape = text_shapes[1] if len(text_shapes) > 1 else None
        body_shape = text_shapes[2] if len(text_shapes) > 2 else None
        
        # 设置标题
        set_text_in_shape(title_shape, title, bold=True, size=28, color=RGBColor(0xFF,0xFF,0xFF))
        
        # 设置标签
        if badge_shape:
            set_text_in_shape(badge_shape, badge, bold=True, size=16, color=RGBColor(0xFF,0xFF,0xFF))
        
        # 设置正文
        if body_shape:
            set_multiline_text(body_shape, content_lines, size=13, color=RGBColor(0xFF,0xFF,0xFF))

def set_text_in_shape(shape, text, bold=False, size=14, color=RGBColor(0xFF,0xFF,0xFF)):
    """设置shape的文本"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 清除
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = '微软雅黑'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def set_multiline_text(shape, lines, size=13, color=RGBColor(0xFF,0xFF,0xFF)):
    """设置多行文本"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        run = p.add_run()
        run.text = line
        run.font.name = '微软雅黑'
        run.font.size = Pt(size)
        run.font.color.rgb = color
        
        if line.startswith('▎') or line.startswith('▶') or line.startswith('🔴') or line.startswith('⚠'):
            run.font.bold = True

# ================================================================
# 1. 修改第2页（培训议程） - slide index 1
# ================================================================
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '1号文' in para.text:
                # 替换1号文和2号文引用
                for run in para.runs:
                    run.text = run.text.replace('1号文·2号文·', '')
                    run.text = run.text.replace('国务院核心文件解读（', '国务院及成都市政策解读（')
print('✅ 第2页 议程文字已更新')

# ================================================================
# 2. 修改第3页（第一部分总览） - slide index 2
# ================================================================
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '1号文 → 2号文 → 15号文 → 46号令' in para.text:
                for run in para.runs:
                    run.text = run.text.replace(
                        '1号文 → 2号文 → 15号文 → 46号令：从技术底座到兜底保障的完整政策链条',
                        '15号文 → 46号令 → 成都市配套制度：国家顶层设计 × 地方落地执行'
                    )
print('✅ 第3页 总览文字已更新')

# ================================================================
# 3. 替换第4页（原1号文 → SW-2026-1196） - slide index 3
# ================================================================
slide4 = prs.slides[3]
text_shapes = [s for s in slide4.shapes if s.has_text_frame]
# text_shapes[0]=标题, [1]=badge标签, [2]=正文
if len(text_shapes) >= 3:
    set_text_in_shape(text_shapes[0], 'SW-2026-1196：《成都市属国有企业违规经营投资责任追究实施办法》',
                      bold=True, size=28)
    set_text_in_shape(text_shapes[1], '市级配套', bold=True, size=16)
    set_multiline_text(text_shapes[2], [
        '▎追责范围（13大类）',
        '集团管控 · 风险管理 · 购销管理 · 工程承包 · 金融业务 · 科技管理',
        '资金管理 · 担保活动 · 产权管理 · 固定资产投资 · 股权投资 · 资产管理 · 境外投资',
        '',
        '▎追责方式',
        '批评诫勉 / 组织处理（停职·免职·降职）/ 扣减薪酬 / 禁入限制 / 纪律处分 / 移送司法',
        '',
        '▎损失认定标准',
        '<100万元：一般损失  |  100万~1000万：较大损失  |  >1000万：重大损失',
        '',
        '▎核心原则',
        '🔴 重大决策终身问责  🟢 集体决策中提出反对意见可免责',
        '⚠ 损失未达标准但导致企业无法持续经营 → 视为重大损失',
        '',
        '▶ 与你的关系',
        '这是成都市国资委针对市属国企的专项追责办法，天府广场项目的每一项审计问题——',
        '停车场收费不清、商户租金拖欠、物业用房违规使用——都在这13大类追责范围之内',
    ], size=12)
    print('✅ 第4页 替换为 SW-2026-1196')

# ================================================================
# 4. 替换第5页（原2号文 → 成国资发〔2025〕15号） - slide index 4
# ================================================================
slide5 = prs.slides[4]
text_shapes5 = [s for s in slide5.shapes if s.has_text_frame]
if len(text_shapes5) >= 3:
    set_text_in_shape(text_shapes5[0], '成国资发〔2025〕15号：进一步规范市属国企资产租赁管理',
                      bold=True, size=26)
    set_text_in_shape(text_shapes5[1], '业务直击', bold=True, size=16)
    set_multiline_text(text_shapes5[2], [
        '▎公开招租（信息披露）',
        '📋 必须在市政府网+市国资委网+产权交易所+集团官网同步发布',
        '正式披露≥10个工作日 | 结果公示≥3个工作日',
        '',
        '▎租金定价',
        '💰 以评估/估价/询价/公开市场价为依据，综合考虑市场需求、区位、业态',
        '',
        '▎五条红线（严禁行为）',
        '❌ 拆分资产规避公开招租   ❌ 未经批准擅自出租',
        '❌ 擅自改变用途或租赁期限   ❌ 截留、挪用或私分租金收入',
        '❌ 违规转租（有违约行为的承租方不得同意续租）',
        '',
        '▎监管机制',
        '🔍 "三重一大"集体决策 → 报市国资委备案 → 年度自查自纠 → 审计/专项检查',
        '',
        '▶ 与你的关系',
        '天府广场商业租赁的核心问题——合同续租、租金定价、商户转租——都在这个文件里',
        '把红线画清楚了：什么能做、什么不能做、做了有什么后果',
    ], size=12)
    print('✅ 第5页 替换为 成国资发〔2025〕15号')

# ================================================================
# 5. 修改第8页（政策全景图） - slide index 7
# 将原来的1号文和2号文改为两个成都政策
# ================================================================
slide8 = prs.slides[7]
target_texts = {
    '1号文': '1196号\n(追责)',  # 原1号文位→1196号
    '看见': '定位',
    '数智化系统让异常无处遁形': '13大类追责情形全覆盖\n终身追责不设时限',
    '2号文': '15号\n(租赁)',  # 原2号文位→15号租赁
    '看准': '规范',
    '穿透监管框架界定查什么、查多深': '公开招租·定价机制\n租赁管理全流程规范',
    '四位一体：政策全景图': '六位一体：政策全景图（国家+市级）',
    '技术 ──→ 制度 ──→ 执行 ──→ 问责  =  完整闭环': '15号文→46号令→1196号→15号  =  国家·市级双层闭环',
}

for shape in slide8.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for old, new in target_texts.items():
                if old in para.text:
                    for run in para.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)

print('✅ 第8页 政策全景图已更新')

# ================================================================
# 保存
# ================================================================
prs.save(dst)
print(f'\n✅ 完成！新文件: {dst}')
print(f'   文件大小: {os.path.getsize(dst)//1024}KB')
