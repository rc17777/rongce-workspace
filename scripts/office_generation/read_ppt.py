#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""读取PPT每页文字"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation(r'D:\openclaw-workspace\temp\轨道培训意见征求稿_616.pptx')
for i in range(1, 11):
    slide = prs.slides[i]
    print(f'===== 第{i+1}页 =====')
    for si, shape in enumerate(slide.shapes):
        print(f'  [Shape{si}] name={shape.name} left={shape.left} top={shape.top} w={shape.width} h={shape.height}')
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                t = para.text.strip()
                if t:
                    print(f'    P{pi}: {t}')
        if shape.has_table:
            for ri, row in enumerate(shape.table.rows):
                cells = [c.text.strip() for c in row.cells]
                print(f'    Table R{ri}: {chr(124).join(cells)}')
    print()
