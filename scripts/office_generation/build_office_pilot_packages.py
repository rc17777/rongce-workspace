#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pathlib import Path
from datetime import datetime

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(r'C:\Users\scrccpa\Documents\Obsidian Vault\融策AI知识中枢\交付模板雏形\Office三件套试点交付包')

PACKS = {
    '工程审计': {
        'focus': '政府投资项目、招投标、工程程序、合同执行、资金支付、现场实施情况',
        'tagline': '从立项到现场，从合同到资金，形成工程审计闭环证据链。',
        'objectives': ['查清工程项目立项、招投标、合同执行、资金支付和现场实施情况。', '揭示应招未招、先建后招、虚假招标、工程量虚增、变更签证依据不足等问题。', '推动建设单位和主管部门完善工程管理制度，压实整改责任。'],
        'data': ['项目立项批复', '招标公告与招标文件', '投标文件及评标资料', '中标通知书', '合同及补充协议', '工程计量支付资料', '设计变更和签证资料', '竣工验收资料', '财政评审/结算审核资料', '现场照片和踏勘记录'],
        'risks': ['应招未招、先建后招、虚假招标', '投标人串通投标或围标陪标', '评标专家异常打分', '违规确定中标人', '合同条款与招标文件不一致', '工程变更签证依据不足', '工程量虚增或重复计量', '整改程序倒置或责任不清'],
        'methods': ['招投标全链条拆解', '项目台账与交易台账比对', '合同与支付凭证穿透核查', '现场踏勘验证工程量', '问题整改闭环跟踪'],
    },
    '医保卫健数据审计': {
        'focus': '医保基金、医院诊疗行为、第三方检验、公共卫生服务、老年人健康管理服务真实性',
        'tagline': '用多源数据把服务真实性、资金真实性和监管责任串起来。',
        'objectives': ['查清医保结算、诊疗服务、检验检查和公共卫生服务数据真实性。', '揭示串换诊疗项目、过度诊疗、死亡后服务记录、虚构体检服务套取补助等问题。', '推动主管部门完善数据治理、服务监管和资金绩效考核机制。'],
        'data': ['医保结算明细', '医院HIS诊疗明细', '检验检查项目明细', '第三方检测机构报告', '医保目录与项目编码表', '居民健康档案', '体检服务记录', '死亡人员和户籍注销数据', '基层医疗机构补助资金台账', '绩效考核资料'],
        'risks': ['串换诊疗项目', '过度诊疗和无指征检查', '重复收费或分解收费', '死亡后仍产生服务记录', '未满年龄提前纳入服务', '虚构体检服务套取补助', '规律性虚假电话号码或身份信息', '服务记录与资金拨付不匹配'],
        'methods': ['多源数据清洗与字段统一', '医保结算与诊疗记录交叉比对', '年龄/死亡/服务频次逻辑校验', 'SQL疑点筛查', '疑点回溯到原始病历和服务记录'],
    },
    '政策落实审计': {
        'focus': '重大政策部署、专项资金、区域协同、“两重”“两新”、项目落地、部门协同和整改闭环',
        'tagline': '沿政策链、项目链、资金链、责任链穿透政策落地效果。',
        'objectives': ['查清重大政策落实、资金安排、项目推进和部门协同情况。', '揭示政策执行打折扣、项目申报不实、资金闲置沉淀、监管缺位和整改不实等问题。', '推动政策落实由项目整改向制度完善和系统治理延伸。'],
        'data': ['政策文件和实施方案', '项目申报资料', '资金下达文件', '资金拨付明细', '项目实施台账', '主管部门审核资料', '部门协同会议纪要', '绩效目标和绩效评价资料', '整改台账', '回头看资料'],
        'risks': ['政策执行打折扣', '项目申报不实或审核不严', '资金闲置沉淀或截留挪用', '重复补贴或监管缺位', '部门信息共享不足', '项目落地慢或绩效不明显', '整改敷衍或虚假整改', '政策最后一公里堵点'],
        'methods': ['政策链条对标', '项目申报到资金拨付全流程穿透', '多部门数据交叉验证', '绩效目标与实际结果比对', '整改回头看'],
    },
}

BLUE = '0A1F3F'
TEAL = '1A5C6E'
GOLD = 'C5955C'
GRAY = 'F5F2EC'
LIGHT_TEAL = 'E8F1F2'
WHITE = 'FFFFFF'
DARK = '1F2933'


def hex_rgb(value):
    value = value.strip('#')
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def ppt_rgb(value):
    value = value.strip('#')
    return PptRGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def shade_cell(cell, color):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:fill'), color)
    tc_pr.append(shd)


def set_doc_style(doc):
    sec = doc.sections[0]
    sec.top_margin = Cm(2.4)
    sec.bottom_margin = Cm(2.0)
    sec.left_margin = Cm(2.6)
    sec.right_margin = Cm(2.6)
    style = doc.styles['Normal']
    style.font.name = '宋体'
    style._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    style.font.size = Pt(11)
    style.paragraph_format.line_spacing = 1.25


def add_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        run.font.name = '微软雅黑'
        run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
        run.font.color.rgb = hex_rgb(BLUE if level == 1 else TEAL)
        run.font.bold = True
    return p


def add_cover(doc, name, cfg):
    for _ in range(2):
        doc.add_paragraph('')
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('融策试点交付包')
    r.font.name = '微软雅黑'
    r._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
    r.font.size = Pt(15)
    r.font.color.rgb = hex_rgb(GOLD)
    r.bold = True

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f'{name}\n审计实施方案与专报模板')
    r.font.name = '微软雅黑'
    r._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
    r.font.size = Pt(24)
    r.font.color.rgb = hex_rgb(BLUE)
    r.bold = True

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(cfg['tagline'])
    r.font.name = '微软雅黑'
    r._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
    r.font.size = Pt(12)
    r.font.color.rgb = hex_rgb(TEAL)

    table = doc.add_table(rows=3, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'
    rows = [('适用场景', cfg['focus']), ('交付内容', '实施方案、取数清单、问题清单、专报模板、汇报提纲'), ('生成时间', datetime.now().strftime('%Y-%m-%d %H:%M'))]
    for i, (k, v) in enumerate(rows):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v
    for row in table.rows:
        for j, cell in enumerate(row.cells):
            shade_cell(cell, BLUE if j == 0 else GRAY)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.name = '微软雅黑'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                    run.font.size = Pt(10)
                    run.font.color.rgb = hex_rgb(WHITE if j == 0 else DARK)
                    run.font.bold = j == 0
    doc.add_page_break()


def bullet(doc, text):
    p = doc.add_paragraph(text, style='List Bullet')
    for run in p.runs:
        run.font.name = '宋体'
        run._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    return p


def make_docx(name, cfg, folder):
    doc = Document()
    set_doc_style(doc)
    add_cover(doc, name, cfg)

    add_heading(doc, '一、审计实施方案', 1)
    add_heading(doc, '（一）项目背景', 2)
    doc.add_paragraph(f'本项目围绕{cfg["focus"]}开展审计，重点揭示政策执行、业务真实性、程序合规、资金安全和绩效结果等方面存在的问题，推动整改和治理提升。')
    add_heading(doc, '（二）审计目标', 2)
    for x in cfg['objectives']:
        bullet(doc, x)
    add_heading(doc, '（三）审计重点', 2)
    for x in cfg['risks']:
        bullet(doc, x)
    add_heading(doc, '（四）审计方法', 2)
    for x in cfg['methods']:
        bullet(doc, x)

    add_heading(doc, '（五）审计作业路径', 2)
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Table Grid'
    hdr = table.rows[0].cells
    for i, h in enumerate(['阶段', '重点动作', '关键资料', '输出成果']):
        hdr[i].text = h
    steps = [
        ('审前调查', '收集政策、制度、台账、数据和业务资料', '政策文件、业务台账、资金明细', '审前调查记录'),
        ('风险画像', '按业务链条梳理风险点，形成疑点清单', '项目/对象/资金明细', '风险画像表'),
        ('数据分析', '对关键字段、金额、时间、对象进行比对筛查', '原始电子表、系统导出数据', '疑点清单'),
        ('现场核查', '对重点疑点开展访谈、踏勘、资料复核', '凭证、合同、现场照片', '取证单'),
        ('成果输出', '形成汇报材料、审计报告或专报', '证据链、整改台账', '专报/报告/汇报PPT'),
    ]
    for row in steps:
        cells = table.add_row().cells
        for i, v in enumerate(row):
            cells[i].text = v
    for r, row in enumerate(table.rows):
        for cell in row.cells:
            shade_cell(cell, BLUE if r == 0 else (GRAY if r % 2 == 0 else WHITE))
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.name = '微软雅黑'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                    run.font.size = Pt(9)
                    run.font.color.rgb = hex_rgb(WHITE if r == 0 else DARK)
                    run.font.bold = r == 0

    add_heading(doc, '二、审计专报初稿模板', 1)
    add_heading(doc, '（一）基本情况', 2)
    doc.add_paragraph(f'根据工作安排，审计组围绕{cfg["focus"]}开展审计，重点关注政策执行、业务真实性、资金安全和绩效结果等情况。')
    add_heading(doc, '（二）审计发现的主要问题', 2)
    doc.add_paragraph('1. 问题一：……')
    doc.add_paragraph('经审计发现，……反映出相关单位在制度执行、过程管控和责任落实方面存在薄弱环节。')
    doc.add_paragraph('2. 问题二：……')
    doc.add_paragraph('经比对相关数据和资料，发现……疑似存在业务真实性不足、资金管理不规范或政策执行不到位等问题。')
    add_heading(doc, '（三）审计建议', 2)
    for x in ['建议相关单位对照问题清单限期整改。', '建议主管部门完善制度流程，加强业务审核和数据监管。', '建议建立长效机制，推动问题整改从个案处理向系统治理延伸。']:
        bullet(doc, x)
    path = folder / f'{name}审计实施方案与专报模板-美化增强版.docx'
    doc.save(path)
    return path


def style_sheet(ws):
    header_fill = PatternFill('solid', fgColor=BLUE)
    sub_fill = PatternFill('solid', fgColor=LIGHT_TEAL)
    header_font = Font(name='微软雅黑', color=WHITE, bold=True, size=11)
    body_font = Font(name='微软雅黑', color=DARK, size=10)
    thin = Side(style='thin', color=GOLD)
    ws.freeze_panes = 'A2'
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
    for ridx, row in enumerate(ws.iter_rows(min_row=2), 2):
        for cell in row:
            cell.fill = sub_fill if ridx % 2 == 0 else PatternFill('solid', fgColor=WHITE)
            cell.font = body_font
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
            cell.alignment = Alignment(vertical='top', wrap_text=True)
    for row in ws.iter_rows():
        for cell in row:
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for col in range(1, ws.max_column + 1):
        width = 12 if col == 1 else 22
        ws.column_dimensions[get_column_letter(col)].width = width
    ws.row_dimensions[1].height = 26


def make_xlsx(name, cfg, folder):
    wb = Workbook()
    ws = wb.active
    ws.title = '取数清单'
    ws.append(['序号', '资料/数据名称', '提供单位', '格式要求', '用途', '备注'])
    for i, d in enumerate(cfg['data'], 1):
        ws.append([i, d, '主管部门/实施单位', '原始电子表或盖章纸质件', '用于核实业务真实性、资金流向或程序合规', '保留原始字段，不只提供汇总表'])
    style_sheet(ws)

    ws2 = wb.create_sheet('问题清单')
    ws2.append(['序号', '问题类别', '问题事实', '涉及金额/数量', '依据', '责任单位', '风险影响', '整改建议', '状态'])
    for i, risk in enumerate(cfg['risks'][:8], 1):
        ws2.append([i, risk, '', '', '', '', '', '', '待核实'])
    style_sheet(ws2)

    ws3 = wb.create_sheet('取证单样表')
    ws3.append(['字段', '内容'])
    for k in ['项目名称', '被审计单位', '取证事项', '取证日期', '审计人员', '被取证人员', '问题事实', '证据材料', '被审计单位说明', '审计组意见']:
        ws3.append([k, ''])
    style_sheet(ws3)

    ws4 = wb.create_sheet('汇报要点')
    ws4.append(['模块', '内容提示'])
    for k, v in [('项目背景', cfg['focus']), ('审计方法', '、'.join(cfg['methods'])), ('主要风险', '、'.join(cfg['risks'][:5])), ('整改建议', '限期整改、制度完善、长效治理')]:
        ws4.append([k, v])
    style_sheet(ws4)

    path = folder / f'{name}取数清单与问题取证样表-美化增强版.xlsx'
    wb.save(path)
    return path


def add_band(slide, y, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, Inches(13.333), h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(color)
    shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, size=24, color=BLUE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = '微软雅黑'
    p.font.size = PptPt(size)
    p.font.bold = bold
    p.font.color.rgb = ppt_rgb(color)
    return box


def add_pill(slide, x, y, text, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.45), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = '微软雅黑'
    p.font.size = PptPt(12)
    p.font.bold = True
    p.font.color.rgb = ppt_rgb(WHITE)
    return shape


def title_slide(prs, name, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, Inches(7.5), GRAY)
    add_band(slide, 0, Inches(0.42), BLUE)
    add_band(slide, Inches(6.95), Inches(0.55), GOLD)
    add_text(slide, Inches(0.75), Inches(1.55), Inches(11.8), Inches(0.45), '融策试点交付包', 18, GOLD, True)
    add_text(slide, Inches(0.75), Inches(2.15), Inches(11.8), Inches(1.0), f'{name}审计汇报提纲', 36, BLUE, True)
    add_text(slide, Inches(0.78), Inches(3.25), Inches(10.8), Inches(0.7), cfg['tagline'], 19, TEAL, False)
    add_pill(slide, Inches(0.78), Inches(4.35), '实施方案')
    add_pill(slide, Inches(3.45), Inches(4.35), '取数清单')
    add_pill(slide, Inches(6.12), Inches(4.35), '问题取证')
    add_pill(slide, Inches(8.79), Inches(4.35), '汇报专报')
    add_text(slide, Inches(0.8), Inches(6.2), Inches(8), Inches(0.4), datetime.now().strftime('%Y-%m-%d'), 13, DARK)


def content_slide(prs, title, bullets, accent=TEAL):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, Inches(0.28), BLUE)
    add_text(slide, Inches(0.58), Inches(0.58), Inches(11.5), Inches(0.55), title, 26, BLUE, True)
    add_band(slide, Inches(1.32), Inches(0.05), GOLD)
    y = Inches(1.65)
    for i, b in enumerate(bullets[:7], 1):
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.78), y + Inches(0.02), Inches(0.26), Inches(0.26))
        marker.fill.solid()
        marker.fill.fore_color.rgb = ppt_rgb(accent)
        marker.line.fill.background()
        add_text(slide, Inches(1.18), y - Inches(0.05), Inches(10.8), Inches(0.45), b, 19, DARK)
        y += Inches(0.62)
    return slide


def make_pptx(name, cfg, folder):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide(prs, name, cfg)
    content_slide(prs, '一、项目背景', [f'围绕{cfg["focus"]}开展审计', '关注政策执行、业务真实性、程序合规、资金安全和绩效结果', '输出问题清单、取证单、专报初稿和整改建议'], TEAL)
    content_slide(prs, '二、审计重点', cfg['risks'], GOLD)
    content_slide(prs, '三、审计方法', cfg['methods'], TEAL)
    content_slide(prs, '四、主要发现', ['问题一：……', '问题二：……', '问题三：……', '问题四：……'], GOLD)
    content_slide(prs, '五、风险影响', ['资金风险', '管理风险', '政策绩效风险', '廉政风险'], TEAL)
    content_slide(prs, '六、整改建议', ['限期整改', '完善制度', '压实责任', '建立长效机制'], GOLD)
    path = folder / f'{name}审计汇报PPT提纲-美化增强版.pptx'
    prs.save(path)
    return path


def write_readme(made):
    lines = [
        '# Office三件套试点交付包 - 美化增强版',
        '',
        f'生成时间：{datetime.now().strftime("%Y-%m-%d %H:%M")}',
        '',
        '## 本次增强',
        '',
        '- Word：新增正式封面、项目摘要表、作业路径表、统一标题和表格配色。',
        '- Excel：新增冻结首行、隔行底色、字段宽度、问题清单预置风险项、汇报要点页。',
        '- PPT：改为16:9定制版式，加入封面、色带、重点页和统一视觉体系。',
        '- 命名：全部以 `-美化增强版` 结尾，便于和旧版区分。',
        '',
        '## 文件清单',
        '',
    ]
    for p in made:
        lines.append(f'- `{p}`')
    (OUT / 'README-美化增强版.md').write_text('\n'.join(lines), encoding='utf-8')


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    made = []
    for name, cfg in PACKS.items():
        folder = OUT / name
        folder.mkdir(parents=True, exist_ok=True)
        made.append(make_docx(name, cfg, folder))
        made.append(make_xlsx(name, cfg, folder))
        made.append(make_pptx(name, cfg, folder))
    write_readme(made)
    print('OUT', OUT)
    for p in made:
        print(p)


if __name__ == '__main__':
    main()
