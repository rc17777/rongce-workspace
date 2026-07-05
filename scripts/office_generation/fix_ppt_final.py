#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""修复PPT - 保留模板格式替换文字"""
import sys,os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

src_orig = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_616.pptx'
dst_fixed = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx'

# 从原始文件重新开始（确保Logo等图片完好）
prs = Presentation(src_orig)

def safe_color(run):
    """安全获取run的字体颜色"""
    try:
        if run.font.color and run.font.color.rgb:
            return run.font.color.rgb
    except:
        pass
    return RGBColor(0xFF, 0xFF, 0xFF)

# ================================================================
# Slide 2 (index 1): 培训议程
# ================================================================
for shape in prs.slides[1].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('1号文·2号文·', '')
                run.text = run.text.replace('国务院核心文件解读', '国务院及成都市政策解读')
print('1. Slide 2 - 议程 updated')

# ================================================================
# Slide 3 (index 2): 第一部分总览
# ================================================================
for shape in prs.slides[2].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '1号文 → 2号文' in run.text:
                    run.text = '15号文 → 46号令 → 成都市配套制度：国家顶层设计 × 地方落地执行'
print('2. Slide 3 - 总览 updated')

# ================================================================
# Slide 4 (index 3): 1号文 → SW-2026-1196
# ================================================================
slide4 = prs.slides[3]
title_shape = badge_shape = body_shape = None
for s in slide4.shapes:
    if s.name == 'TextBox 5':
        title_shape = s
    elif s.name == 'Rectangle 6':
        badge_shape = s
    elif s.name == 'TextBox 7':
        body_shape = s

if title_shape:
    para = title_shape.text_frame.paragraphs[0]
    para.text = ''
    run = para.add_run()
    run.text = 'SW-2026-1196：成都市属国企违规经营投资责任追究实施办法'
    run.font.name = '微软雅黑'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

if badge_shape:
    para = badge_shape.text_frame.paragraphs[0]
    para.text = ''
    run = para.add_run()
    run.text = '市级配套'
    run.font.name = '微软雅黑'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

if body_shape:
    lines = [
        '▎追责范围 — 13大类全覆盖',
        '集团管控 · 风险管理 · 购销管理 · 工程承包 · 金融业务 · 科技创新',
        '资金管理 · 担保活动 · 产权管理 · 固定资产投资 · 股权投资 · 资产管理 · 境外投资',
        '',
        '▎追责方式（可合并使用）',
        '批评诫勉 / 组织处理（停职·免职·降职）/ 扣减薪酬 / 禁入限制(5年~终身)',
        '纪律处分 / 移送纪检监察或司法机关',
        '',
        '▎损失认定三档',
        '一般损失 <100万  |  较大损失 100万~1000万  |  重大损失 >1000万',
        '',
        '▎三条硬杠杠',
        '🔴 重大决策终身问责 — 退休≠安全，离职≠免责',
        '🟢 集体决策中明确反对意见 → 可免予责任',
        '⚠ 损失未达标准但导致企业无法持续经营 → 视为重大损失',
        '',
        '▶ 与你的关系',
        '停车场收费不清、商户租金拖欠、物业用房违规——天府广场13项问题全在追责范围内',
    ]
    tf = body_shape.text_frame
    # 确保段落数量足够
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, line in enumerate(lines):
        para = tf.paragraphs[i]
        para.text = ''
        run = para.add_run()
        run.text = line
        run.font.name = '微软雅黑'
        if line.startswith('▎') or line.startswith('▶') or line.startswith('🔴') or line.startswith('🟢') or line.startswith('⚠'):
            run.font.bold = True
            run.font.size = Pt(14)
        else:
            run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

print('3. Slide 4 - SW-2026-1196')

# ================================================================
# Slide 5 (index 4): 2号文 → 成国资发〔2025〕15号
# ================================================================
slide5 = prs.slides[4]
t5 = b5 = body5 = None
for s in slide5.shapes:
    if s.name == 'TextBox 5':
        t5 = s
    elif s.name == 'Rectangle 6':
        b5 = s
    elif s.name == 'TextBox 7':
        body5 = s

if t5:
    para = t5.text_frame.paragraphs[0]
    para.text = ''
    run = para.add_run()
    run.text = '成国资发〔2025〕15号：市属国企资产租赁管理制度'
    run.font.name = '微软雅黑'
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

if b5:
    para = b5.text_frame.paragraphs[0]
    para.text = ''
    run = para.add_run()
    run.text = '业务直击'
    run.font.name = '微软雅黑'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

if body5:
    lines5 = [
        '▎核心要求：公开招租',
        '必须在市政府网+市国资委网+产权交易所+集团官网同步发布',
        '正式披露≥10个工作日 | 结果公示≥3个工作日',
        '',
        '▎租金定价机制',
        '以评估/估价/询价/公开市场价为依据，综合市场需求、区位布局、功能业态确定',
        '',
        '▎五条红线',
        '❌ 拆分资产规避公开招租   ❌ 未经批准擅自出租',
        '❌ 擅自改变用途或租赁期限   ❌ 截留/挪用/私分租金收入',
        '❌ 违规转租 · 有违约行为不得续租 · 短期租赁到期不得续租',
        '',
        '▎监管机制',
        '"三重一大"集体决策 → 报市国资委备案 → 年度自查自纠 → 审计+专项检查',
        '重点查：未公开租赁 · 租金底价偏高 · 超长租期 · 违规转租',
        '',
        '▶ 与你的关系',
        '天府广场商业租赁的合同续租、商户转租、租金定价——全部在这个文件里画了红线',
    ]
    tf = body5.text_frame
    while len(tf.paragraphs) < len(lines5):
        tf.add_paragraph()
    for i, line in enumerate(lines5):
        para = tf.paragraphs[i]
        para.text = ''
        run = para.add_run()
        run.text = line
        run.font.name = '微软雅黑'
        if line.startswith('▎') or line.startswith('▶') or line.startswith('❌') or line.startswith('"'):
            run.font.bold = True
            run.font.size = Pt(14)
        else:
            run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

print('4. Slide 5 - 成国资发15号')

# ================================================================
# Slide 8 (index 7): 政策全景图
# ================================================================
slide8 = prs.slides[7]
replacements = {
    '1号文': '1196号\n(追责)',
    '看见': '定位',
    '数智化系统让异常无处遁形': '13大类追责全覆盖 | 终身追责不设时限',
    '2号文': '15号\n(租赁)',
    '看准': '规范',
    '穿透监管框架界定查什么、查多深': '公开招租 | 定价机制 | 全流程规范',
    '四位一体：政策全景图': '六位一体：政策全景图（国家+市级）',
    '技术 ──→ 制度 ──→ 执行 ──→ 问责  =  完整闭环': '15号文→46号令→1196号→15号  =  国家·市级双层闭环',
}
for shape in slide8.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for old, new in replacements.items():
                if old in para.text:
                    for run in para.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
print('5. Slide 8 - 全景图 updated')

# ================================================================
# Slide 27 (index 26): 关键启示
# ================================================================
for shape in prs.slides[26].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '1号文 → 用系统减少' in para.text:
                for run in para.runs:
                    if '1号文' in run.text:
                        run.text = run.text.replace(
                            '1号文 → 用系统减少"人"的随意性       2号文 → 用框架消除"管"的盲区',
                            '1196号文 → 13大类追责全覆盖, 终身问责不设限    15号租赁文 → 公开招租五条红线, 画清合规边界'
                        )
print('6. Slide 27 - 关键启示 updated')

# ================================================================
# Slide 28 (index 27): 合规自查要点
# ================================================================
swap28 = {
    '对照1号文': '对照1196号文',
    '信息系统打通了吗？预警自动了吗？': '追责红线画清了吗？投融资合规吗？',
    '对照2号文': '对照15号租赁文',
    '穿透到底了吗？资金可追溯了吗？': '公开招租执行了吗？转租合规了吗？',
}
for shape in prs.slides[27].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for old, new in swap28.items():
                if old == para.text.strip():
                    for run in para.runs:
                        if old in run.text:
                            run.text = new
print('7. Slide 28 - 自查要点 updated')

# ================================================================
# 保存
# ================================================================
prs.save(dst_fixed)
print(f'\nDONE: {dst_fixed}')
print(f'Size: {os.path.getsize(dst_fixed)//1024}KB')
