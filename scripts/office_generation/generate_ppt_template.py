"""
审计/咨询汇报PPT标准模板生成器
基于中银证券研报分析，统一配色与排版规范
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

# ===== 配色方案 =====
C = {
    'primary':        RGBColor(0x1A, 0x3A, 0x6E),
    'primary_light':  RGBColor(0xD6, 0xE0, 0xF0),
    'primary_pale':   RGBColor(0xEE, 0xF2, 0xF9),
    'secondary':      RGBColor(0xE8, 0x6A, 0x17),
    'text':           RGBColor(0x33, 0x33, 0x33),
    'text_light':     RGBColor(0x88, 0x88, 0x88),
    'white':          RGBColor(0xFF, 0xFF, 0xFF),
    'bg_dark':        RGBColor(0x0F, 0x24, 0x48),
    'bg_slide':       RGBColor(0xF7, 0xF8, 0xFA),
    'accent_blue':    RGBColor(0x34, 0x98, 0xDB),
    'accent_green':   RGBColor(0x27, 0xAE, 0x60),
    'risk_high':      RGBColor(0xCC, 0x33, 0x33),
    'risk_medium':    RGBColor(0xE8, 0x6A, 0x17),
    'risk_low':       RGBColor(0x33, 0x99, 0x33),
}

def hex_from_rgb(c):
    return f'{c[0]:02x}{c[1]:02x}{c[2]:02x}'

def add_slide_number(slide, num, total):
    """页脚页码"""
    left = Inches(8.6)
    top = Inches(7.1)
    width = Inches(1.2)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f'{num} / {total}'
    run.font.size = Pt(8)
    run.font.color.rgb = C['text_light']

def add_slide_title(slide, title, subtitle=None):
    """在给定slide上添加标题内容"""
    # 顶部色带
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C['primary']
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.color.rgb = C['primary']
    run.font.bold = True

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.85), Inches(8.8), Inches(0.4))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(13)
        run2.font.color.rgb = C['text_light']

def add_body_text(slide, text, left=0.6, top=1.5, width=8.8, size=14, color=None, bold=False, align=PP_ALIGN.LEFT):
    """添加正文文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color or C['text']
    run.font.bold = bold
    return tf

def add_bullet_list(slide, items, left=0.8, top=1.6, width=8.4, step=0.4):
    """添加要点列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(step * len(items) + 0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # 项目符号
        run_bullet = p.add_run()
        run_bullet.text = '\u25b8  '
        run_bullet.font.size = Pt(13)
        run_bullet.font.color.rgb = C['primary']
        run_bullet.font.bold = True
        # 内容
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(12)
        run_text.font.color.rgb = C['text']
        p.space_after = Pt(6)
    return tf

def add_table(slide, headers, rows, left=0.6, top=1.5, col_widths=None):
    """添加表格"""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    tbl_width = Inches(8.8)
    tbl_height = Inches(0.4 * n_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), tbl_width, tbl_height)
    table = table_shape.table

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = C['white']
                run.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = C['primary']

    # 数据行
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = C['text']
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C['primary_pale']

    # 列宽
    if col_widths:
        for i, w in enumerate(col_widths):
            if i < n_cols:
                table.columns[i].width = Inches(w)

    return table

def add_source_note(slide, text, top=6.8):
    """数据来源脚注"""
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'数据来源：{text}'
    run.font.size = Pt(8)
    run.font.color.rgb = C['text_light']
    run.font.italic = True

def add_kpi_box(slide, label, value, x, y, color=None):
    """KPI指标卡片"""
    w = Inches(2.6)
    h = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C['primary_pale']
    shape.line.fill.background()

    # 数值
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_val = p.add_run()
    run_val.text = str(value)
    run_val.font.size = Pt(28)
    run_val.font.color.rgb = color or C['primary']
    run_val.font.bold = True

    # 标签
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run_label = p2.add_run()
    run_label.text = label
    run_label.font.size = Pt(10)
    run_label.font.color.rgb = C['text_light']

def generate_ppt_template():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== 幻灯片母版背景 =====
    blank_layout = prs.slide_layouts[6]  # blank layout

    TOTAL = 12

    # ===== 1. 封面页 =====
    slide = prs.slides.add_slide(blank_layout)
    # 深色背景块
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(4.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = C['bg_dark']
    bg_shape.line.fill.background()

    # 标题
    add_body_text(slide, '审计 / 咨询报告', left=0.6, top=0.6, size=14, color=C['text_light'])
    tf_main = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(2.0))
    tf_main.text_frame.word_wrap = True
    p_main = tf_main.text_frame.paragraphs[0]
    run_main = p_main.add_run()
    run_main.text = '[请输入报告标题]'
    run_main.font.size = Pt(36)
    run_main.font.color.rgb = C['white']
    run_main.font.bold = True

    # 底部信息
    add_body_text(slide, f'四川融策会计师事务所 / 四川融策工程咨询公司', left=0.6, top=5.0, size=12, color=C['text_light'])
    add_body_text(slide, datetime.date.today().strftime('%Y年%m月'), left=0.6, top=5.4, size=11, color=C['text_light'])
    add_body_text(slide, '机密 · 仅限内部使用', left=0.6, top=6.8, size=9, color=C['text_light'])
    add_slide_number(slide, 1, TOTAL)

    # ===== 2. 目录页 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '目  录', 'CONTENTS')
    toc_items = [
        '项目概述与目标',
        '项目总体情况分析',
        '重点问题深度剖析',
        '问题发现汇总',
        '改进建议与方案',
        '结论与总体评价',
    ]
    for i, item in enumerate(toc_items):
        y = 1.6 + i * 0.8
        # 编号圆
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = C['primary']
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = f'0{i+1}'
        cr.font.size = Pt(13)
        cr.font.color.rgb = C['white']
        cr.font.bold = True

        add_body_text(slide, item, left=1.7, top=y, size=15, color=C['text'])
    add_slide_number(slide, 2, TOTAL)

    # ===== 3. 项目概述 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '一、项目概述', 'Project Overview')
    add_bullet_list(slide, [
        '委托单位：[请输入委托单位全称]',
        '被审计单位：[请输入被审计单位全称]',
        '项目类型：绩效评价 / 资产清查 / 专项债审计 / 工程结算 / 其他',
        '实施周期：[开始日期] 至 [结束日期]，历时 [天数] 天',
        '投入力量：累计 [人数] 人次，查阅凭证[数量]份、合同[数量]份',
    ], top=1.5, step=0.45)
    add_slide_number(slide, 3, TOTAL)

    # ===== 4. 核心KPI指标 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '一、核心指标概览', 'Key Indicators')

    add_kpi_box(slide, '项目总预算', '2,000万元', 0.6, 1.5, C['primary'])
    add_kpi_box(slide, '实际支出', '1,645.67万元', 3.5, 1.5, C['accent_blue'])
    add_kpi_box(slide, '总体执行率', '82.3%', 6.4, 1.5, C['accent_green'])
    add_kpi_box(slide, '发现问题', '19项', 0.6, 3.0, C['risk_high'])
    add_kpi_box(slide, '高风险问题', '3项', 3.5, 3.0, C['risk_high'])
    add_kpi_box(slide, '中低风险问题', '16项', 6.4, 3.0, C['risk_medium'])

    add_source_note(slide, '根据项目财务台账及银行对账单汇总，统计截至202X年X月X日', top=4.8)
    add_slide_number(slide, 4, TOTAL)

    # ===== 5. 预算执行对比 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '二、预算执行情况对比', 'Budget Execution Comparison')
    budget_data = [
        ('项目A', '1,000.00', '856.32', '85.6%', '正常'),
        ('项目B', '500.00', '423.15', '84.6%', '正常'),
        ('项目C', '300.00', '298.70', '99.6%', '正常'),
        ('项目D', '200.00', '67.50', '33.8%', '⚠ 偏低'),
        ('合计', '2,000.00', '1,645.67', '82.3%', '—'),
    ]
    add_table(slide,
        ['项目名称', '预算金额(万元)', '实际支出(万元)', '执行率', '状态'],
        budget_data,
        left=0.6, top=1.4, col_widths=[2.0, 1.8, 1.8, 1.5, 1.7]
    )
    add_source_note(slide, '[数据来源说明]', top=3.6)

    # 说明文字
    add_body_text(slide, '关键发现：', left=0.6, top=4.1, size=12, color=C['primary'], bold=True)
    add_bullet_list(slide, [
        '项目D预算执行率仅33.8%，存在重大执行偏差',
        '项目C执行率接近100%，需关注是否存在年末突击花钱',
        '总体执行率82.3%，处于合理区间',
    ], left=0.6, top=4.6, step=0.35)
    add_slide_number(slide, 5, TOTAL)

    # ===== 6. 资金构成（饼图说明） =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '二、项目资金构成分布', 'Fund Distribution')
    add_body_text(slide, '[此处插入饼图/环形图]', left=0.6, top=1.5, size=13, color=C['text_light'], align=PP_ALIGN.LEFT)
    add_body_text(slide, '图表配色建议：深蓝系渐变调色盘 #1a3a6e / #2c5f9e / #4a90c4 / #7ab4d6 / #a8d0e6', left=0.6, top=2.0, size=10, color=C['text_light'])
    add_body_text(slide, '饼图扇区不超过5个，标注占比百分比，核心扇区略突出', left=0.6, top=2.4, size=10, color=C['text_light'])
    add_slide_number(slide, 6, TOTAL)

    # ===== 7. 预算执行趋势（柱状图） =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '二、月度预算执行进度趋势', 'Monthly Execution Trend')
    add_body_text(slide, '[此处插入柱状图+折线图组合]', left=0.6, top=1.5, size=13, color=C['text_light'])
    add_body_text(slide, '图表配色建议：柱状图 月度支出额 浅蓝色 #3498DB，折线图 累计执行率 深蓝色 #1a3a6e加粗，计划线 红色虚线 #CC3333', left=0.6, top=2.0, size=10, color=C['text_light'])
    add_slide_number(slide, 7, TOTAL)

    # ===== 8. 重点问题分析 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '三、重点问题分析', 'Key Issues Analysis')
    issues = [
        ('预算执行', '项目D执行率仅33.8%，前期论证不充分，立项条件不成熟', '高'),
        ('资金管理', '专项资金存在短期挪用，内控机制存在漏洞', '中'),
        ('项目管理', '项目A未经批准擅自调整建设内容，涉及85万元', '高'),
        ('合同规范', '部分合同签订不规范，要素缺失、条款模糊', '中'),
        ('档案管理', '过程性文件缺失，档案管理制度不健全', '低'),
    ]
    for i, (cat, desc, risk) in enumerate(issues):
        y = 1.5 + i * 1.05
        # 分类标签
        tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(1.8), Inches(0.45)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = C['primary']
        tag.line.fill.background()
        ttf = tag.text_frame
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = cat
        tr.font.size = Pt(11)
        tr.font.color.rgb = C['white']

        # 风险标签
        risk_color = {'高': C['risk_high'], '中': C['risk_medium'], '低': C['risk_low']}
        risk_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.55), Inches(y), Inches(0.6), Inches(0.45)
        )
        risk_shape.fill.solid()
        risk_shape.fill.fore_color.rgb = risk_color.get(risk, C['text_light'])
        risk_shape.line.fill.background()
        rtf = risk_shape.text_frame
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.CENTER
        rr = rp.add_run()
        rr.text = f'{risk}风险'
        rr.font.size = Pt(9)
        rr.font.color.rgb = C['white']
        rr.font.bold = True

        # 描述
        add_body_text(slide, desc, left=3.3, top=y + 0.05, size=11, color=C['text'])

    add_slide_number(slide, 8, TOTAL)

    # ===== 9. 问题汇总统计 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '四、问题发现汇总', 'Findings Summary')

    # 风险统计表
    risk_data = [
        ('高风险', '3', '15.8%', '需立即整改'),
        ('中风险', '8', '42.1%', '限期整改'),
        ('低风险', '8', '42.1%', '持续关注'),
        ('合计', '19', '100%', '—'),
    ]
    add_table(slide,
        ['风险等级', '问题数量', '占比', '处理建议'],
        risk_data,
        left=0.6, top=1.4, col_widths=[2.2, 2.2, 2.2, 2.2]
    )
    add_slide_number(slide, 9, TOTAL)

    # ===== 10. 问题清单 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '四、重点问题清单（部分）', 'Key Finding Details')
    finding_data = [
        ('F-001', '项目D预算执行率仅33.8%', '132.50', '高'),
        ('F-002', '部分支出未取得合规发票', '45.30', '高'),
        ('F-003', '项目A擅自调整建设内容', '85.00', '高'),
        ('F-004', '合同签订不规范', '—', '中'),
        ('F-005', '专项资金短期挪用', '120.00', '中'),
    ]
    add_table(slide,
        ['编号', '问题描述', '涉及金额(万元)', '风险等级'],
        finding_data,
        left=0.6, top=1.4, col_widths=[1.5, 4.0, 2.0, 1.3]
    )
    add_source_note(slide, '以上为部分重点问题摘录，完整清单详见正式报告')
    add_slide_number(slide, 10, TOTAL)

    # ===== 11. 改进建议 =====
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, '五、改进建议', 'Recommendations')
    suggestions = [
        ('加强前期论证', '增加可行性研究深度，引入专家评审，从源头杜绝"先天不足"项目'),
        ('完善资金管理', '建立专项资金全流程监控，专户管理、专款专用'),
        ('规范合同变更', '统一合同模板，建立变更审批层级管理制度'),
        ('强化绩效闭环', '目标设定→过程监控→结果评价→反馈应用全链条管理'),
    ]
    for i, (title, detail) in enumerate(suggestions):
        y = 1.5 + i * 1.3
        # 编号
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(0.5), Inches(0.5)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = C['primary']
        num_box.line.fill.background()
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        nr = np.add_run()
        nr.text = str(i + 1)
        nr.font.size = Pt(16)
        nr.font.color.rgb = C['white']
        nr.font.bold = True

        # 标题
        add_body_text(slide, title, left=1.3, top=y + 0.05, size=14, color=C['primary'], bold=True)
        # 详情
        add_body_text(slide, detail, left=1.3, top=y + 0.45, size=11, color=C['text'])

    add_slide_number(slide, 11, TOTAL)

    # ===== 12. 结束页 =====
    slide = prs.slides.add_slide(blank_layout)
    bg_end = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg_end.fill.solid()
    bg_end.fill.fore_color.rgb = C['bg_dark']
    bg_end.line.fill.background()

    add_body_text(slide, '感谢聆听', left=0.6, top=2.2, width=8.8, size=40, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
    add_body_text(slide, 'THANK YOU', left=0.6, top=3.0, width=8.8, size=16, color=C['text_light'], align=PP_ALIGN.CENTER)
    add_body_text(slide, '四川融策会计师事务所 / 四川融策工程咨询公司', left=0.6, top=4.5, width=8.8, size=13, color=C['text_light'], align=PP_ALIGN.CENTER)
    add_body_text(slide, '地址：[公司地址]  |  电话：[联系电话]  |  邮箱：[联系邮箱]', left=0.6, top=6.5, width=8.8, size=10, color=C['text_light'], align=PP_ALIGN.CENTER)
    add_slide_number(slide, 12, TOTAL)

    # ===== 保存 =====
    output_path = r'D:\openclaw-workspace\审计汇报PPT标准模板.pptx'
    prs.save(output_path)
    print(f'[OK] PPT模板已生成：{output_path}')
    return output_path

if __name__ == '__main__':
    generate_ppt_template()
