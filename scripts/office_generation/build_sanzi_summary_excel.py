import re
from pathlib import Path
from collections import Counter

import pdfplumber
from docx import Document
from pptx import Presentation
import xlsxwriter

BASE = Path(r"C:\Users\scrccpa\Desktop\三资三化课件")
PROJECT = Path(r"D:\openclaw-workspace\projects\西昌三资三化投标")
TEMP_SANZI = Path(r"D:\openclaw-workspace\temp\三资三化")
TEMP_ASSET = Path(r"D:\openclaw-workspace\temp\国有资产")
OUT = BASE / "三资三化资料清单+可引用观点摘要.xlsx"

KEYWORDS = [
    "三资三化", "三资", "资源资产化", "资产资本化", "资本杠杆化", "资产证券化", "存量资产", "盘活",
    "国有资本", "国有资产", "国有资源", "经营权", "有偿配置", "确权", "确值", "确管", "收益统筹",
    "全口径预算", "大财政", "债务化解", "融资", "市场化", "风险", "国企", "平台公司",
]


def clean(text):
    return re.sub(r"\s+", " ", text or "").strip()


def extract_pdf(path):
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages[:80]:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                pass
    return clean("\n".join(parts))


def extract_docx(path):
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if vals:
                paras.append(" | ".join(vals))
    return clean("\n".join(paras))


def extract_pptx(path):
    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            chunks.append(f"第{i}页：" + "；".join(texts))
    return clean("\n".join(chunks))


def extract_md(path):
    return clean(path.read_text(encoding="utf-8", errors="ignore"))


def split_sentences(text):
    text = re.sub(r"([。！？；])", r"\1\n", text)
    parts = [clean(x) for x in text.splitlines()]
    return [p for p in parts if len(p) >= 16]


def score_sentence(s):
    score = 0
    for kw in KEYWORDS:
        if kw in s:
            score += 3 if kw in ("三资三化", "资源资产化", "资产资本化", "资本杠杆化", "存量资产", "经营权") else 1
    if re.search(r"\d+(\.\d+)?\s*(万亿|亿元|万元|%|个|项|年|类|批)", s):
        score += 2
    if any(x in s for x in ["一是", "二是", "三是", "应当", "必须", "重点", "核心", "路径", "风险", "建议"]):
        score += 1
    return score


def top_sentences(text, n=6):
    ss = split_sentences(text)
    ranked = sorted(ss, key=lambda x: (score_sentence(x), len(x)), reverse=True)
    picked = []
    for s in ranked:
        if len(s) > 180:
            s = s[:177] + "…"
        if all(s[:40] not in p for p in picked):
            picked.append(s)
        if len(picked) >= n:
            break
    return picked


def infer_theme(name, text):
    rules = [
        ("专项行动实施", ["实施方案", "专项行动", "工作目标", "组织保障"]),
        ("存量资产盘活", ["存量资产", "盘活", "低效", "闲置"]),
        ("国有资本运营", ["国有资本", "资本运营", "国有企业", "资本"]),
        ("债务化解与融资", ["债务", "融资", "杠杆", "证券化"]),
        ("资源资产化路径", ["资源资产化", "确权", "确值", "经营权"]),
        ("平台与数字化管理", ["平台", "数据库", "信息化", "监管"]),
    ]
    blob = name + text[:5000]
    hits = []
    for theme, kws in rules:
        if any(k in blob for k in kws):
            hits.append(theme)
    return "、".join(hits[:3]) or "综合资料"


def abstract(text, name):
    ts = top_sentences(text, 5)
    if ts:
        return "；".join(ts[:3])
    return f"{name}：文本可提取内容较少，建议作为原始附件查阅。"


def key_terms(text):
    counts = Counter()
    for kw in KEYWORDS:
        c = text.count(kw)
        if c:
            counts[kw] = c
    return "、".join([k for k, _ in counts.most_common(8)])


def file_type(path):
    return path.suffix.lower().lstrip(".").upper()


sources = []
for folder in [BASE, PROJECT, TEMP_SANZI, TEMP_ASSET]:
    if folder.exists():
        for p in folder.rglob("*"):
            if not p.is_file() or p.name.startswith("~$"):
                continue
            if p.suffix.lower() not in [".pdf", ".pptx", ".docx", ".md"]:
                continue
            if p == OUT:
                continue
            if folder == PROJECT:
                keep = {"服务方案-正式版.md", "三篇文献深度分析.md", "国有资产文献深度分析.md", "网上查询成果.md", "背景分析与素材需求.md", "西昌三资三化投标服务方案-完整版-V8f.docx"}
                if p.name not in keep:
                    continue
            sources.append(p)

records = []
for p in sources:
    try:
        if p.suffix.lower() == ".pdf":
            text = extract_pdf(p)
        elif p.suffix.lower() == ".docx":
            text = extract_docx(p)
        elif p.suffix.lower() == ".pptx":
            text = extract_pptx(p)
        else:
            text = extract_md(p)
        err = ""
    except Exception as e:
        text = ""
        err = str(e)
    records.append({
        "path": str(p),
        "name": p.name,
        "type": file_type(p),
        "source_group": "最新讲义课件" if BASE in p.parents or p.parent == BASE else ("既有项目成果" if PROJECT in p.parents or p.parent == PROJECT else "原始参考文献"),
        "theme": infer_theme(p.name, text),
        "chars": len(text),
        "abstract": abstract(text, p.name) if text else f"未能提取正文：{err}",
        "terms": key_terms(text),
        "quotes": top_sentences(text, 8),
    })

quote_rows = []
for rec in records:
    for i, q in enumerate(rec["quotes"], 1):
        quote_rows.append({
            "source": rec["name"],
            "theme": rec["theme"],
            "quote": q,
            "use": "可用于政策背景、实施路径、盘活模式、风险控制或案例支撑" if i <= 3 else "可作为补充论据或方案细化素材",
            "path": rec["path"],
        })

summary_points = [
    ["核心逻辑", "“三资三化”可概括为把沉睡的国有资源、资产、资金纳入统一清查、确权、确值、运营体系，通过资源资产化、资产资本化、资本杠杆化形成财政资源统筹能力。"],
    ["工作主线", "先起底清查形成底数，再分类确权确值，再分批制定盘活方案，最后通过经营权配置、租赁、转让、证券化、资本运营等方式实现收益。"],
    ["政策价值", "该模式服务“大财政”“大统筹”和全口径预算管理，有助于把依托行政权力、政府信用、国有资源资产取得的收入纳入预算统筹。"],
    ["项目抓手", "可围绕闲置低效资产资源、特许经营权、公共数据资源、行政事业单位资产、国企经营性资产等建立资产清单、收益清单、问题清单和项目清单。"],
    ["风险边界", "重点防止虚假盘活、重复融资、收益高估、权属不清、公益属性弱化、隐性债务新增、国有资产流失和程序合规瑕疵。"],
    ["咨询服务打法", "咨询成果应体现“政策研究—底数摸排—分类评价—方案设计—交易配置—风险控制—后续运营”全流程闭环。"],
]

workbook = xlsxwriter.Workbook(str(OUT))
workbook.set_properties({"title": "三资三化资料清单+可引用观点摘要", "author": "融策右护卫"})
fmt_title = workbook.add_format({"font_name": "微软雅黑", "font_size": 16, "bold": True, "font_color": "#0A1F3F", "align": "left", "valign": "vcenter"})
fmt_sub = workbook.add_format({"font_name": "微软雅黑", "font_size": 10, "italic": True, "font_color": "#666666", "valign": "vcenter"})
fmt_header = workbook.add_format({"font_name": "微软雅黑", "font_size": 11, "bold": True, "font_color": "#FFFFFF", "bg_color": "#0A1F3F", "border": 1, "border_color": "#C5955C", "align": "center", "valign": "vcenter", "text_wrap": True})
fmt_body = workbook.add_format({"font_name": "微软雅黑", "font_size": 10, "border": 1, "border_color": "#C5955C", "valign": "top", "text_wrap": True})
fmt_alt = workbook.add_format({"font_name": "微软雅黑", "font_size": 10, "bg_color": "#F5F2EC", "border": 1, "border_color": "#C5955C", "valign": "top", "text_wrap": True})
fmt_num = workbook.add_format({"font_name": "微软雅黑", "font_size": 10, "border": 1, "border_color": "#C5955C", "align": "center", "valign": "top"})
fmt_num_alt = workbook.add_format({"font_name": "微软雅黑", "font_size": 10, "bg_color": "#F5F2EC", "border": 1, "border_color": "#C5955C", "align": "center", "valign": "top"})


def add_table(sheet_name, title, subtitle, headers, rows, widths, row_height=96):
    ws = workbook.add_worksheet(sheet_name)
    ws.hide_gridlines(2)
    ws.merge_range(0, 0, 0, len(headers) - 1, title, fmt_title)
    ws.merge_range(1, 0, 1, len(headers) - 1, subtitle, fmt_sub)
    ws.set_row(0, 28)
    ws.set_row(1, 24)
    ws.set_row(2, 32)
    for i, width in enumerate(widths):
        ws.set_column(i, i, width)
    for col, header in enumerate(headers):
        ws.write(2, col, header, fmt_header)
    for r, row in enumerate(rows, 3):
        ws.set_row(r, row_height)
        alt = (r - 3) % 2 == 1
        for c, value in enumerate(row):
            fmt = (fmt_num_alt if alt else fmt_num) if c == 0 else (fmt_alt if alt else fmt_body)
            ws.write(r, c, value, fmt)
    ws.autofilter(2, 0, max(2, len(rows) + 2), len(headers) - 1)
    ws.freeze_panes(3, 0)
    return ws

rows1 = [[idx, rec["source_group"], rec["name"], rec["type"], rec["theme"], rec["terms"], rec["abstract"], rec["path"]] for idx, rec in enumerate(records, 1)]
add_table("资料清单", "三资三化资料清单", "来源：最新讲义课件 + 既有西昌三资三化项目资料 + 原始参考文献", ["序号", "资料类别", "资料名称", "格式", "主题归类", "关键词", "内容提炼", "文件位置"], rows1, [6, 14, 34, 8, 24, 28, 78, 80], 168)

rows2 = [[idx, q["theme"], q["quote"], q["use"], q["source"], q["path"]] for idx, q in enumerate(quote_rows, 1)]
add_table("可引用观点摘要", "可引用观点摘要", "已从课件、方案、文献中筛选可直接转化为方案表述的观点句", ["序号", "主题", "可引用观点/摘要句", "建议用途", "来源资料", "文件位置"], rows2, [6, 24, 86, 34, 34, 80], 90)

rows3 = [[i + 1, k, v] for i, (k, v) in enumerate(summary_points)]
add_table("课件提炼总结", "课件提炼总结", "面向投标方案、汇报材料和咨询服务方案的综合提炼", ["序号", "提炼维度", "核心总结"], rows3, [6, 22, 110], 86)

rows4 = [
    [1, "投标服务方案", "优先引用“全流程闭环、四张清单、分批盘活、风险控制、后续运营”表述，体现咨询服务不是写材料，而是协助形成可落地项目包。"],
    [2, "政策依据章节", "引用大财政、大统筹、全口径预算、国有资产收益纳入预算管理等内容，形成政策必要性。"],
    [3, "实施路径章节", "按照“清查摸底—确权确值—分类评价—盘活设计—交易实施—收益统筹—监督闭环”组织。"],
    [4, "案例支撑章节", "湖北21.5万亿元三资清理、江苏闲置资产盘活、无锡三资盘活专班等案例可作为经验借鉴。"],
    [5, "风险控制章节", "重点写权属、估值、交易、债务、收益、合规、公益属性七类风险，体现专业审慎。"],
    [6, "后续加工", "如需形成PPT或正式报告，可基于本表第二张“可引用观点摘要”筛选20条核心观点重组。"],
]
add_table("使用建议", "使用建议", "把资料转化为咨询方案、标书和汇报材料时的建议", ["序号", "应用场景", "建议"], rows4, [6, 24, 112], 76)

workbook.close()
print(str(OUT))
print(f"records={len(records)}, quotes={len(quote_rows)}")
