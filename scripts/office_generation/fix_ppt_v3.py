#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""修复PPT - 用替换run文本的方式（不clear段落，不丢失格式）"""
import sys,os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
from lxml import etree

src = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_616.pptx'
dst = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx'

prs = Presentation(src)

def replace_run_text_safe(para, old_text, new_text):
    """安全替换run文本 - 不修改段落结构"""
    for run in para.runs:
        if old_text in run.text:
            run.text = run.text.replace(old_text, new_text)
            return True
    return False

def replace_para_all(para, replacements):
    """替换段落中所有匹配的文本"""
    changed = False
    for run in para.runs:
        for old, new in replacements.items():
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    return changed

# ================================================================
# Slide 4 (index 3): 1号文 → 成国资发〔2026〕14号
# 策略: 只替换文本，不动段落结构，不clear
# ================================================================
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if not shape.has_text_frame:
        continue
    tf = shape.text_frame
    
    # 找到每个shape中第一个非空段落
    first_para = None
    for para in tf.paragraphs:
        if para.text.strip():
            first_para = para
            break
    
    if first_para is None:
        continue
    
    # 替换标题
    if 'TextBox 5' == shape.name or '1号文' in first_para.text:
        for run in first_para.runs:
            if '1号文' in run.text:
                run.text = '成国资发〔2026〕14号：《成都市属国企违规经营投资责任追究实施办法》'
    
    # 替换标签
    if 'Rectangle 6' == shape.name or '技术底座' == first_para.text.strip():
        for run in first_para.runs:
            if '技术底座' in run.text:
                run.text = '市级配套'

    # 替换正文
    if 'TextBox 7' == shape.name:
        new_lines = [
            '▎追责范围 — 13大类全覆盖',
            '集团管控 · 风险管理 · 购销管理 · 工程承包 · 金融业务 · 科技创新',
            '资金管理 · 担保活动 · 产权管理 · 固定资产投资 · 股权投资 · 资产管理 · 境外投资',
            '',
            '▎追责方式（可合并使用）',
            '批评诫勉 / 组织处理（停职·免职·降职）/ 扣减薪酬 / 禁入限制(5年~终身)',
            '纪律处分 / 移送纪检监察或司法机关',
            '',
            '▎损失认定三档',
            '一般损失 <100万  |  较大损失 100万~1000万  |  重大损失 >1000万',
            '',
            '▎三条硬杠杠',
            '🔴 重大决策终身问责 — 退休≠安全，离职≠免责',
            '🟢 集体决策中明确反对意见 → 可免予责任',
            '⚠ 损失未达标准但导致企业无法持续经营 → 视为重大损失',
            '',
            '▶ 与你的关系',
            '停车场收费不清、商户租金拖欠、物业用房违规——天府广场13项问题全在追责范围内',
        ]
        
        # 用第一个非空段落的run格式作为模板
        template_para = None
        template_run = None
        for p in tf.paragraphs:
            if p.runs and p.text.strip() and '▎' not in p.text:
                template_para = p
                template_run = p.runs[0]
                break
        if template_para is None:
            template_para = first_para
            template_run = first_para.runs[0] if first_para.runs else None
        
        # 逐行替换
        for i, line in enumerate(new_lines):
            if i < len(tf.paragraphs):
                para = tf.paragraphs[i]
            else:
                # 复制上一个段落的结构（用XML深拷贝）
                last_p = tf.paragraphs[-1]
                new_p = deepcopy(last_p._element)
                tf._txBody.append(new_p)
                para = tf.paragraphs[-1]
            
            # 替换现有run的文本
            if para.runs and i < len(tf.paragraphs):
                para.runs[0].text = line
            elif para.runs:
                for r in para.runs:
                    r.text = ''
                para.runs[0].text = line
            else:
                r = para.add_run()
                r.text = line
                if template_run:
                    r.font.name = template_run.font.name
                    r.font.size = template_run.font.size
                    try:
                        r.font.color.rgb = template_run.font.color.rgb
                    except:
                        r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        
        # 清除多余段落
        extra = len(tf.paragraphs) - len(new_lines)
        if extra > 0:
            for j in range(extra):
                p = tf.paragraphs[len(new_lines)]
                p.getparent().remove(p)

print('Slide 4 - 14号文 OK')

# ================================================================
# Slide 5 (index 4): 2号文 → 成国资发〔2025〕15号
# ================================================================
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if not shape.has_text_frame:
        continue
    tf = shape.text_frame
    first_para = None
    for para in tf.paragraphs:
        if para.text.strip():
            first_para = para
            break
    if first_para is None:
        continue
    
    if 'TextBox 5' == shape.name or '2号文' in first_para.text:
        for run in first_para.runs:
            if '2号文' in run.text:
                run.text = '成国资发〔2025〕15号：市属国企资产租赁管理制度'
    
    if 'Rectangle 6' == shape.name or '总纲领' == first_para.text.strip():
        for run in first_para.runs:
            if '总纲领' in run.text:
                run.text = '业务直击'

    if 'TextBox 7' == shape.name:
        new_lines5 = [
            '▎核心要求：公开招租',
            '必须在市政府网+市国资委网+产权交易所+集团官网同步发布',
            '正式披露≥10个工作日 | 结果公示≥3个工作日',
            '',
            '▎租金定价机制',
            '以评估/估价/询价/公开市场价为依据，综合市场需求、区位布局、功能业态确定',
            '',
            '▎五条红线',
            '❌ 拆分资产规避公开招租   ❌ 未经批准擅自出租',
            '❌ 擅自改变用途或租赁期限   ❌ 截留/挪用/私分租金收入',
            '❌ 违规转租 · 有违约行为不得续租 · 短期租赁到期不得续租',
            '',
            '▎监管机制',
            '"三重一大"集体决策 → 报市国资委备案 → 年度自查自纠 → 审计+专项检查',
            '',
            '▶ 与你的关系',
            '天府广场商业租赁的合同续租、商户转租、租金定价——全部在这个文件里画了红线',
        ]
        
        while len(tf.paragraphs) < len(new_lines5):
            last = tf.paragraphs[-1]
            new_p = deepcopy(last._element)
            tf._txBody.append(new_p)
        
        for i, line in enumerate(new_lines5):
            para = tf.paragraphs[i]
            if para.runs:
                para.runs[0].text = line
            else:
                r = para.add_run()
                r.text = line
        
        extra = len(tf.paragraphs) - len(new_lines5)
        if extra > 0:
            for j in range(extra):
                p = tf.paragraphs[len(new_lines5)]
                p.getparent().remove(p)

print('Slide 5 - 15号文 OK')

# ================================================================
# 其他页面更新
# ================================================================

# Slide 2
for shape in prs.slides[1].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('1号文·2号文·', '')
                run.text = run.text.replace('国务院核心文件解读', '国务院及成都市政策解读')
print('Slide 2 OK')

# Slide 3
for shape in prs.slides[2].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '1号文 → 2号文' in run.text:
                    run.text = '15号文 → 46号令 → 成都市配套制度：国家顶层设计 × 地方落地执行'
print('Slide 3 OK')

# Slide 8
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '1号文' == run.text.strip():
                    run.text = '14号(追责)'
                elif '2号文' == run.text.strip():
                    run.text = '15号(租赁)'
                elif '看见' == run.text.strip():
                    run.text = '定位'
                elif '看准' == run.text.strip():
                    run.text = '规范'
                elif '数智化系统让异常无处遁形' == run.text.strip():
                    run.text = '13大类追责全覆盖 | 终身追责'
                elif '穿透监管框架界定查什么、查多深' == run.text.strip():
                    run.text = '公开招租 | 五条红线 | 全流程规范'
                elif '四位一体：政策全景图' in run.text:
                    run.text = run.text.replace('四位一体：政策全景图', '六位一体：政策全景图（国家+市级）')
                elif '技术 ──→ 制度 ──→ 执行 ──→ 问责  =  完整闭环' in run.text:
                    run.text = '15号文→46号令→14号→15号  =  国家·市级双层闭环'
print('Slide 8 OK')

# Slide 27
for shape in prs.slides[26].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '1号文 → 用系统减少' in run.text:
                    run.text = run.text.replace(
                        '1号文 → 用系统减少"人"的随意性       2号文 → 用框架消除"管"的盲区',
                        '14号文 → 13大类追责全覆盖, 终身问责不设限    15号文 → 公开招租五条红线, 画清合规边界'
                    )
print('Slide 27 OK')

# Slide 28
for shape in prs.slides[27].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '对照1号文' == run.text.strip():
                    run.text = '对照14号文'
                elif '信息系统打通了吗？预警自动了吗？' == run.text.strip():
                    run.text = '追责红线画清了吗？投融资合规吗？'
                elif '对照2号文' == run.text.strip():
                    run.text = '对照15号租赁文'
                elif '穿透到底了吗？资金可追溯了吗？' == run.text.strip():
                    run.text = '公开招租执行了吗？转租合规了吗？'
print('Slide 28 OK')

# 保存
prs.save(dst)
print(f'\nDone: {dst}')
print(f'Size: {os.path.getsize(dst)//1024}KB')
