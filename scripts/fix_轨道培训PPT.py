#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""修复PPT：第4/5页空白 + 第27/28页替换 + Logo恢复"""
import sys,os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy

# 从原始文件重建4/5页（保留模板自带的格式），同时修改其他页
src_orig = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_616.pptx'
dst_new = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_更新版.pptx'
dst_fixed = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx'

prs = Presentation(dst_new)
prs_orig = Presentation(src_orig)

def safe_replace_runs(shape, old_new_pairs):
    """安全替换run文本，保留原有格式"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in old_new_pairs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

def replace_textbox_content(shape, new_lines):
    """替换文本框内容，保留每行原有run的格式"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 保留原有段落结构，只替换文字
    for i, para in enumerate(tf.paragraphs):
        if i < len(new_lines):
            # 保存第一个run的格式
            if para.runs:
                first_run = para.runs[0]
                # 用新的run替代旧的
                para.clear()
                new_run = para.add_run()
                new_run.text = new_lines[i]
                # 保持字体设置
                new_run.font.name = '微软雅黑'
                new_run.font.size = Pt(12)
                new_run.font.color.rgb = first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else RGBColor(0xFF,0xFF,0xFF)
        else:
            if para.runs:
                para.clear()

# ================================================================
# 方案：直接从原始PPT复制第4/5页的XML结构，再"填充"新文字
# 这样保留所有原始格式（字体颜色、大小、模板样式）
# ================================================================

# 1. 重新处理第4页（slide index 3）
slide4 = prs.slides[3]
text_shapes = [s for s in slide4.shapes if s.has_text_frame]

# 获取原始第4页的文本样式作为参考
orig_slide4 = prs_orig.slides[3]
orig_text_shapes = [s for s in orig_slide4.shapes if s.has_text_frame]

# 方案：保留slide4的shape结构，逐run替换文本
# TextBox 5 (title): 替换为SW-2026-1196标题
title_shape = None
badge_shape = None
body_shape = None
for s in text_shapes:
    if s.name == 'TextBox 5':
        title_shape = s
    elif s.name == 'Rectangle 6':
        badge_shape = s
    elif s.name == 'TextBox 7':
        body_shape = s

# 用原始文件来参考格式 — 获取原始文件同位置shape的run格式
if title_shape:
    # 获取原始标题的run
    orig_title_shape = None
    for s in orig_text_shapes:
        if s.name == 'TextBox 5':
            orig_title_shape = s
            break
    
    # 替换标题
    if orig_title_shape and orig_title_shape.text_frame.paragraphs[0].runs:
        orig_run = orig_title_shape.text_frame.paragraphs[0].runs[0]
        para = title_shape.text_frame.paragraphs[0]
        para.clear()
        new_run = para.add_run()
        new_run.text = 'SW-2026-1196：成都市属国企违规经营投资责任追究实施办法'
        new_run.font.name = orig_run.font.name
        new_run.font.size = orig_run.font.size
        new_run.font.bold = orig_run.font.bold
        try:
                    new_run.font.color.rgb = orig_run.font.color.rgb
                except:
                    pass
        print('Title format preserved from original')

if badge_shape:
    orig_badge_shape = None
    for s in orig_text_shapes:
        if s.name == 'Rectangle 6':
            orig_badge_shape = s
            break
    if orig_badge_shape and orig_badge_shape.text_frame.paragraphs[0].runs:
        orig_run = orig_badge_shape.text_frame.paragraphs[0].runs[0]
        para = badge_shape.text_frame.paragraphs[0]
        para.clear()
        new_run = para.add_run()
        new_run.text = '市级配套'
        new_run.font.name = orig_run.font.name
        new_run.font.size = orig_run.font.size
        new_run.font.bold = orig_run.font.bold
        try:
                    new_run.font.color.rgb = orig_run.font.color.rgb
                except:
                    pass
        print('Badge format preserved')

if body_shape:
    orig_body_shape = None
    for s in orig_text_shapes:
        if s.name == 'TextBox 7':
            orig_body_shape = s
            break
    
    body_lines = [
        '▎追责范围 — 13大类全覆盖',
        '集团管控 · 风险管理 · 购销管理 · 工程承包 · 金融业务 · 科技创新',
        '资金管理 · 担保活动 · 产权管理 · 固定资产投资 · 股权投资 · 资产管理 · 境外投资',
        '',
        '▎追责方式（可合并使用）',
        '批评诫勉 / 组织处理（停职·免职·降职）/ 扣减薪酬 / 禁入限制(5年~终身)',
        '纪律处分 / 移送纪检监察或司法机关',
        '',
        '▎损失认定三档',
        '一般损失 <100万  |  较大损失 100万~1000万  |  重大损失 >1000万',
        '',
        '▎三条硬杠杠',
        '🔴 重大决策终身问责 — 退休≠安全，离职≠免责',
        '🟢 集体决策中明确反对意见 → 可免予责任',
        '⚠ 损失未达标准但导致企业无法持续经营 → 视为重大损失',
        '',
        '▶ 与你的关系',
        '停车场收费不清、商户租金拖欠、物业用房违规——天府广场13项问题全在追责范围内',
    ]
    
    if orig_body_shape and orig_body_shape.text_frame.paragraphs[0].runs:
        orig_run = orig_body_shape.text_frame.paragraphs[0].runs[0]
        tf = body_shape.text_frame
        
        # 保留第一段，清空后重写
        for pi in range(len(tf.paragraphs)):
            if tf.paragraphs[pi].runs:
                tf.paragraphs[pi].clear()
        
        # 确保有足够的段落
        while len(tf.paragraphs) < len(body_lines):
            tf.add_paragraph()
        
        for i, line in enumerate(body_lines):
            para = tf.paragraphs[i]
            para.clear()
            new_run = para.add_run()
            new_run.text = line
            new_run.font.name = orig_run.font.name or '微软雅黑'
            new_run.font.size = orig_run.font.size or Pt(13)
            if line.startswith('▎') or line.startswith('▶') or line.startswith('🔴') or line.startswith('⚠'):
                new_run.font.bold = True
                new_run.font.size = orig_run.font.size or Pt(14)
            else:
                new_run.font.bold = False
                new_run.font.size = orig_run.font.size or Pt(13)
            try:
                    new_run.font.color.rgb = orig_run.font.color.rgb
                except:
                    pass
        print('Body format preserved from original')

print('✅ 第4页 (SW-2026-1196) 修复完成（保留原始模板格式）')

# ================================================================
# 2. 重新处理第5页（slide index 4） — 同样保留原始格式
# ================================================================
slide5 = prs.slides[4]
text_shapes5 = [s for s in slide5.shapes if s.has_text_frame]
orig_slide5 = prs_orig.slides[4]
orig_text_shapes5 = [s for s in orig_slide5.shapes if s.has_text_frame]

title5 = None
badge5 = None
body5 = None
for s in text_shapes5:
    if s.name == 'TextBox 5':
        title5 = s
    elif s.name == 'Rectangle 6':
        badge5 = s
    elif s.name == 'TextBox 7':
        body5 = s

if title5:
    orig_title5 = next((s for s in orig_text_shapes5 if s.name == 'TextBox 5'), None)
    if orig_title5 and orig_title5.text_frame.paragraphs[0].runs:
        orig_run = orig_title5.text_frame.paragraphs[0].runs[0]
        para = title5.text_frame.paragraphs[0]
        para.clear()
        new_run = para.add_run()
        new_run.text = '成国资发〔2025〕15号：市属国企资产租赁管理制度'
        new_run.font.name = orig_run.font.name
        new_run.font.size = orig_run.font.size
        new_run.font.bold = orig_run.font.bold
        try:
                    new_run.font.color.rgb = orig_run.font.color.rgb
                except:
                    pass

if badge5:
    orig_badge5 = next((s for s in orig_text_shapes5 if s.name == 'Rectangle 6'), None)
    if orig_badge5 and orig_badge5.text_frame.paragraphs[0].runs:
        orig_run = orig_badge5.text_frame.paragraphs[0].runs[0]
        para = badge5.text_frame.paragraphs[0]
        para.clear()
        new_run = para.add_run()
        new_run.text = '业务直击'
        new_run.font.name = orig_run.font.name
        new_run.font.size = orig_run.font.size
        new_run.font.bold = orig_run.font.bold
        try:
                    new_run.font.color.rgb = orig_run.font.color.rgb
                except:
                    pass

if body5:
    orig_body5 = next((s for s in orig_text_shapes5 if s.name == 'TextBox 7'), None)
    body5_lines = [
        '▎核心要求：公开招租',
        '必须在市政府网+市国资委网+产权交易所+集团官网同步发布',
        '正式披露≥10个工作日 | 结果公示≥3个工作日',
        '',
        '▎租金定价机制',
        '以评估/估价/询价/公开市场价为依据，综合市场需求、区位布局、功能业态确定',
        '',
        '▎五条红线',
        '❌ 拆分资产规避公开招租   ❌ 未经批准擅自出租',
        '❌ 擅自改变用途或租赁期限   ❌ 截留/挪用/私分租金收入',
        '❌ 违规转租 · 有违约行为不得续租 · 短期租赁到期不得续租',
        '',
        '▎监管机制',
        '"三重一大"集体决策 → 报市国资委备案 → 年度自查自纠 → 审计+专项检查',
        '重点查：未公开租赁 · 租金底价偏高 · 超长租期 · 违规转租',
        '',
        '▶ 与你的关系',
        '天府广场商业租赁的合同续租、商户转租、租金定价——全部在这个文件里画了红线',
    ]
    
    if orig_body5 and orig_body5.text_frame.paragraphs[0].runs:
        orig_run = orig_body5.text_frame.paragraphs[0].runs[0]
        tf = body5.text_frame
        
        for pi in range(len(tf.paragraphs)):
            if tf.paragraphs[pi].runs:
                tf.paragraphs[pi].clear()
        
        while len(tf.paragraphs) < len(body5_lines):
            tf.add_paragraph()
        
        for i, line in enumerate(body5_lines):
            para = tf.paragraphs[i]
            para.clear()
            new_run = para.add_run()
            new_run.text = line
            new_run.font.name = orig_run.font.name or '微软雅黑'
            new_run.font.size = orig_run.font.size or Pt(13)
            if line.startswith('▎') or line.startswith('▶') or line.startswith('❌') or line.startswith('"'):
                new_run.font.bold = True
                new_run.font.size = orig_run.font.size or Pt(14)
            else:
                new_run.font.bold = False
                new_run.font.size = orig_run.font.size or Pt(13)
            try:
                    new_run.font.color.rgb = orig_run.font.color.rgb
                except:
                    pass

print('✅ 第5页 (成国资发15号) 修复完成（保留原始模板格式）')

# ================================================================
# 3. 修复Logo（第4页Picture 4）— 从原始文件复制
# ================================================================
orig_slide4 = prs_orig.slides[3]
curr_slide4 = prs.slides[3]

# 找到原始文件和当前文件的Picture 4
orig_pic4 = None
curr_pic4 = None
for s in orig_slide4.shapes:
    if s.name == 'Picture 4':
        orig_pic4 = s
        break
for s in curr_slide4.shapes:
    if s.name == 'Picture 4':
        curr_pic4 = s
        break

if orig_pic4 and curr_pic4:
    # 复制图片blob
    orig_img = orig_pic4.image
    # 删除当前图片shape并重建
    # 获取当前元素的类型和位置
    sp = curr_pic4._element
    blipFill = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
    if blipFill is not None:
        # 获取关系
        rId = blipFill.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
        if rId:
            # 替换图片part
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            img_part_orig = orig_pic4.part
            # 获取原始图片blob
            import io
            orig_blob = orig_img.blob
            # 在当前slide的part中替换图片
            curr_part = slide4.part
            try:
                img_part = curr_part.related_part(rId)
                # 直接用新blob覆盖
                img_part._blob = orig_blob
                print(f'✅ Logo图片已恢复 (size={len(orig_blob)}bytes)')
            except Exception as e:
                print(f'⚠ Logo恢复失败: {e}')
                print('  rId={}'.format(rId))

# ================================================================
# 4. 修改第27页（关键启示） - slide index 26
# ================================================================
slide27 = prs.slides[26]
for shape in slide27.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full_text = para.text
            if '1号文 → 用系统减少' in full_text and '2号文 → 用框架消除' in full_text:
                for run in para.runs:
                    if '1号文' in run.text:
                        run.text = run.text.replace(
                            '1号文 → 用系统减少"人"的随意性       2号文 → 用框架消除"管"的盲区',
                            '1196号文 → 13大类追责全覆盖, 终身问责不设限   15号租赁文 → 公开招租五条红线'
                        )
                        print(f'✅ 第27页: 1号文2号文内容已替换')

# ================================================================
# 5. 修改第28页（合规自查要点） - slide index 27
# ================================================================
slide28 = prs.slides[27]
for shape in slide28.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # 替换标签和问题
                if '对照1号文' == run.text.strip():
                    run.text = '对照1196号文'
                    print('✅ 第28页: 对照1号文 → 对照1196号文')
                elif '信息系统打通了吗？预警自动了吗？' == run.text.strip():
                    run.text = '追责红线画清楚了吗？投融资决策合规吗？'
                elif '对照2号文' == run.text.strip():
                    run.text = '对照15号租赁文'
                    print('✅ 第28页: 对照2号文 → 对照15号租赁文')
                elif '穿透到底了吗？资金可追溯了吗？' == run.text.strip():
                    run.text = '公开招租执行了吗？转租合规了吗？'

print('✅ 第28页 合规自查要点 已更新')

# ================================================================
# 保存
# ================================================================
prs.save(dst_fixed)
print(f'\n✅ 最终版: {dst_fixed}')
print(f'   文件大小: {os.path.getsize(dst_fixed)//1024}KB')
