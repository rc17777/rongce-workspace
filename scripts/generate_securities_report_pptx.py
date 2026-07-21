# -*- coding: utf-8 -*-
"""
融策·可编辑研报图文协同版 PPTX 生成器
====================================
目标：不是导出静态图片，而是生成 PowerPoint 原生可编辑对象：
- 标题栏、KPI、正文观点、右侧判断框均为文本框
- 图表使用 Office 原生 chart，可编辑数据和样式
- 每页遵循券商研报逻辑：页标题 -> 核心指标 -> 图表证据 -> 右侧判断 -> 底部来源

用法：
    python -X utf8 scripts/generate_securities_report_pptx.py output/融策_研报图文协同版.pptx
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.xmlchemy import OxmlElement

NAVY = RGBColor(6, 26, 51)
NAVY2 = RGBColor(10, 42, 74)
GOLD = RGBColor(184, 138, 68)
GOLD2 = RGBColor(214, 176, 113)
TEAL = RGBColor(26, 111, 120)
INK = RGBColor(31, 41, 51)
GRAY = RGBColor(83, 97, 109)
MUTED = RGBColor(135, 147, 160)
PANEL = RGBColor(244, 246, 248)
WHITE = RGBColor(255, 255, 255)
GRID = RGBColor(217, 222, 229)
RED = RGBColor(184, 64, 58)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_font(run, size=10, color=INK, bold=False):
    run.font.name = 'Microsoft YaHei'
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, x, y, w, h, text='', size=10, color=INK, bold=False, align=PP_ALIGN.LEFT, fill=None, border=None):
    shape = slide.shapes.add_textbox(x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return shape


def add_header(slide, title, subtitle, tag='融策·审盾研究'):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    gold = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.62), SLIDE_W, Inches(0.07))
    gold.fill.solid(); gold.fill.fore_color.rgb = GOLD; gold.line.fill.background()
    add_textbox(slide, Inches(0.36), Inches(0.08), Inches(3), Inches(0.18), tag, 8.5, GOLD2, True)
    add_textbox(slide, Inches(0.34), Inches(0.27), Inches(8.8), Inches(0.28), title, 18, WHITE, True)
    add_textbox(slide, Inches(9.6), Inches(0.29), Inches(3.35), Inches(0.24), subtitle, 9, RGBColor(201, 210, 220), False, PP_ALIGN.RIGHT)


def add_kpi_row(slide, kpis):
    x0, y, total_w, h = Inches(0.38), Inches(0.88), Inches(12.55), Inches(0.72)
    gap = Inches(0.08)
    box_w = (total_w - gap * (len(kpis) - 1)) / len(kpis)
    for i, kpi in enumerate(kpis):
        x = x0 + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, h)
        box.fill.solid(); box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = RGBColor(213, 219, 227); box.line.width = Pt(0.75)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.10), box_w - Inches(0.2), Inches(0.18), kpi['label'], 8, GRAY)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.32), box_w - Inches(0.2), Inches(0.28), kpi['value'], 17, kpi.get('color', NAVY), True)


def add_note_panel(slide, notes, x=Inches(8.35), y=Inches(1.95), w=Inches(4.25), h=Inches(4.35)):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(250, 251, 252)
    panel.line.color.rgb = RGBColor(205, 213, 223); panel.line.width = Pt(0.8)
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.38))
    head.fill.solid(); head.fill.fore_color.rgb = NAVY2; head.line.fill.background()
    add_textbox(slide, x + Inches(0.14), y + Inches(0.08), w - Inches(0.25), Inches(0.22), '核心判断', 10, WHITE, True)
    body = slide.shapes.add_textbox(x + Inches(0.20), y + Inches(0.58), w - Inches(0.38), h - Inches(0.72))
    body.line.fill.background()
    tf = body.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, note in enumerate(notes, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run(); r1.text = f'{i}.  '; set_font(r1, 8.8, GOLD, True)
        r2 = p.add_run(); r2.text = note; set_font(r2, 8.8, INK, False)


def add_source(slide, source):
    add_textbox(slide, Inches(0.38), Inches(7.08), Inches(6.3), Inches(0.18), f'资料来源：{source}', 7.5, MUTED)
    add_textbox(slide, Inches(8.1), Inches(7.08), Inches(4.8), Inches(0.18), '制图：融策AI审计中台 | 2026-07-21', 7.5, MUTED, False, PP_ALIGN.RIGHT)


def style_chart(chart):
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.font.name = 'Microsoft YaHei'
    chart.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = GRAY
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.color.rgb = GRAY
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = GRID
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)


def set_series_colors(chart, colors):
    for idx, series in enumerate(chart.series):
        color = colors[idx % len(colors)]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.25)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.data_labels.font.size = Pt(7.5)
        series.data_labels.font.color.rgb = INK


def add_line_chart(slide):
    add_header(slide, '财政收入修复斜率放缓，税收收入弹性仍是关键变量', '宏观财政专题 | 年度趋势跟踪')
    add_kpi_row(slide, [
        {'label': '2025E预算收入', 'value': '22.8万亿', 'color': NAVY},
        {'label': '2024-2025E增量', 'value': '+0.8万亿', 'color': GOLD},
        {'label': '税收占比', 'value': '83.3%', 'color': TEAL},
    ])
    add_textbox(slide, Inches(0.55), Inches(1.75), Inches(6.6), Inches(0.28), '图1：全国一般公共预算收入与税收收入走势（万亿元）', 12, INK, True)
    data = CategoryChartData()
    data.categories = ['2020', '2021', '2022', '2023', '2024', '2025E']
    data.add_series('一般公共预算收入', (18.29, 20.25, 20.37, 21.68, 22.00, 22.80))
    data.add_series('税收收入', (15.43, 17.27, 16.66, 18.11, 18.20, 19.00))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.55), Inches(2.15), Inches(7.25), Inches(4.15), data).chart
    style_chart(chart)
    chart.value_axis.minimum_scale = 14
    chart.value_axis.maximum_scale = 24
    for s in chart.series:
        s.has_data_labels = True
        s.data_labels.show_value = False
    set_series_colors(chart, [NAVY, GOLD])
    add_note_panel(slide, [
        '财政收入修复并非线性扩张，税基质量和房地产链条回暖仍决定后续弹性。',
        '税收收入占比维持高位，说明非税收入拉动空间有限，审计应重点关注税源真实性。',
        '若预算收入增速明显高于经济增速，应回查一次性收入、非税缴库和跨期调节。',
    ])
    add_source(slide, '财政部，Wind，融策会计师事务所整理')


def add_bar_chart(slide):
    add_header(slide, '项目结构决定审计资源配置，绩效评价与工程决算是主战场', '融策业务结构分析 | 项目类型对比')
    add_kpi_row(slide, [
        {'label': '项目总量', 'value': '168个', 'color': NAVY},
        {'label': '优势业务占比', 'value': '43.5%', 'color': GOLD},
        {'label': '可AI复核项目', 'value': '100+', 'color': TEAL},
    ])
    add_textbox(slide, Inches(0.55), Inches(1.75), Inches(6.6), Inches(0.28), '图2：2025年度审计咨询项目类型分布（个）', 12, INK, True)
    data = CategoryChartData()
    data.categories = ['绩效评价', '经责审计', '工程决算', '专项审计', '资产清查', '预算执行']
    data.add_series('融策承接', (45, 32, 28, 25, 20, 18))
    data.add_series('行业均值', (38, 35, 22, 30, 25, 22))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(2.15), Inches(7.25), Inches(4.15), data).chart
    style_chart(chart)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 50
    for s in chart.series:
        s.has_data_labels = True
        s.data_labels.show_value = True
        s.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    set_series_colors(chart, [NAVY, RGBColor(154, 166, 178)])
    add_note_panel(slide, [
        '绩效评价和工程竣工决算是最适合作为审盾一期验证样板的业务线。',
        '项目结构越标准化，越适合沉淀复核清单、法规引用和图表模板。',
        '后续应按业务线建立“数据口径-指标体系-报告模板”三件套。',
    ])
    add_source(slide, '融策项目台账，行业访谈，融策AI审计中台')


def add_report_section(slide):
    add_header(slide, '报告与图表协同：用正文结论牵引数据证据', '融策研报模板 | 图文协同页')
    add_textbox(slide, Inches(0.55), Inches(1.05), Inches(7.2), Inches(0.36), '1. 审计研报页不应只放图，而要形成“观点-证据-判断-行动”的闭环', 15, NAVY, True)
    add_textbox(slide, Inches(0.55), Inches(1.55), Inches(7.25), Inches(1.25),
                '券商研报的图表不是正文的装饰，而是论证链条中的证据单元。每一张图应回答一个明确问题：趋势是什么、异常在哪里、管理动作是什么。融策报告应将图表嵌入章节逻辑，而不是把图表集中堆在附录。',
                10.5, INK)
    add_textbox(slide, Inches(0.55), Inches(3.0), Inches(3.4), Inches(0.42), '建议页内结构', 12, WHITE, True, fill=NAVY2)
    add_textbox(slide, Inches(0.55), Inches(3.48), Inches(3.4), Inches(2.15),
                '① 章节结论：一句话先给判断\n② KPI：用3个数字压住场\n③ 主图：呈现趋势/对比/结构\n④ 右侧判断：解释为什么重要\n⑤ 来源：交代数据可信度',
                10, INK, fill=PANEL, border=RGBColor(213, 219, 227))
    add_textbox(slide, Inches(4.35), Inches(3.0), Inches(3.4), Inches(0.42), '适用报告类型', 12, WHITE, True, fill=NAVY2)
    add_textbox(slide, Inches(4.35), Inches(3.48), Inches(3.4), Inches(2.15),
                '绩效评价报告\n经责审计结果报告\n财政专项资金分析报告\n工程决算审计专题报告\n招投标串标风险分析报告',
                10, INK, fill=PANEL, border=RGBColor(213, 219, 227))
    add_note_panel(slide, [
        '后续模板应以“页”为最小生产单元，每页绑定一个核心判断。',
        '图表数据、正文结论、资料来源要同步生成，避免图文两张皮。',
        '可编辑PPTX适合内部打磨和客户汇报，Word适合正式归档和出具。',
    ])
    add_source(slide, '长江证券/国金证券研报结构拆解，融策模板化复刻')


def build(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_report_section(slide)

    slide = prs.slides.add_slide(blank)
    add_line_chart(slide)

    slide = prs.slides.add_slide(blank)
    add_bar_chart(slide)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f'✅ 可编辑研报图文协同版已生成: {output_path}')


if __name__ == '__main__':
    output = sys.argv[1] if len(sys.argv) > 1 else 'output/融策_研报图文协同版_可编辑.pptx'
    build(output)
