"""Remove empty placeholder shapes from PPT"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.shapes.placeholder import PlaceholderPicture

SRC = r'D:\openclaw-workspace\output\v5_final.pptx'
DST = r'D:\openclaw-workspace\output\v5_clean.pptx'
DESK = r'C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v5_模板.pptx'

prs = Presentation(SRC)

removed = 0
for si, slide in enumerate(prs.slides):
    shapes_to_remove = []
    for shape in slide.shapes:
        # Check if it's a placeholder
        is_ph = hasattr(shape, 'is_placeholder') and shape.is_placeholder
        # Check if it has empty/default text
        has_text = shape.has_text_frame
        if has_text:
            text = shape.text_frame.text.strip()
            is_empty = not text or text in ['点击此处添加标题', '点击此处添加文本', '单击此处添加标题', '单击此处添加文本', 'Click to add title', 'Click to add text']
        else:
            is_empty = False
        
        if is_ph and has_text and is_empty:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        removed += 1
        print(f'  Slide {si+1}: removed "{shape.name}"')

print(f'\nTotal removed: {removed} shapes')
prs.save(DST)

# Copy to desktop
import shutil, os
try:
    if os.path.exists(DESK):
        os.remove(DESK)
    shutil.copy2(DST, DESK)
    print(f'Copied to desktop: {DESK}')
except Exception as e:
    # Try with suffix
    base = r'C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v5_模板'
    for i in range(2, 10):
        alt = f'{base}_{i}.pptx'
        if not os.path.exists(alt):
            shutil.copy2(DST, alt)
            print(f'Copied to: {alt}')
            break
