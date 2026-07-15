"""
v5b: 全新方案 — 从模板加载，用BytesIO保存绕过文件锁，重写Content_Types
"""
import sys, os, shutil, io, zipfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ── Paths ──
TMPL = r'C:\Users\scrccpa\Desktop\物资管理制度培训-2025.1.20.pptx'
OUT  = r'D:\openclaw-workspace\output\v5_final.pptx'
DESK = r'C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v5_模板.pptx'
ASSETS = r'D:\openclaw-workspace\scripts\ppt_assets'

# ── Colors ──
CLR = {
    'dark':     RGBColor(0x33, 0x33, 0x33),
    'title':    RGBColor(0x1A, 0x3C, 0x6E),
    'accent':   RGBColor(0xC0, 0x39, 0x2B),
    'blue':     RGBColor(0x29, 0x80, 0xB9),
    'green':    RGBColor(0x27, 0xAE, 0x60),
    'amber':    RGBColor(0xE6, 0x7E, 0x22),
    'gray':     RGBColor(0x7F, 0x8C, 0x8D),
    'white':    RGBColor(0xFF, 0xFF, 0xFF),
}

FONT = '微软雅黑'
FONT_TITLE = '微软雅黑'
SW = Emu(12192000)
SH = Emu(6858000)

# ── Assets ──
BANNER = os.path.join(ASSETS, 'slide2_图片 3.png')
LOGO   = os.path.join(ASSETS, 'slide2_图片 8.png')
COVER_BG1 = os.path.join(ASSETS, 'slide1_图片 1.png')
COVER_BG2 = os.path.join(ASSETS, 'slide1_图片 2.png')
COVER_LOGO = os.path.join(ASSETS, 'slide1_图片 7.png')

def emu(inches):
    return int(inches * 914400)

# ── Load & clean ──
print('Loading template...')
prs = Presentation(TMPL)

# Remove all slides by deleting from sldIdLst AND their rels
sldIdLst = prs.slides._sldIdLst
to_remove = list(sldIdLst)
for sldId in to_remove:
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    sldIdLst.remove(sldId)
    if rId:
        try:
            prs.part.drop_rel(rId)
        except:
            pass

print(f'Slides after cleanup: {len(prs.slides)}')

# ── Helpers ──
def add_slide(layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_decor(slide):
    slide.shapes.add_picture(BANNER, 0, SH - emu(1.3), SW, emu(1.3))
    slide.shapes.add_picture(LOGO, emu(0.4), emu(0.3), emu(1.4), emu(0.6))

def add_title_box(slide, text, top=emu(0.5), left=emu(0.8), width=None, height=emu(0.8),
                  font_size=30, color=None, bold=True, align=PP_ALIGN.LEFT):
    if width is None: width = SW - emu(1.6)
    if color is None: color = CLR['title']
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.name = FONT_TITLE; p.font.size = Pt(font_size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return txBox

def add_text_box(slide, lines, top, left=emu(0.8), width=None, height=None,
                 font_size=16, color=None, line_spacing=1.5):
    if width is None: width = SW - emu(1.6)
    if height is None: height = SH - top - emu(1.8)
    if color is None: color = CLR['dark']
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, str):
            p.text = line; p.font.bold = False; p.font.color.rgb = color
        elif isinstance(line, tuple):
            p.text = line[0]
            p.font.bold = line[1] if len(line) > 1 else False
            p.font.color.rgb = line[2] if len(line) > 2 else color
        p.font.name = FONT; p.font.size = Pt(font_size)
        p.space_after = Pt(4); p.level = 0
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
        spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox

def add_label(slide, text, top, left, width, height, bg_color=CLR['blue'], font_size=14, font_color=None):
    if font_color is None: font_color = CLR['white']
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = bg_color; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.name = FONT; p.font.size = Pt(font_size)
    p.font.bold = True; p.font.color.rgb = font_color; p.alignment = PP_ALIGN.CENTER
    return shape

# ═══════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════
print('Slide 1: Cover')
s = add_slide(0)
s.shapes.add_picture(COVER_BG1, emu(-0.0), emu(0.9), emu(13.3), emu(6.6))
s.shapes.add_picture(COVER_BG2, emu(-0.3), emu(0.9), emu(13.9), emu(6.7))
s.shapes.add_picture(COVER_LOGO, emu(0.4), emu(0.3), emu(1.4), emu(0.6))
add_title_box(s, '提升公司全员审计风险意识', top=emu(2.0), font_size=40, color=CLR['white'],
              align=PP_ALIGN.CENTER, width=SW - emu(2.0), left=emu(1.0))
tb = add_text_box(s, ['四川融策会计师事务所  |  2026年6月  |  四川轨道公司专题培训'],
                  top=emu(3.5), font_size=18, color=RGBColor(0xDD,0xDD,0xDD),
                  left=emu(1.0), width=SW - emu(2.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════
print('Slide 2: Agenda')
s = add_slide(1); add_decor(s)
add_title_box(s, '培训议程', top=emu(0.6))
items = [
    ('第一部分', '国务院核心文件解读（1号文·2号文·15号文·46号令）', '⏱ 40分钟', CLR['blue']),
    ('第二部分', '天府广场审计实战案例（4大领域·13项问题）', '⏱ 45分钟', CLR['accent']),
    ('第三部分', '总结与建议（制度落地六步法·合规自查）', '⏱ 5分钟', CLR['green']),
]
y = emu(1.8)
for label, desc, time, clr in items:
    add_label(s, label, y, emu(1.0), emu(2.0), emu(0.5), bg_color=clr, font_size=14)
    add_text_box(s, [desc], top=y+emu(0.05), left=emu(3.3), width=emu(7.5), height=emu(0.5), font_size=18)
    add_text_box(s, [time], top=y+emu(0.05), left=emu(11.0), width=emu(1.8), height=emu(0.5), font_size=14, color=CLR['gray'])
    y += emu(1.1)

# ═══════════════════════════════════
# SLIDE 3: PART 1 DIVIDER
# ═══════════════════════════════════
print('Slide 3: Part 1 Divider')
s = add_slide(2); add_decor(s)
add_title_box(s, '第一部分', top=emu(2.0), font_size=44, color=CLR['blue'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
add_title_box(s, '国务院核心文件解读', top=emu(3.2), font_size=28, color=CLR['dark'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0), bold=False)
tb = add_text_box(s, ['1号文 → 2号文 → 15号文 → 46号令：从技术底座到兜底保障的完整政策链条'],
                  top=emu(4.2), font_size=16, color=CLR['gray'], left=emu(1.5), width=SW-emu(3.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 4: DOC #1
# ═══════════════════════════════════
print('Slide 4: Doc #1')
s = add_slide(1); add_decor(s)
add_title_box(s, '1号文：《推动央企加快财务数智化转型升级》', top=emu(0.6), font_size=26)
add_label(s, '技术底座', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['blue'])
add_text_box(s, [
    ('▎建系统', True, CLR['title']),
    '统一财务数智化平台，打通数据孤岛。停车场收费、商业租赁、物业运营、设备维保——系统不通，数据就是一座座孤岛',
    ('▎强支撑', True, CLR['title']),
    '数据治理、安全防护、人才保障。系统建好了数据质量跟不上，后续分析决策都跑偏',
    ('▎给工具', True, CLR['title']),
    '部署智能化财务分析工具，实现风险自动预警',
    ('▶ 与你的关系', True, CLR['accent']),
    '每一笔收费、每一份合同、每一次审批，将来都要能被系统追踪。人在做，系统在看，审计在复盘',
], top=emu(1.9), font_size=15, line_spacing=1.25)

# ═══════════════════════════════════
# SLIDE 5: DOC #2
# ═══════════════════════════════════
print('Slide 5: Doc #2')
s = add_slide(1); add_decor(s)
add_title_box(s, '2号文：《加强央企穿透式监管指导意见（试行）》', top=emu(0.6), font_size=26)
add_label(s, '总纲领', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['blue'])
add_text_box(s, [
    ('四层穿透', True, CLR['title']),
    ('① 层级穿透', True, CLR['dark']), '穿透股权结构到底层实体。子公司、代管单位、合作经营方——监管一穿到底',
    ('② 资金穿透', True, CLR['dark']), '追踪资金流向全链条。谁收的、怎么收的、进了哪个账户、什么时候入的账',
    ('③ 业务穿透', True, CLR['dark']), '穿透业务全流程。商户实际经营面积是否超合同范围？用制度和数据把缝隙堵上',
    ('④ 追责穿透', True, CLR['dark']), '穿透责任链条至个人。出问题不再是"集体负责"，具体的人要扛',
    ('▶ 与你的关系', True, CLR['accent']), '监管的天网正在织密，以前管不到的灰色地带，以后都会被穿透',
], top=emu(1.9), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 6: DOC #15
# ═══════════════════════════════════
print('Slide 6: Doc #15')
s = add_slide(1); add_decor(s)
add_title_box(s, '15号文：《2026年央企内控体系建设与监督工作通知》', top=emu(0.6), font_size=26)
add_label(s, '操作手册', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['amber'])
add_text_box(s, [
    ('五大核心要求', True, CLR['title']),
    '❶ 完善内控制度体系 — 不相容职务分离、手动放行审批流程，有没有？',
    '❷ 强化重大风险防控 — 收费异常、商户拖欠租金…能不能被及时发现？',
    '❸ 落实监督评价机制 — 是自查自评走过场，还是真正过硬的第三方评价？',
    '❹ 推进数智化赋能 — 系统不行就要建，数据不通就要通，不是"将就"是"建设"',
    '❺ 建立整改问责闭环 — 审计发现问题，限期整改，整改不力问责',
    ('▶ 与你的关系', True, CLR['accent']), '每一条要求，都能在天府广场审计报告中找到对应的反面案例',
], top=emu(1.9), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 7: ORDER #46
# ═══════════════════════════════════
print('Slide 7: Order #46')
s = add_slide(1); add_decor(s)
add_title_box(s, '46号令：《违规经营投资责任追究制度》', top=emu(0.6), font_size=26)
add_label(s, '兜底保障', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['accent'])
add_text_box(s, [
    ('三句话总结：程序合规 · 尽职免责 · 失职追责', True, CLR['accent']),
    ('🔴 明底线', True, CLR['dark']), '违规行为红线画清楚。资产损失、违规经营、内控失效、隐瞒不报——白纸黑字',
    ('🔴 严追责', True, CLR['dark']), '资产损失必问责，重大决策终身追责。今天签的审批，十年后出问题，追责链条不断',
    ('🟢 宽容错', True, CLR['dark']), '改革创新中的探索失误、非主观过失、履行了必要程序且未谋私利——可从轻或免予追责',
    '不是让你不敢干事，是让你干事走程序',
], top=emu(1.9), font_size=14, line_spacing=1.25)

# ═══════════════════════════════════
# SLIDE 8: PANORAMA
# ═══════════════════════════════════
print('Slide 8: Panorama')
s = add_slide(1); add_decor(s)
add_title_box(s, '四位一体：政策全景图', top=emu(0.6))
cols = [
    ('1号文', '看见', '数智化系统\n让异常无处遁形', CLR['blue']),
    ('2号文', '看准', '穿透监管框架\n界定查什么、查多深', CLR['blue']),
    ('15号文', '做到', '内控操作手册\n把制度变成行动', CLR['amber']),
    ('46号令', '守住', '追责底线\n让违规必付代价', CLR['accent']),
]
cw, g = emu(2.7), emu(0.3)
for i, (label, verb, desc, clr) in enumerate(cols):
    x = emu(0.7) + i*(cw+g)
    add_label(s, label, emu(1.8), x+emu(0.3), emu(2.0), emu(0.6), bg_color=clr, font_size=22)
    tb = add_text_box(s, [verb], top=emu(2.6), left=x, width=cw, height=emu(0.5), font_size=28, color=clr)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = add_text_box(s, [desc], top=emu(3.2), left=x+emu(0.1), width=cw-emu(0.2), height=emu(1.2), font_size=14)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
tb = add_text_box(s, ['技术 ──→ 制度 ──→ 执行 ──→ 问责  =  完整闭环'],
                  top=emu(4.6), font_size=18, color=CLR['title'], left=emu(1.0), width=SW-emu(2.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 9: ORDER #46 DEEP DIVE
# ═══════════════════════════════════
print('Slide 9: Order #46 Deep Dive')
s = add_slide(1); add_decor(s)
add_title_box(s, '46号令深度解读：追责倒查链', top=emu(0.6), font_size=26)
add_label(s, '🆕 重点', emu(1.2), emu(0.8), emu(1.5), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('🔴 13大类98种违规情形', True, CLR['accent']),
    '投资决策失误 · 产权交易违规 · 工程超概 · 围标串标 · 资金挪用 · 违规签约',
    '财务造假 · 违规担保 · 境外投资失控 · 薪酬福利违规……白纸黑字列出来',
    ('🔴 追责对象全覆盖', True, CLR['accent']),
    '决策人（签字）→ 执行人（经办）→ 审核人（审核）→ 监管人（监督）→ 退休人员',
    '退休 ≠ 安全，离职 ≠ 免责',
    ('🔴 追责时间轴', True, CLR['accent']),
    '重大违规：无期限、终身追责 | 较大多数：倒查20~30年 | 一般违规：倒查10~15年',
    ('⚠ 今天签的字、经手的事，20年后可能被翻出来追责——不是吓唬人，是白纸黑字', True, CLR['accent']),
], top=emu(1.8), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 10: ROADMAP
# ═══════════════════════════════════
print('Slide 10: Roadmap')
s = add_slide(1); add_decor(s)
add_title_box(s, '穿透式监管落地时间表', top=emu(0.6))
milestones = [
    ('2026', '央企全面达标', '100%建成穿透式\n智能监管平台\n十大领域全覆盖', CLR['accent']),
    ('2027', '省级国企推广', '60%以上省属国企\n建成监管平台\n大监督协同机制', CLR['amber']),
    ('2028', '地市县区覆盖', '穿透式监管\n从央企到省属到市属\n到区县全部覆盖', CLR['blue']),
]
cw, g = emu(3.6), emu(0.4)
for i, (year, title, desc, clr) in enumerate(milestones):
    x = emu(0.7) + i*(cw+g)
    add_label(s, year, emu(1.6), x, emu(1.8), emu(0.55), bg_color=clr, font_size=24)
    tb = add_text_box(s, [title], top=emu(2.3), left=x, width=cw, height=emu(0.5), font_size=16, color=CLR['title'])
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = add_text_box(s, [desc], top=emu(2.9), left=x+emu(0.2), width=cw-emu(0.4), height=emu(1.5), font_size=13)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('四大支柱', True, CLR['title']),
    '制度：统一监管制度+标准化规程库    数据：数据治理五大域、打破孤岛',
    '工具：智能监管平台、人防→技防→智防    协同：大监督+内控协同+智能预警',
    ('▶ 四川轨道公司作为省级国企，2027年就是截止日期——不是"未来要做"，是"明年完成"', True, CLR['accent']),
], top=emu(4.5), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 11: PART 2 DIVIDER
# ═══════════════════════════════════
print('Slide 11: Part 2 Divider')
s = add_slide(2); add_decor(s)
add_title_box(s, '第二部分', top=emu(2.0), font_size=44, color=CLR['accent'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
add_title_box(s, '天府广场独立商业项目·审计实战', top=emu(3.2), font_size=28, color=CLR['dark'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0), bold=False)
tb = add_text_box(s, ['审计期间：2024.01 — 2026.03 ｜ 4大领域 · 13项问题'],
                  top=emu(4.2), font_size=16, color=CLR['gray'], left=emu(1.5), width=SW-emu(3.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 12: AUDIT OVERVIEW
# ═══════════════════════════════════
print('Slide 12: Audit Overview')
s = add_slide(1); add_decor(s)
add_title_box(s, '审计概况', top=emu(0.6))
domains = [
    ('🅿 停车场管理', '车位底数、收费内控\n减免审批、计费规则', CLR['accent']),
    ('🏢 经营管理', '房间使用、租赁范围\n调换审批、退租交接', CLR['amber']),
    ('🏠 物业用房', '台账管理、内部改造\n外单位使用', CLR['blue']),
    ('👁 现场管理', '维保记录、设备台账\n年度考评', CLR['green']),
]
cw, g = emu(2.7), emu(0.3)
for i, (title, desc, clr) in enumerate(domains):
    x = emu(0.7) + i*(cw+g)
    add_label(s, title, emu(1.5), x, cw, emu(0.55), bg_color=clr, font_size=15)
    tb = add_text_box(s, [desc], top=emu(2.2), left=x+emu(0.1), width=cw-emu(0.2), height=emu(1.2), font_size=13)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('13项问题 · 涉及金额22,671元', True, CLR['title']),
    '金额不算大，但问题的性质比金额重要——停车场内控缺失、物业用房台账失范，都是系统漏洞而非个案偶然',
    '今天22,671元，明天如果漏洞没堵上，可能就是后面加两个零',
], top=emu(3.8), font_size=15, line_spacing=1.3)

# ═══════════════════════════════════
# SLIDE 13: PARKING I
# ═══════════════════════════════════
print('Slide 13: Parking I')
s = add_slide(1); add_decor(s)
add_title_box(s, '停车场管理问题（一）', top=emu(0.6), font_size=26)
add_label(s, '问题1', emu(1.4), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('停车位底数不清，账实不符，设施失效', True, CLR['accent']),
    '备案379个 → 系统330个 → 画线331个 → 有效320个  四个数字，同一个停车场！',
    '现场停车感应装置基本全部损坏。备案后无动态更新、系统数据无定期核对、硬件坏了没人修',
], top=emu(1.4), left=emu(2.2), width=emu(9.5), height=emu(1.5), font_size=14, line_spacing=1.3)
add_label(s, '问题2', emu(3.2), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('停车场内控不健全 ⚠ 13项问题中性质最严重', True, CLR['accent']),
    '❶ 不相容职务未分离：同一人管审核+收费+对账+系统录入+充值+撤销——三权集于一身',
    '❷ 手动放行管控薄弱：351次手动抬杆，绝大部分是临时车辆',
    '❸ 三缺：抬杆记录台账缺失 + 保安收款记录缺失 + 收费员核对记录缺失',
    ('灵魂拷问：351次手动放行，应收多少、实收多少、有没有少交的——谁能回答？', True, CLR['accent']),
], top=emu(3.2), left=emu(2.2), width=emu(9.5), height=emu(2.5), font_size=13, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 14: PARKING II
# ═══════════════════════════════════
print('Slide 14: Parking II')
s = add_slide(1); add_decor(s)
add_title_box(s, '停车场管理问题（二）', top=emu(0.6), font_size=26)
add_label(s, '问题3', emu(1.4), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('减免优惠车辆审批不严，准入管控流于形式', True, CLR['accent']),
    '9辆公务车/员工车未履行报备手续，免费停车2-3月 → 涉及9,813元',
    '"三人行"优惠演变成1带多：67人中老客户仅9人，新客户58人 → 审批依据仅为团购名单',
    '流程走过场 → 审核松一松 → 费就漏一漏',
], top=emu(1.4), left=emu(2.2), width=emu(9.5), height=emu(1.8), font_size=14, line_spacing=1.3)
add_label(s, '问题4', emu(3.5), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('计费规则管理不到位，收费标准备案变更不及时', True, CLR['accent']),
    '两条规则设置依据查不到：缴费后30分钟不离场重计费、新能源车多15分钟宽限期',
    '收费标准2017年调整后，至今未做备案变更',
    '合规不是"我觉得合理就行"，是"文件怎么说就怎么做"',
], top=emu(3.5), left=emu(2.2), width=emu(9.5), height=emu(1.8), font_size=14, line_spacing=1.3)

# ═══════════════════════════════════
# SLIDE 15: OPERATIONS
# ═══════════════════════════════════
print('Slide 15: Operations')
s = add_slide(1); add_decor(s)
add_title_box(s, '经营管理问题（问题5-8）', top=emu(0.6), font_size=26)
issues = [
    ('问题5', '房间使用超出协议范围', '今站购物中心17间房，10间未签协议。管理人员觉得"配套用房不用签"——审计不认"默认"'),
    ('问题6', '经营区域超出协议范围', '5家商户+外摆区域超出租赁范围，存在无偿占用风险'),
    ('问题7', '房间调换未审批', '商户B247(40㎡)→B213(66㎡)，调换无审批，多出26㎡未计价 → 少收11,656元'),
    ('问题8', '退租交接存在疏漏', 'B220退租后今站仍持钥匙，退租返还单未注明钥匙返还。丢了东西算谁的？'),
]
y = emu(1.4)
for label, title, desc in issues:
    add_label(s, label, y, emu(0.8), emu(1.1), emu(0.35), bg_color=CLR['amber'], font_size=12)
    add_text_box(s, [(title, True, CLR['amber']), desc],
                 top=y, left=emu(2.1), width=emu(9.6), height=emu(1.0), font_size=13, line_spacing=1.2)
    y += emu(1.2)
add_text_box(s, [('共同病根：合同管理粗放 + 审批形同虚设 + 交接缺少闭环', True, CLR['accent'])],
             top=emu(5.8), font_size=14)

# ═══════════════════════════════════
# SLIDE 16: PROPERTY + ON-SITE
# ═══════════════════════════════════
print('Slide 16: Property + On-site')
s = add_slide(1); add_decor(s)
add_title_box(s, '物业用房与现场管理问题（问题9-13）', top=emu(0.6), font_size=24)
add_label(s, '物业用房', emu(1.3), emu(0.8), emu(1.8), emu(0.4), bg_color=CLR['blue'])
add_text_box(s, [
    ('问题9：台账不全', True, CLR['dark']),
    '5个房间仅1个入台账，门牌缺失或标错——资产连台账都不全，拿什么管？',
    ('问题10：未经审批被占用', True, CLR['dark']),
    '内部改造加装隔断未经审批，7间房被外单位使用无借用函件。你的房子被外人用了，连一张纸的凭据都拿不出来',
], top=emu(1.8), left=emu(0.8), width=emu(5.5), height=emu(3.0), font_size=13, line_spacing=1.2)
add_label(s, '现场管理', emu(1.3), emu(6.8), emu(1.8), emu(0.4), bg_color=CLR['green'])
add_text_box(s, [
    ('问题12：记录逻辑失真', True, CLR['dark']),
    '同一标段同一天两份维保记录，处理事项截然不同',
    '维保记录提前签字、不同内容用相同总结语',
    '台账空调多2台、消火栓少18个——设备实物与账面差20个',
    '⚠ 出安全事故查设备台账发现是假的——这不是审计问题，是法律问题',
    ('问题13：年度考评缺位', True, CLR['dark']),
    '仅1,200元，但反映管理链条末端失控：上级标准到基层变成一纸空文',
], top=emu(1.8), left=emu(6.8), width=emu(5.5), height=emu(3.5), font_size=13, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 17: ISSUE PANORAMA
# ═══════════════════════════════════
print('Slide 17: Issue Panorama')
s = add_slide(1); add_decor(s)
add_title_box(s, '13项问题全景图', top=emu(0.6))
nums = [
    ('13项', '问题', CLR['accent']),
    ('4大', '管理领域', CLR['amber']),
    ('22,671元', '涉及金额', CLR['blue']),
    ('5大', '核心法规被违反', CLR['accent']),
]
cw = emu(2.8)
for i, (num, label, clr) in enumerate(nums):
    x = emu(0.7) + i*(cw+emu(0.3))
    tb = add_text_box(s, [num], top=emu(1.5), left=x, width=cw, height=emu(0.6), font_size=32, color=clr)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = add_text_box(s, [label], top=emu(2.1), left=x, width=cw, height=emu(0.4), font_size=14, color=CLR['gray'])
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('共同特征链：内控体系"形似而神不似"', True, CLR['title']),
    '制度"有"但执行"空" → 台账"有"但信息"缺" → 流程"有"但留痕"无" → 审批"过"但审查"松"',
    '这就是15号文要着力解决的，也是46号令要追责的那种状态',
    ('你不是没制度，你是有了制度没当回事——这才是最可怕的', True, CLR['accent']),
], top=emu(2.8), font_size=16, line_spacing=1.35)

# ═══════════════════════════════════

# ═══════════════════════════════════
# SLIDE 18: PART 3 DIVIDER (EXPANDED)
# ═══════════════════════════════════
print('Slide 18: Part 3 Divider')
s = add_slide(2); add_decor(s)
add_title_box(s, '第三部分', top=emu(1.5), font_size=44, color=CLR['green'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
add_title_box(s, '从问题到提升·管理升华', top=emu(2.8), font_size=28, color=CLR['dark'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0), bold=False)
tb = add_text_box(s, ['13项问题 → 4大根因 → 5条建议 → 6步落地 → 1张路线图'],
                  top=emu(4.0), font_size=16, color=CLR['gray'], left=emu(1.5), width=SW-emu(3.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 19: ROOT CAUSE
# ═══════════════════════════════════
print('Slide 19: Root Cause')
s = add_slide(1); add_decor(s)
add_title_box(s, '13项问题的底层逻辑：管理闭环的四个断裂点', top=emu(0.5), font_size=24)
add_text_box(s, [
    ('13项问题表面各不相同，但根因指向同一个事实：管理闭环断裂', True, CLR['accent']),
    '',
    ('\U0001f534 断裂点一：制度与执行脱节 — 墙上制度', True, CLR['accent']),
    '不相容职务分离制度写了，执行时一人兼审核+收费+对账 \u2192 问题2',
    '退租交接流程定了，退租后钥匙还在你手里 \u2192 问题8',
    '',
    ('\U0001f534 断裂点二：台账与实物不符 — 纸上台账', True, CLR['accent']),
    '车位备案379个，实际有效320个，四种口径四个数 \u2192 问题1',
    '物业用房5间仅1间入台账，门牌缺失 \u2192 问题9',
    '空调多2台、消火栓少18个，账面与实物差20个 \u2192 问题12',
    '',
    ('\U0001f534 断裂点三：审批与管控虚设 — 橡皮图章', True, CLR['accent']),
    '三人行优惠变1带多，67人中仅9人是老客户 \u2192 问题3',
    '商户从40\u33a1调到66\u33a1无审批无调价 \u2192 问题7',
    '17间房10间没签协议，默认\uff1c配套用房不用签\uff1e \u2192 问题5',
    '',
    ('\U0001f534 断裂点四：监督与闭环缺失 — 形式检查', True, CLR['accent']),
    '维保记录提前签字、两份记录内容截然不同 \u2192 问题12',
    '年度考评走过场，扣1,200元走个形式 \u2192 问题13',
    '351次手动放行无台账无核对无追溯 \u2192 问题2',
], top=emu(1.1), font_size=11.5, line_spacing=1.1)

# ═══════════════════════════════════
# SLIDE 20: THREE-LEVEL IC
# ═══════════════════════════════════
print('Slide 20: Three-Level IC')
s = add_slide(1); add_decor(s)
add_title_box(s, '内控体系的三阶进化', top=emu(0.5), font_size=28)
levels = [
    ('第一阶', '合规型内控', '制度齐全\n满足监管最低要求\n\u201c该有的都有了\u201d', '及格线', CLR['gray']),
    ('第二阶', '有效型内控', '制度落地执行\n流程闭环留痕\n\u201c做了的都有记录\u201d', '优良线', CLR['amber']),
    ('第三阶', '卓越型内控', '数据驱动预警\n主动风险防控\n\u201c没发生的也能预见\u201d', '标杆线', CLR['green']),
]
cw, g = emu(3.8), emu(0.2)
for i, (stage, title, desc, level, clr) in enumerate(levels):
    x = emu(0.7) + i*(cw+g)
    add_label(s, f'{stage}\uff1a{title}', emu(1.3), x, cw, emu(0.5), bg_color=clr, font_size=15)
    add_label(s, level, emu(1.9), x+emu(1.0), emu(1.8), emu(0.35), bg_color=CLR['accent'], font_size=11)
    tb = add_text_box(s, [desc], top=emu(2.5), left=x+emu(0.2), width=cw-emu(0.4), height=emu(1.5), font_size=13)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('\u25b6 现状诊断', True, CLR['title']),
    '天府广场项目管理处于\u201c第一阶\u2192第二阶\u201d之间\u2014\u2014制度有但执行不实',
    ('\u25b6 目标', True, CLR['title']),
    '2027年底前达到第二阶\uff08有效型\uff09\uff0c向第三阶\uff08卓越型\uff09迈进',
], top=emu(4.2), font_size=14, line_spacing=1.3)

# ═══════════════════════════════════
# SLIDE 21: THREE LINES OF DEFENSE
# ═══════════════════════════════════
print('Slide 21: Three Lines of Defense')
s = add_slide(1); add_decor(s)
add_title_box(s, '风险管理的三道防线', top=emu(0.5), font_size=28)
lines_data = [
    ('第一道', '业务部门', '风险所有者', '日常操作中的\n自我管控与自查', '停车场收费员、商户管理员\n维保人员、审批经办人', CLR['blue']),
    ('第二道', '风控/合规部门', '风险管理者', '制度建设、监督检查\n风险预警与报告', '内控部门、财务部门\n法务部门、安质部门', CLR['amber']),
    ('第三道', '审计部门', '独立保证者', '独立评价、问题发现\n改进建议与追责', '内部审计、外部审计\n纪检、巡视', CLR['accent']),
]
cw, g = emu(3.8), emu(0.2)
for i, (num, dept, role, duty, who, clr) in enumerate(lines_data):
    x = emu(0.7) + i*(cw+g)
    add_label(s, f'{num}防线\uff1a{dept}', emu(1.2), x, cw, emu(0.4), bg_color=clr, font_size=13)
    add_text_box(s, [
        (role, True, clr),
        duty,
    ], top=emu(1.8), left=x+emu(0.1), width=cw-emu(0.2), height=emu(1.2), font_size=12, line_spacing=1.2)
    add_text_box(s, [
        ('典型岗位', True, CLR['gray']),
        who,
    ], top=emu(3.2), left=x+emu(0.1), width=cw-emu(0.2), height=emu(1.0), font_size=11, line_spacing=1.2, color=CLR['gray'])
add_text_box(s, [
    ('\u26a0 现状问题', True, CLR['accent']),
    '三道防线之间信息断裂、各自为战\uff1a第一道觉得\u201c反正有人查\u201d\uff0c第二道觉得\u201c等审计来\u201d\uff0c第三道一年只来一次',
    ('\u25b6 改进方向', True, CLR['title']),
    '建立三道防线信息共享与协同机制 \u2192 季度联席风险分析会 \u2192 问题台账动态共享',
], top=emu(4.5), font_size=13, line_spacing=1.15)

# ═══════════════════════════════════
# SLIDE 22: RECOMMENDATIONS (EXPANDED)
# ═══════════════════════════════════
print('Slide 22: Recommendations')
s = add_slide(1); add_decor(s)
add_title_box(s, '综合管理建议', top=emu(0.5), font_size=28)
recs_expanded = [
    ('\U0001f17f 停车场管理', [
        '底数清查\uff1a逐位实测+备案变更\uff0c确保一个口径一个数',
        '岗位分设\uff1a审核/收费/对账/系统录入四岗分离',
        '手动放行管控\uff1a每次放行\u2192即时登记\u2192双人复核\u2192周周核对',
    ], CLR['accent']),
    ('\U0001f3e2 租赁与物业用房', [
        '权属排查\uff1a补充协议全覆盖\uff0c不认\u201c默认\u201d只认签字',
        '面积核验\uff1a定期抽查商户实际使用面积vs合同面积',
        '退租闭环\uff1a钥匙交还+现场核验+双方签字+归档',
    ], CLR['amber']),
    ('\U0001f441 现场运营', [
        '维保标准化\uff1a统一记录模板+拍照留证+双人签字',
        '设备盘点\uff1a季度实物盘点+与台账比对+差异即时上报',
        '考评闭环\uff1a年度考评结果与绩效挂钩+整改复查',
    ], CLR['blue']),
    ('\U0001f504 常态化机制', [
        '月度自查\uff1a各部门按清单自查\uff0c问题即查即改即报',
        '季度联席会\uff1a风控+业务+审计三方数据比对',
        '年度外审\uff1a独立第三方穿透式评价',
    ], CLR['green']),
]
y = emu(1.2)
for title, items, clr in recs_expanded:
    add_label(s, title, y, emu(0.8), emu(3.2), emu(0.4), bg_color=clr, font_size=12)
    add_text_box(s, items, top=y, left=emu(4.2), width=emu(8.0), height=emu(1.0), font_size=11, line_spacing=1.15)
    y += emu(1.1)
add_text_box(s, [('四条建议四个关键词\uff1a分离\u00b7核验\u00b7标准化\u00b7闭环\u00b7协同', True, CLR['title'])],
             top=emu(5.7), font_size=14)

# ═══════════════════════════════════
# SLIDE 23: SIX-STEP (EXPANDED)
# ═══════════════════════════════════
print('Slide 23: Six-Step')
s = add_slide(1); add_decor(s)
add_title_box(s, '制度落地六步法\uff1a从天府广场案例看每一步怎么走', top=emu(0.4), font_size=22)
steps_exp = [
    ('\u2460建制度', '停车场无手动放行制度\uff1f\u2192 立即补建', '制度不在多\uff0c在管用', CLR['blue']),
    ('\u2461明职责', '谁审月租\uff1f谁管收费\uff1f谁核对账\uff1f\u2192 落实到人', '一岗一责\uff0c白纸黑字', CLR['blue']),
    ('\u2462讲培训', '新制度出台谁教过\uff1f\u2192 每人签字确认培训记录', '没培训=没发布', CLR['amber']),
    ('\u2463建台账', '手动放行登记了吗\uff1f退租交接签了吗\uff1f\u2192 留痕', '没记录=没发生', CLR['amber']),
    ('\u2464强检查', '检查不是翻翻本子签字走人\u2192 交叉互查+飞行检查', '检查也是要留痕的', CLR['accent']),
    ('\u2465严问责', '问题整改谁盯\uff1f到期未改谁扛\uff1f\u2192 闭环到人', '没有问责=没有闭环', CLR['accent']),
]
y = emu(1.1)
for action, example, principle, clr in steps_exp:
    add_label(s, action, y, emu(0.7), emu(2.0), emu(0.4), bg_color=clr, font_size=13)
    add_text_box(s, [example], top=y, left=emu(2.9), width=emu(5.5), height=emu(0.35), font_size=12, line_spacing=1.1)
    add_text_box(s, [f'\u2192 {principle}'], top=y+emu(0.05), left=emu(8.5), width=emu(3.5), height=emu(0.35), font_size=12, color=CLR['accent'])
    y += emu(0.75)
add_text_box(s, [
    ('核心结论', True, CLR['title']),
    '打通\u201c最后一公里\u201d的关键\uff0c不在于制度有多少\uff0c而在于执行有多实',
    '六步法的本质\uff1a把\u201c应该做\u201d变成\u201c必须做\u201d\uff0c把\u201c做了\u201d变成\u201c能证明做了\u201d',
], top=emu(5.4), font_size=13, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 24: DIGITAL AUDIT PREP
# ═══════════════════════════════════
print('Slide 24: Digital Audit Prep')
s = add_slide(1); add_decor(s)
add_title_box(s, '数字化审计时代\uff1a企业必须做好的四项准备', top=emu(0.5), font_size=24)
add_text_box(s, [
    ('2026年起\uff0c穿透式智能监管平台全面上线\u3002审计不再是翻凭证\u3001看台账\u2014\u2014', True, CLR['accent']),
    ('而是系统自动抓数据\u3001自动比对\u3001自动预警\u3002企业必须提前准备\u3002', True, CLR['accent']),
    '',
    ('\u2776 财务系统标准化与互联互通', True, CLR['title']),
    '停车场收费系统\u3001租赁管理系统\u3001物业管理系统\u3001财务系统\u2014\u2014四系统必须打通',
    '打破信息孤岛\uff1a一个商户的租金收缴\u3001面积使用\u3001合同变更\uff0c在一个视图里看到全貌',
    '',
    ('\u2777 业务数据全流程数字化留痕', True, CLR['title']),
    '每一笔收费从收\u2192核\u2192存全链路可追溯\uff1b每一次审批从提\u2192审\u2192批有完整时间戳',
    '手动放行\u3001现金收款这类\u201c离线操作\u201d必须纳入系统管控',
    '',
    ('\u2778 建立数据治理机制', True, CLR['title']),
    '数据质量是数字化审计的生命线\uff1a录错了\u3001漏了\u3001改了\u2014\u2014都会被自动标记',
    '指定数据责任人+数据录入标准+异常数据自动报警',
    '',
    ('\u2779 培养全员的数字化思维', True, CLR['title']),
    '不是IT部门的事\uff0c是每个岗位的事\u3002你在系统里录的每一条数据\uff0c都会被审计追踪',
    '\u201c以前都这么干的\u201d在数字化审计面前不再是免责理由',
], top=emu(1.1), font_size=11, line_spacing=1.1)

# ═══════════════════════════════════
# SLIDE 25: PROACTIVE DEFENSE
# ═══════════════════════════════════
print('Slide 25: Proactive Defense')
s = add_slide(1); add_decor(s)
add_title_box(s, '从\u201c被动整改\u201d到\u201c主动防控\u201d\uff1a管理思维的跃迁', top=emu(0.5), font_size=24)
add_text_box(s, [
    ('传统模式\uff08被动\uff09\uff1a', True, CLR['accent']),
    '审计发现问题 \u2192 写整改报告 \u2192 下次审计再发现新问题 \u2192 再写整改报告',
    '永远在追着问题跑\uff0c永远在\u201c不及格\u2192及格\u2192又不及格\u201d的循环中',
    '',
    ('目标模式\uff08主动\uff09\uff1a', True, CLR['green']),
    '风险自查 \u2192 自动预警 \u2192 即时纠偏 \u2192 审计确认 \u2192 制度优化',
    '问题在萌芽阶段被识别和解决\uff0c审计变成\u201c确认\u201d而非\u201c发现\u201d',
    '',
    ('实现路径\uff1a三件事', True, CLR['title']),
    '',
    ('第一件\uff1a建立部门风险自查清单', True, CLR['dark']),
    '每个业务部门有一张\u201c可能出问题的清单\u201d\u2014\u2014对标今天的13项问题',
    '停车场\uff1a车位变动登记了吗\uff1f手动放行有没有记录\uff1f优惠政策有没有超范围\uff1f',
    '',
    ('第二件\uff1a季度风险扫描', True, CLR['dark']),
    '每季度由第二道防线\uff08风控/合规\uff09牵头\uff0c用\u201c审计视角\u201d扫一遍关键流程',
    '发现问题 \u2192 即时整改 \u2192 不再等到外部审计来',
    '',
    ('第三件\uff1a问题台账动态清零', True, CLR['dark']),
    '建立\u201c问题库\u201d\u2014\u2014每一项问题有编号\u3001有责任人\u3001有整改期限',
    '整改一项销号一项\uff0c超期未销号自动升级到追责',
    '',
    ('一句话\uff1a最好的审计结果是\u201c审不出大问题\u201d\u2014\u2014不是查不到\uff0c是真没问题', True, CLR['accent']),
], top=emu(1.1), font_size=11, line_spacing=1.05)

# ═══════════════════════════════════
# SLIDE 26: ACTION ROADMAP
# ═══════════════════════════════════
print('Slide 26: Roadmap')
s = add_slide(1); add_decor(s)
add_title_box(s, '2026-2027 行动路线图\uff1a从天府广场整改到监管达标', top=emu(0.5), font_size=24)
phases = [
    ('2026\nQ3', '制度补缺\n与整改', '完成13项问题整改\n补建缺失制度\n完成停车位备案变更', CLR['accent']),
    ('2026\nQ4', '流程重构\n与培训', '不相容职务分离落地\n全员制度培训\n台账模板标准化', CLR['amber']),
    ('2027\nQ1', '系统打通\n与数据治理', '四大系统互联互通\n历史数据清洗\n数据责任人指定', CLR['blue']),
    ('2027\nQ2', '机制常态化\n与自检', '季度风险扫描机制运行\n自查清单迭代优化\n问题台账动态清零', CLR['blue']),
    ('2027\nQ3-Q4', '迎接验收\n对标达标', '穿透式监管平台对接\n第三方独立评价\n达到内控二阶标准', CLR['green']),
]
cw, g = emu(2.3), emu(0.1)
for i, (time, title, desc, clr) in enumerate(phases):
    x = emu(0.5) + i*(cw+g)
    add_label(s, time, emu(1.2), x, cw, emu(0.65), bg_color=clr, font_size=11)
    add_text_box(s, [(title, True, clr)], top=emu(2.0), left=x, width=cw, height=emu(0.35), font_size=12)
    s.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_box(s, [desc], top=emu(2.5), left=x+emu(0.05), width=cw-emu(0.1), height=emu(2.0), font_size=10.5, line_spacing=1.15)
add_text_box(s, [
    ('关键节点', True, CLR['title']),
    '2027年底\uff1a四川省国资委对省级国企穿透式监管平台达标验收\u2014\u201460%省属企业必须达标',
    '天府广场项目作为轨道公司核心商业资产\uff0c应作为首批达标单位',
], top=emu(4.8), font_size=13, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 27: KEY INSIGHTS
# ═══════════════════════════════════
print('Slide 27: Key Insights')
s = add_slide(1); add_decor(s)
add_title_box(s, '关键启示', top=emu(0.5), font_size=28)
add_text_box(s, [
    ('三句话', True, CLR['title']),
    '',
    '\u2776 内控不是墙上的制度\uff0c是做出来的',
    '    \u2014\u2014天府广场13项问题\uff0c每一项都有制度\uff0c每一项都没执行到位',
    '',
    '\u2777 台账不是纸上的记录\uff0c是出了事能查到的',
    '    \u2014\u2014351次手动放行无记录=出了事无人能答\uff0c这不是侥幸\uff0c是隐患',
    '',
    '\u2778 审批不是流程走过场\uff0c是责任落实到人',
    '    \u2014\u201467人的\u201c三人行\u201d变成\u201c1带多\u201d\uff0c审批松一松\uff0c后果自己扛',
    '',
    '',
    ('四位一体框架\u2014\u2014从\u201c知道\u201d到\u201c做到\u201d的完整逻辑', True, CLR['title']),
    '',
    '1号文 \u2192 用系统减少\u201c人\u201d的随意性       2号文 \u2192 用框架消除\u201c管\u201d的盲区',
    '15号文 \u2192 用执行填补\u201c制\u201d与\u201c行\u201d之间的鸿沟   46号令 \u2192 用问责守住该守的底线',
], top=emu(1.2), font_size=14, line_spacing=1.15)

# ═══════════════════════════════════
# SLIDE 28: SELF-CHECK
# ═══════════════════════════════════
print('Slide 28: Self-Check')
s = add_slide(1); add_decor(s)
add_title_box(s, '合规自查要点\uff08课后作业\uff09', top=emu(0.5), font_size=28)
checks = [
    ('对照1号文', '财务系统打通了吗\uff1f预警自动了吗\uff1f'),
    ('对照2号文', '穿透到底了吗\uff1f资金可追溯了吗\uff1f'),
    ('对照15号文', '内控制度全覆盖了吗\uff1f整改闭环了吗\uff1f'),
    ('对照46号令', '红线全员知晓了吗\uff1f容错边界明确了吗\uff1f'),
]
y = emu(1.4)
for label, question in checks:
    add_label(s, label, y, emu(0.8), emu(2.5), emu(0.4), bg_color=CLR['title'], font_size=13)
    add_text_box(s, [question], top=y+emu(0.05), left=emu(3.5), width=emu(8.5), height=emu(0.4), font_size=16)
    y += emu(0.85)
add_text_box(s, [
    '',
    ('核心原则八字诀', True, CLR['title']),
    '制度有 \u00b7 执行实  |  台账有 \u00b7 信息全  |  流程有 \u00b7 留痕清  |  审批严 \u00b7 追责准',
], top=emu(4.5), font_size=18, line_spacing=1.3)

# ═══════════════════════════════════
# SLIDE 29: CLOSING
# ═══════════════════════════════════
print('Slide 29: Closing')
s = add_slide(1); add_decor(s)
add_title_box(s, '结语', top=emu(0.5), font_size=28)
add_text_box(s, [
    ('送大家三句话', True, CLR['title']),
    '',
    ('\U0001f50d 每一次审计都是一次体检', True, CLR['dark']),
    '体检不是为难你\uff0c是帮你发现隐患\u2014\u2014早发现\uff0c早治疗\uff0c早安心',
    '',
    ('\U0001fa9e 每一个问题都是一面镜子', True, CLR['dark']),
    '照出来的不是你一个人的问题\uff0c是整个管理体系的问题\u2014\u2014改一个点\uff0c堵一个面',
    '',
    ('\U0001f4c8 每一项整改都是一次升级', True, CLR['dark']),
    '把问题改到位了\uff0c你的管理水平就上了一个台阶\u2014\u2014不是应付检查\uff0c是提升自己',
    '',
    '',
    ('今天开始\uff0c从\u201c不会查到我\u201d到\u201c查了我也不怕\u201d', True, CLR['accent']),
    '\u2014\u2014这就是\u201c全员审计风险意识\u201d的真正含义',
], top=emu(1.2), font_size=15, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 30: THANK YOU
# ═══════════════════════════════════
print('Slide 30: Thank You')
s = add_slide(1); add_decor(s)
add_title_box(s, '谢谢大家', top=emu(1.5), font_size=48, color=CLR['title'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
tb = add_text_box(s, [
    '\u201c提升全员审计风险意识\u201d\u2014\u2014核心词不是\u201c审计\u201d\uff0c是\u201c全员\u201d',
    '审计团队来查是一年一次\uff0c但风险发生是每时每刻',
    '真正防住风险的人\uff0c是你们每一个业务岗位上的每一个人',
    '',
    '四川融策会计师事务所  |  2026年6月',
    '欢迎会后交流提问',
], top=emu(3.0), font_size=16, color=CLR['dark'], left=emu(1.5), width=SW-emu(3.0))
for p in tb.text_frame.paragraphs:
    p.alignment = PP_ALIGN.CENTER

# ── Save ──
OUT = r'D:\\openclaw-workspace\\output\\v5_expanded.pptx'
print(f'\\nSaving 30 slides to {OUT}...')
prs.save(OUT)
print(f'Saved to {OUT}')
