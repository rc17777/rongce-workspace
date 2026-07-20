#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""终版修复：不增不减段落，只替换文字"""
import sys,os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

src = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_616.pptx'
dst = r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx'

prs = Presentation(src)

# ========== Slide 4: 1号文→14号文 ==========
# 原TextBox7有8个段落(P0-P7)，保持不变，只替换文字
slide4_new = [
    '▎追责范围 — 13大类全覆盖：集团管控·风险管理·购销·工程·金融·科技·资金·担保·产权·投资·改制·资产管理·境外投资',
    '▎追责方式（可合并使用）：批评诫勉 / 组织处理（停职·免职·降职）/ 扣减薪酬 / 禁入限制(5年~终身) / 纪律处分 / 移送司法',
    '▎损失认定三档：一般损失<100万 | 较大损失100万~1000万 | 重大损失>1000万',
    '',
    '▎三条硬杠杠',
    '🔴 重大决策终身问责 — 退休≠安全，离职≠免责      🟢 集体决策中明确反对意见→免予责任',
    '⚠ 损失未达标准但导致企业无法持续经营 → 视为重大损失',
    '▶ 与你的关系：停车场收费不清、商户租金拖欠、物业用房违规——天府广场13项问题全在追责范围内',
]

slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.name == 'TextBox 5':  # 标题
        for run in shape.text_frame.paragraphs[0].runs:
            if '1号文' in run.text:
                run.text = '成国资发〔2026〕14号：《成都市属国企违规经营投资责任追究实施办法》'
    elif shape.name == 'Rectangle 6':  # 标签
        for run in shape.text_frame.paragraphs[0].runs:
            if '技术底座' in run.text:
                run.text = '市级配套'
    elif shape.name == 'TextBox 7':  # 正文
        tf = shape.text_frame
        for i in range(min(len(slide4_new), len(tf.paragraphs))):
            if tf.paragraphs[i].runs:
                tf.paragraphs[i].runs[0].text = slide4_new[i]

print('Slide 4 OK')

# ========== Slide 5: 2号文→15号文 ==========
# 原TextBox7有11个段落(P0-P10)
slide5_new = [
    '▎核心要求：公开招租 — 市政府网+市国资委网+产权交易所+集团官网同步发布',
    '正式披露≥10个工作日 | 结果公示≥3个工作日',
    '▎租金定价机制：评估/估价/询价/公开市场价为依据，综合市场需求、区位布局、功能业态确定',
    '',
    '▎五条红线',
    '❌ 拆分资产规避公开招租   ❌ 未经批准擅自出租   ❌ 擅自改变用途或租赁期限',
    '❌ 截留/挪用/私分租金收入   ❌ 违规转租（有违约行为不得续租）',
    '',
    '▎监管机制："三重一大"集体决策→报市国资委备案→年度自查自纠→审计+专项检查',
    '▶ 与你的关系',
    '天府广场商业租赁的合同续租、商户转租、租金定价——全部在这个文件里画了红线',
]

slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.name == 'TextBox 5':
        for run in shape.text_frame.paragraphs[0].runs:
            if '2号文' in run.text:
                run.text = '成国资发〔2025〕15号：市属国企资产租赁管理制度'
    elif shape.name == 'Rectangle 6':
        for run in shape.text_frame.paragraphs[0].runs:
            if '总纲领' in run.text:
                run.text = '业务直击'
    elif shape.name == 'TextBox 7':
        tf = shape.text_frame
        for i in range(len(slide5_new)):
            if i < len(tf.paragraphs) and tf.paragraphs[i].runs:
                tf.paragraphs[i].runs[0].text = slide5_new[i]

print('Slide 5 OK')

# ========== Slide 27: 关键启示 ==========
slide27 = prs.slides[26]
for shape in slide27.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text
            if '1号文 →' in t and '2号文 →' in t:
                for run in para.runs:
                    if '1号文' in run.text:
                        run.text = run.text.replace('1号文 → 用系统减少"人"的随意性', '14号文 → 13大类追责全覆盖, 终身问责不设限')
                    if '2号文 → 用框架消除"管"的盲区' in run.text:
                        run.text = run.text.replace('2号文 → 用框架消除"管"的盲区', '15号文 → 公开招租五条红线, 画清合规边界')
print('Slide 27 OK')

# ========== Slide 28: 合规自查 ==========
slide28 = prs.slides[27]
for shape in slide28.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == '对照1号文':
                    run.text = '对照14号文'
                elif run.text.strip() == '信息系统打通了吗？预警自动了吗？':
                    run.text = '追责红线画清了吗？投融资合规吗？'
                elif run.text.strip() == '对照2号文':
                    run.text = '对照15号租赁文'
                elif run.text.strip() == '穿透到底了吗？资金可追溯了吗？':
                    run.text = '公开招租执行了吗？转租合规了吗？'
print('Slide 28 OK')

# ========== Other slides ==========
# Slide 2
for shape in prs.slides[1].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('1号文·2号文·', '')
                run.text = run.text.replace('国务院核心文件解读', '国务院及成都市政策解读')

# Slide 3
for shape in prs.slides[2].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '1号文 → 2号文' in run.text:
                    run.text = '15号文 → 46号令 → 成都市配套制度：国家顶层设计 × 地方落地执行'

# Slide 8
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == '1号文': run.text = '14号(追责)'
                elif run.text.strip() == '2号文': run.text = '15号(租赁)'
                elif run.text.strip() == '看见': run.text = '定位'
                elif run.text.strip() == '看准': run.text = '规范'
                elif '数智化系统让异常无处遁形' in run.text: run.text = '13大类追责全覆盖 | 终身追责'
                elif '穿透监管框架界定查什么、查多深' in run.text: run.text = '公开招租 | 五条红线 | 全流程规范'
                elif '四位一体' in run.text: run.text = run.text.replace('四位一体', '六位一体')
                elif '技术 ──→ 制度' in run.text: run.text = '15号文→46号令→14号→15号  =  国家·市级双层闭环'

print('Slides 2,3,8 OK')

prs.save(dst)
print(f'\nDone: {dst}')
print(f'Size: {os.path.getsize(dst)//1024}KB')
