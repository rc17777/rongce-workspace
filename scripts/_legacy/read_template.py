import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'C:\Users\scrccpa\Desktop\物资管理制度培训-2025.1.20.pptx')
print(f'幻灯片尺寸: {prs.slide_width} x {prs.slide_height} ({prs.slide_width/914400:.1f}x{prs.slide_height/914400:.1f}英寸)')
print(f'幻灯片数: {len(prs.slides)}')
print()

print('=== Slide Layouts ===')
for li, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {li}: "{layout.name}"')
    for ph in layout.placeholders:
        print(f'    PH idx={ph.placeholder_format.idx} name="{ph.name}" type={ph.placeholder_format.type}')

print()
print('=== 每页内容 ===')
for si, slide in enumerate(prs.slides):
    print(f'--- 第{si+1}页 (layout: "{slide.slide_layout.name}") ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text[:200].replace('\n', ' | ')
            print(f'  [文本] {shape.name}: "{text}"')
        elif shape.shape_type == 13:
            print(f'  [图片] {shape.name} ({shape.width/914400:.1f}x{shape.height/914400:.1f}英寸)')
