#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""撤掉配图 → 换极简装饰元素"""
import sys,os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation(r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx')

# ========== 1. 删除第24/25/27/28/29页新加的图片 ==========
for slide_idx in [23, 24, 26, 27, 28]:
    slide = prs.slides[slide_idx]
    # 找到最后添加的图片shape（Picture类型，且名称是'Picture'开头的，也是最大的那个）
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13 and shape.name.startswith('Picture') and int(shape.name.split()[-1]) >= 10:  # PICTURE type
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    print(f'Slide {slide_idx+1}: removed {len(shapes_to_remove)} images')

# ========== 2. 添加极简装饰：右侧淡金细竖线 ==========
# 使用和封面一样的金线装饰风格
GOLD = RGBColor(0xC5, 0x95, 0x5C)
GOLD_DIM = RGBColor(0x3D, 0x2D, 0x1C)  # 暗金（低调）

for slide_idx in [23, 24, 26, 27, 28]:
    slide = prs.slides[slide_idx]
    
    # 右侧淡金竖线装饰（3px宽，从顶部到底部）
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(32.2), Cm(2.0), Cm(0.08), Cm(14)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()  # 无边框
    # 设置透明度（如果是支持的）
    try:
        # 尝试设置透明度
        from pptx.oxml.ns import qn
        solidFill = line.fill._fill
        srgb = solidFill.find(qn('a:solidFill'))
        if srgb is not None:
            alpha = srgb.find(qn('a:alpha'))
            if alpha is None:
                alpha = srgb.makeelement(qn('a:alpha'), {'val': '30000'})
                srgb.append(alpha)
    except:
        pass
    
    # 右侧底部淡金小横线（装饰收尾）
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(30.0), Cm(16.5), Cm(2.5), Cm(0.06)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = GOLD
    line2.line.fill.background()
    
    # 右上角小三角装饰
    tri = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Cm(31.5), Cm(1.0), Cm(1.2), Cm(0.8)
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x4A)  # 比背景稍亮的深蓝
    tri.line.fill.background()
    
    print(f'Slide {slide_idx+1}: added subtle decorations')

prs.save(r'D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx')
print(f'\nDone: {os.path.getsize(r"D:\openclaw-workspace\temp\轨道培训意见征求稿_最终版.pptx")//1024}KB')
