"""
v5b: 全新方案 — 从模板加载，用BytesIO保存绕过文件锁，重写Content_Types
"""
import sys, os, shutil, io, zipfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ── Paths ──
TMPL = r'C:\Users\scrccpa\Desktop\物资管理制度培训-2025.1.20.pptx'
OUT  = r'D:\openclaw-workspace\output\v5_final.pptx'
DESK = r'C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v5_模板.pptx'
ASSETS = r'D:\openclaw-workspace\scripts\ppt_assets'

# ── Colors ──
CLR = {
    'dark':     RGBColor(0x33, 0x33, 0x33),
    'title':    RGBColor(0x1A, 0x3C, 0x6E),
    'accent':   RGBColor(0xC0, 0x39, 0x2B),
    'blue':     RGBColor(0x29, 0x80, 0xB9),
    'green':    RGBColor(0x27, 0xAE, 0x60),
    'amber':    RGBColor(0xE6, 0x7E, 0x22),
    'gray':     RGBColor(0x7F, 0x8C, 0x8D),
    'white':    RGBColor(0xFF, 0xFF, 0xFF),
}

FONT = '微软雅黑'
FONT_TITLE = '微软雅黑'
SW = Emu(12192000)
SH = Emu(6858000)

# ── Assets ──
BANNER = os.path.join(ASSETS, 'slide2_图片 3.png')
LOGO   = os.path.join(ASSETS, 'slide2_图片 8.png')
COVER_BG1 = os.path.join(ASSETS, 'slide1_图片 1.png')
COVER_BG2 = os.path.join(ASSETS, 'slide1_图片 2.png')
COVER_LOGO = os.path.join(ASSETS, 'slide1_图片 7.png')

def emu(inches):
    return int(inches * 914400)

# ── Load & clean ──
print('Loading template...')
prs = Presentation(TMPL)

# Remove all slides by deleting from sldIdLst AND their rels
sldIdLst = prs.slides._sldIdLst
to_remove = list(sldIdLst)
for sldId in to_remove:
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    sldIdLst.remove(sldId)
    if rId:
        try:
            prs.part.drop_rel(rId)
        except:
            pass

print(f'Slides after cleanup: {len(prs.slides)}')

# ── Helpers ──
def add_slide(layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_decor(slide):
    slide.shapes.add_picture(BANNER, 0, SH - emu(1.3), SW, emu(1.3))
    slide.shapes.add_picture(LOGO, emu(0.4), emu(0.3), emu(1.4), emu(0.6))

def add_title_box(slide, text, top=emu(0.5), left=emu(0.8), width=None, height=emu(0.8),
                  font_size=30, color=None, bold=True, align=PP_ALIGN.LEFT):
    if width is None: width = SW - emu(1.6)
    if color is None: color = CLR['title']
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.name = FONT_TITLE; p.font.size = Pt(font_size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return txBox

def add_text_box(slide, lines, top, left=emu(0.8), width=None, height=None,
                 font_size=16, color=None, line_spacing=1.5):
    if width is None: width = SW - emu(1.6)
    if height is None: height = SH - top - emu(1.8)
    if color is None: color = CLR['dark']
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, str):
            p.text = line; p.font.bold = False; p.font.color.rgb = color
        elif isinstance(line, tuple):
            p.text = line[0]
            p.font.bold = line[1] if len(line) > 1 else False
            p.font.color.rgb = line[2] if len(line) > 2 else color
        p.font.name = FONT; p.font.size = Pt(font_size)
        p.space_after = Pt(4); p.level = 0
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
        spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox

def add_label(slide, text, top, left, width, height, bg_color=CLR['blue'], font_size=14, font_color=None):
    if font_color is None: font_color = CLR['white']
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = bg_color; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.name = FONT; p.font.size = Pt(font_size)
    p.font.bold = True; p.font.color.rgb = font_color; p.alignment = PP_ALIGN.CENTER
    return shape

# ═══════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════
print('Slide 1: Cover')
s = add_slide(0)
s.shapes.add_picture(COVER_BG1, emu(-0.0), emu(0.9), emu(13.3), emu(6.6))
s.shapes.add_picture(COVER_BG2, emu(-0.3), emu(0.9), emu(13.9), emu(6.7))
s.shapes.add_picture(COVER_LOGO, emu(0.4), emu(0.3), emu(1.4), emu(0.6))
add_title_box(s, '提升公司全员审计风险意识', top=emu(2.0), font_size=40, color=CLR['white'],
              align=PP_ALIGN.CENTER, width=SW - emu(2.0), left=emu(1.0))
tb = add_text_box(s, ['四川融策会计师事务所  |  2026年6月  |  四川轨道公司专题培训'],
                  top=emu(3.5), font_size=18, color=RGBColor(0xDD,0xDD,0xDD),
                  left=emu(1.0), width=SW - emu(2.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════
print('Slide 2: Agenda')
s = add_slide(1); add_decor(s)
add_title_box(s, '培训议程', top=emu(0.6))
items = [
    ('第一部分', '国务院核心文件解读（1号文·2号文·15号文·46号令）', '⏱ 40分钟', CLR['blue']),
    ('第二部分', '天府广场审计实战案例（4大领域·13项问题）', '⏱ 45分钟', CLR['accent']),
    ('第三部分', '总结与建议（制度落地六步法·合规自查）', '⏱ 5分钟', CLR['green']),
]
y = emu(1.8)
for label, desc, time, clr in items:
    add_label(s, label, y, emu(1.0), emu(2.0), emu(0.5), bg_color=clr, font_size=14)
    add_text_box(s, [desc], top=y+emu(0.05), left=emu(3.3), width=emu(7.5), height=emu(0.5), font_size=18)
    add_text_box(s, [time], top=y+emu(0.05), left=emu(11.0), width=emu(1.8), height=emu(0.5), font_size=14, color=CLR['gray'])
    y += emu(1.1)

# ═══════════════════════════════════
# SLIDE 3: PART 1 DIVIDER
# ═══════════════════════════════════
print('Slide 3: Part 1 Divider')
s = add_slide(2); add_decor(s)
add_title_box(s, '第一部分', top=emu(2.0), font_size=44, color=CLR['blue'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
add_title_box(s, '国务院核心文件解读', top=emu(3.2), font_size=28, color=CLR['dark'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0), bold=False)
tb = add_text_box(s, ['1号文 → 2号文 → 15号文 → 46号令：从技术底座到兜底保障的完整政策链条'],
                  top=emu(4.2), font_size=16, color=CLR['gray'], left=emu(1.5), width=SW-emu(3.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 4: DOC #1
# ═══════════════════════════════════
print('Slide 4: Doc #1')
s = add_slide(1); add_decor(s)
add_title_box(s, '1号文：《推动央企加快财务数智化转型升级》', top=emu(0.6), font_size=26)
add_label(s, '技术底座', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['blue'])
add_text_box(s, [
    ('▎建系统', True, CLR['title']),
    '统一财务数智化平台，打通数据孤岛。停车场收费、商业租赁、物业运营、设备维保——系统不通，数据就是一座座孤岛',
    ('▎强支撑', True, CLR['title']),
    '数据治理、安全防护、人才保障。系统建好了数据质量跟不上，后续分析决策都跑偏',
    ('▎给工具', True, CLR['title']),
    '部署智能化财务分析工具，实现风险自动预警',
    ('▶ 与你的关系', True, CLR['accent']),
    '每一笔收费、每一份合同、每一次审批，将来都要能被系统追踪。人在做，系统在看，审计在复盘',
], top=emu(1.9), font_size=15, line_spacing=1.25)

# ═══════════════════════════════════
# SLIDE 5: DOC #2
# ═══════════════════════════════════
print('Slide 5: Doc #2')
s = add_slide(1); add_decor(s)
add_title_box(s, '2号文：《加强央企穿透式监管指导意见（试行）》', top=emu(0.6), font_size=26)
add_label(s, '总纲领', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['blue'])
add_text_box(s, [
    ('四层穿透', True, CLR['title']),
    ('① 层级穿透', True, CLR['dark']), '穿透股权结构到底层实体。子公司、代管单位、合作经营方——监管一穿到底',
    ('② 资金穿透', True, CLR['dark']), '追踪资金流向全链条。谁收的、怎么收的、进了哪个账户、什么时候入的账',
    ('③ 业务穿透', True, CLR['dark']), '穿透业务全流程。商户实际经营面积是否超合同范围？用制度和数据把缝隙堵上',
    ('④ 追责穿透', True, CLR['dark']), '穿透责任链条至个人。出问题不再是"集体负责"，具体的人要扛',
    ('▶ 与你的关系', True, CLR['accent']), '监管的天网正在织密，以前管不到的灰色地带，以后都会被穿透',
], top=emu(1.9), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 6: DOC #15
# ═══════════════════════════════════
print('Slide 6: Doc #15')
s = add_slide(1); add_decor(s)
add_title_box(s, '15号文：《2026年央企内控体系建设与监督工作通知》', top=emu(0.6), font_size=26)
add_label(s, '操作手册', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['amber'])
add_text_box(s, [
    ('五大核心要求', True, CLR['title']),
    '❶ 完善内控制度体系 — 不相容职务分离、手动放行审批流程，有没有？',
    '❷ 强化重大风险防控 — 收费异常、商户拖欠租金…能不能被及时发现？',
    '❸ 落实监督评价机制 — 是自查自评走过场，还是真正过硬的第三方评价？',
    '❹ 推进数智化赋能 — 系统不行就要建，数据不通就要通，不是"将就"是"建设"',
    '❺ 建立整改问责闭环 — 审计发现问题，限期整改，整改不力问责',
    ('▶ 与你的关系', True, CLR['accent']), '每一条要求，都能在天府广场审计报告中找到对应的反面案例',
], top=emu(1.9), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 7: ORDER #46
# ═══════════════════════════════════
print('Slide 7: Order #46')
s = add_slide(1); add_decor(s)
add_title_box(s, '46号令：《违规经营投资责任追究制度》', top=emu(0.6), font_size=26)
add_label(s, '兜底保障', emu(1.3), emu(0.8), emu(2.0), emu(0.45), bg_color=CLR['accent'])
add_text_box(s, [
    ('三句话总结：程序合规 · 尽职免责 · 失职追责', True, CLR['accent']),
    ('🔴 明底线', True, CLR['dark']), '违规行为红线画清楚。资产损失、违规经营、内控失效、隐瞒不报——白纸黑字',
    ('🔴 严追责', True, CLR['dark']), '资产损失必问责，重大决策终身追责。今天签的审批，十年后出问题，追责链条不断',
    ('🟢 宽容错', True, CLR['dark']), '改革创新中的探索失误、非主观过失、履行了必要程序且未谋私利——可从轻或免予追责',
    '不是让你不敢干事，是让你干事走程序',
], top=emu(1.9), font_size=14, line_spacing=1.25)

# ═══════════════════════════════════
# SLIDE 8: PANORAMA
# ═══════════════════════════════════
print('Slide 8: Panorama')
s = add_slide(1); add_decor(s)
add_title_box(s, '四位一体：政策全景图', top=emu(0.6))
cols = [
    ('1号文', '看见', '数智化系统\n让异常无处遁形', CLR['blue']),
    ('2号文', '看准', '穿透监管框架\n界定查什么、查多深', CLR['blue']),
    ('15号文', '做到', '内控操作手册\n把制度变成行动', CLR['amber']),
    ('46号令', '守住', '追责底线\n让违规必付代价', CLR['accent']),
]
cw, g = emu(2.7), emu(0.3)
for i, (label, verb, desc, clr) in enumerate(cols):
    x = emu(0.7) + i*(cw+g)
    add_label(s, label, emu(1.8), x+emu(0.3), emu(2.0), emu(0.6), bg_color=clr, font_size=22)
    tb = add_text_box(s, [verb], top=emu(2.6), left=x, width=cw, height=emu(0.5), font_size=28, color=clr)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = add_text_box(s, [desc], top=emu(3.2), left=x+emu(0.1), width=cw-emu(0.2), height=emu(1.2), font_size=14)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
tb = add_text_box(s, ['技术 ──→ 制度 ──→ 执行 ──→ 问责  =  完整闭环'],
                  top=emu(4.6), font_size=18, color=CLR['title'], left=emu(1.0), width=SW-emu(2.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 9: ORDER #46 DEEP DIVE
# ═══════════════════════════════════
print('Slide 9: Order #46 Deep Dive')
s = add_slide(1); add_decor(s)
add_title_box(s, '46号令深度解读：追责倒查链', top=emu(0.6), font_size=26)
add_label(s, '🆕 重点', emu(1.2), emu(0.8), emu(1.5), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('🔴 13大类98种违规情形', True, CLR['accent']),
    '投资决策失误 · 产权交易违规 · 工程超概 · 围标串标 · 资金挪用 · 违规签约',
    '财务造假 · 违规担保 · 境外投资失控 · 薪酬福利违规……白纸黑字列出来',
    ('🔴 追责对象全覆盖', True, CLR['accent']),
    '决策人（签字）→ 执行人（经办）→ 审核人（审核）→ 监管人（监督）→ 退休人员',
    '退休 ≠ 安全，离职 ≠ 免责',
    ('🔴 追责时间轴', True, CLR['accent']),
    '重大违规：无期限、终身追责 | 较大多数：倒查20~30年 | 一般违规：倒查10~15年',
    ('⚠ 今天签的字、经手的事，20年后可能被翻出来追责——不是吓唬人，是白纸黑字', True, CLR['accent']),
], top=emu(1.8), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 10: ROADMAP
# ═══════════════════════════════════
print('Slide 10: Roadmap')
s = add_slide(1); add_decor(s)
add_title_box(s, '穿透式监管落地时间表', top=emu(0.6))
milestones = [
    ('2026', '央企全面达标', '100%建成穿透式\n智能监管平台\n十大领域全覆盖', CLR['accent']),
    ('2027', '省级国企推广', '60%以上省属国企\n建成监管平台\n大监督协同机制', CLR['amber']),
    ('2028', '地市县区覆盖', '穿透式监管\n从央企到省属到市属\n到区县全部覆盖', CLR['blue']),
]
cw, g = emu(3.6), emu(0.4)
for i, (year, title, desc, clr) in enumerate(milestones):
    x = emu(0.7) + i*(cw+g)
    add_label(s, year, emu(1.6), x, emu(1.8), emu(0.55), bg_color=clr, font_size=24)
    tb = add_text_box(s, [title], top=emu(2.3), left=x, width=cw, height=emu(0.5), font_size=16, color=CLR['title'])
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = add_text_box(s, [desc], top=emu(2.9), left=x+emu(0.2), width=cw-emu(0.4), height=emu(1.5), font_size=13)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('四大支柱', True, CLR['title']),
    '制度：统一监管制度+标准化规程库    数据：数据治理五大域、打破孤岛',
    '工具：智能监管平台、人防→技防→智防    协同：大监督+内控协同+智能预警',
    ('▶ 四川轨道公司作为省级国企，2027年就是截止日期——不是"未来要做"，是"明年完成"', True, CLR['accent']),
], top=emu(4.5), font_size=14, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 11: PART 2 DIVIDER
# ═══════════════════════════════════
print('Slide 11: Part 2 Divider')
s = add_slide(2); add_decor(s)
add_title_box(s, '第二部分', top=emu(2.0), font_size=44, color=CLR['accent'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
add_title_box(s, '天府广场独立商业项目·审计实战', top=emu(3.2), font_size=28, color=CLR['dark'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0), bold=False)
tb = add_text_box(s, ['审计期间：2024.01 — 2026.03 ｜ 4大领域 · 13项问题'],
                  top=emu(4.2), font_size=16, color=CLR['gray'], left=emu(1.5), width=SW-emu(3.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 12: AUDIT OVERVIEW
# ═══════════════════════════════════
print('Slide 12: Audit Overview')
s = add_slide(1); add_decor(s)
add_title_box(s, '审计概况', top=emu(0.6))
domains = [
    ('🅿 停车场管理', '车位底数、收费内控\n减免审批、计费规则', CLR['accent']),
    ('🏢 经营管理', '房间使用、租赁范围\n调换审批、退租交接', CLR['amber']),
    ('🏠 物业用房', '台账管理、内部改造\n外单位使用', CLR['blue']),
    ('👁 现场管理', '维保记录、设备台账\n年度考评', CLR['green']),
]
cw, g = emu(2.7), emu(0.3)
for i, (title, desc, clr) in enumerate(domains):
    x = emu(0.7) + i*(cw+g)
    add_label(s, title, emu(1.5), x, cw, emu(0.55), bg_color=clr, font_size=15)
    tb = add_text_box(s, [desc], top=emu(2.2), left=x+emu(0.1), width=cw-emu(0.2), height=emu(1.2), font_size=13)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('13项问题 · 涉及金额22,671元', True, CLR['title']),
    '金额不算大，但问题的性质比金额重要——停车场内控缺失、物业用房台账失范，都是系统漏洞而非个案偶然',
    '今天22,671元，明天如果漏洞没堵上，可能就是后面加两个零',
], top=emu(3.8), font_size=15, line_spacing=1.3)

# ═══════════════════════════════════
# SLIDE 13: PARKING I
# ═══════════════════════════════════
print('Slide 13: Parking I')
s = add_slide(1); add_decor(s)
add_title_box(s, '停车场管理问题（一）', top=emu(0.6), font_size=26)
add_label(s, '问题1', emu(1.4), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('停车位底数不清，账实不符，设施失效', True, CLR['accent']),
    '备案379个 → 系统330个 → 画线331个 → 有效320个  四个数字，同一个停车场！',
    '现场停车感应装置基本全部损坏。备案后无动态更新、系统数据无定期核对、硬件坏了没人修',
], top=emu(1.4), left=emu(2.2), width=emu(9.5), height=emu(1.5), font_size=14, line_spacing=1.3)
add_label(s, '问题2', emu(3.2), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('停车场内控不健全 ⚠ 13项问题中性质最严重', True, CLR['accent']),
    '❶ 不相容职务未分离：同一人管审核+收费+对账+系统录入+充值+撤销——三权集于一身',
    '❷ 手动放行管控薄弱：351次手动抬杆，绝大部分是临时车辆',
    '❸ 三缺：抬杆记录台账缺失 + 保安收款记录缺失 + 收费员核对记录缺失',
    ('灵魂拷问：351次手动放行，应收多少、实收多少、有没有少交的——谁能回答？', True, CLR['accent']),
], top=emu(3.2), left=emu(2.2), width=emu(9.5), height=emu(2.5), font_size=13, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 14: PARKING II
# ═══════════════════════════════════
print('Slide 14: Parking II')
s = add_slide(1); add_decor(s)
add_title_box(s, '停车场管理问题（二）', top=emu(0.6), font_size=26)
add_label(s, '问题3', emu(1.4), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('减免优惠车辆审批不严，准入管控流于形式', True, CLR['accent']),
    '9辆公务车/员工车未履行报备手续，免费停车2-3月 → 涉及9,813元',
    '"三人行"优惠演变成1带多：67人中老客户仅9人，新客户58人 → 审批依据仅为团购名单',
    '流程走过场 → 审核松一松 → 费就漏一漏',
], top=emu(1.4), left=emu(2.2), width=emu(9.5), height=emu(1.8), font_size=14, line_spacing=1.3)
add_label(s, '问题4', emu(3.5), emu(0.8), emu(1.2), emu(0.4), bg_color=CLR['accent'])
add_text_box(s, [
    ('计费规则管理不到位，收费标准备案变更不及时', True, CLR['accent']),
    '两条规则设置依据查不到：缴费后30分钟不离场重计费、新能源车多15分钟宽限期',
    '收费标准2017年调整后，至今未做备案变更',
    '合规不是"我觉得合理就行"，是"文件怎么说就怎么做"',
], top=emu(3.5), left=emu(2.2), width=emu(9.5), height=emu(1.8), font_size=14, line_spacing=1.3)

# ═══════════════════════════════════
# SLIDE 15: OPERATIONS
# ═══════════════════════════════════
print('Slide 15: Operations')
s = add_slide(1); add_decor(s)
add_title_box(s, '经营管理问题（问题5-8）', top=emu(0.6), font_size=26)
issues = [
    ('问题5', '房间使用超出协议范围', '今站购物中心17间房，10间未签协议。管理人员觉得"配套用房不用签"——审计不认"默认"'),
    ('问题6', '经营区域超出协议范围', '5家商户+外摆区域超出租赁范围，存在无偿占用风险'),
    ('问题7', '房间调换未审批', '商户B247(40㎡)→B213(66㎡)，调换无审批，多出26㎡未计价 → 少收11,656元'),
    ('问题8', '退租交接存在疏漏', 'B220退租后今站仍持钥匙，退租返还单未注明钥匙返还。丢了东西算谁的？'),
]
y = emu(1.4)
for label, title, desc in issues:
    add_label(s, label, y, emu(0.8), emu(1.1), emu(0.35), bg_color=CLR['amber'], font_size=12)
    add_text_box(s, [(title, True, CLR['amber']), desc],
                 top=y, left=emu(2.1), width=emu(9.6), height=emu(1.0), font_size=13, line_spacing=1.2)
    y += emu(1.2)
add_text_box(s, [('共同病根：合同管理粗放 + 审批形同虚设 + 交接缺少闭环', True, CLR['accent'])],
             top=emu(5.8), font_size=14)

# ═══════════════════════════════════
# SLIDE 16: PROPERTY + ON-SITE
# ═══════════════════════════════════
print('Slide 16: Property + On-site')
s = add_slide(1); add_decor(s)
add_title_box(s, '物业用房与现场管理问题（问题9-13）', top=emu(0.6), font_size=24)
add_label(s, '物业用房', emu(1.3), emu(0.8), emu(1.8), emu(0.4), bg_color=CLR['blue'])
add_text_box(s, [
    ('问题9：台账不全', True, CLR['dark']),
    '5个房间仅1个入台账，门牌缺失或标错——资产连台账都不全，拿什么管？',
    ('问题10：未经审批被占用', True, CLR['dark']),
    '内部改造加装隔断未经审批，7间房被外单位使用无借用函件。你的房子被外人用了，连一张纸的凭据都拿不出来',
], top=emu(1.8), left=emu(0.8), width=emu(5.5), height=emu(3.0), font_size=13, line_spacing=1.2)
add_label(s, '现场管理', emu(1.3), emu(6.8), emu(1.8), emu(0.4), bg_color=CLR['green'])
add_text_box(s, [
    ('问题12：记录逻辑失真', True, CLR['dark']),
    '同一标段同一天两份维保记录，处理事项截然不同',
    '维保记录提前签字、不同内容用相同总结语',
    '台账空调多2台、消火栓少18个——设备实物与账面差20个',
    '⚠ 出安全事故查设备台账发现是假的——这不是审计问题，是法律问题',
    ('问题13：年度考评缺位', True, CLR['dark']),
    '仅1,200元，但反映管理链条末端失控：上级标准到基层变成一纸空文',
], top=emu(1.8), left=emu(6.8), width=emu(5.5), height=emu(3.5), font_size=13, line_spacing=1.2)

# ═══════════════════════════════════
# SLIDE 17: ISSUE PANORAMA
# ═══════════════════════════════════
print('Slide 17: Issue Panorama')
s = add_slide(1); add_decor(s)
add_title_box(s, '13项问题全景图', top=emu(0.6))
nums = [
    ('13项', '问题', CLR['accent']),
    ('4大', '管理领域', CLR['amber']),
    ('22,671元', '涉及金额', CLR['blue']),
    ('5大', '核心法规被违反', CLR['accent']),
]
cw = emu(2.8)
for i, (num, label, clr) in enumerate(nums):
    x = emu(0.7) + i*(cw+emu(0.3))
    tb = add_text_box(s, [num], top=emu(1.5), left=x, width=cw, height=emu(0.6), font_size=32, color=clr)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = add_text_box(s, [label], top=emu(2.1), left=x, width=cw, height=emu(0.4), font_size=14, color=CLR['gray'])
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('共同特征链：内控体系"形似而神不似"', True, CLR['title']),
    '制度"有"但执行"空" → 台账"有"但信息"缺" → 流程"有"但留痕"无" → 审批"过"但审查"松"',
    '这就是15号文要着力解决的，也是46号令要追责的那种状态',
    ('你不是没制度，你是有了制度没当回事——这才是最可怕的', True, CLR['accent']),
], top=emu(2.8), font_size=16, line_spacing=1.35)

# ═══════════════════════════════════
# SLIDE 18: PART 3 DIVIDER
# ═══════════════════════════════════
print('Slide 18: Part 3 Divider')
s = add_slide(2); add_decor(s)
add_title_box(s, '第三部分', top=emu(2.0), font_size=44, color=CLR['green'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
add_title_box(s, '总结与建议', top=emu(3.2), font_size=28, color=CLR['dark'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0), bold=False)
tb = add_text_box(s, ['希望各位把这几分钟听到的话带回去、落在行动上'],
                  top=emu(4.2), font_size=16, color=CLR['gray'], left=emu(1.5), width=SW-emu(3.0))
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# SLIDE 19: RECOMMENDATIONS
# ═══════════════════════════════════
print('Slide 19: Recommendations')
s = add_slide(1); add_decor(s)
add_title_box(s, '综合管理建议', top=emu(0.6))
recs = [
    ('🅿 停车场管理', '底数清查 + 备案变更 + 岗位分设 + 权限隔离 + 台账 + 核对', CLR['accent']),
    ('🏢 租赁及物业用房', '权属排查 + 补签协议 + 规范台账 + 退租闭环', CLR['amber']),
    ('👁 现场运营', '常态化培训 + 巡查台账 + 统一标准 + 杜绝失真', CLR['blue']),
    ('🔄 常态化核查', '供应商考评闭环 + 自查机制 + 问题整改闭环', CLR['green']),
]
y = emu(1.4)
for title, desc, clr in recs:
    add_label(s, title, y, emu(0.8), emu(3.5), emu(0.45), bg_color=clr, font_size=14)
    add_text_box(s, [desc], top=y+emu(0.05), left=emu(4.5), width=emu(7.5), height=emu(0.5), font_size=15)
    y += emu(1.0)
add_text_box(s, ['每一条都是从问题里"反推"出来的对策，不是拍脑袋写的，是用教训换的'],
             top=emu(5.5), font_size=14, color=CLR['gray'])

# ═══════════════════════════════════
# SLIDE 20: SIX-STEP
# ═══════════════════════════════════
print('Slide 20: Six-Step')
s = add_slide(1); add_decor(s)
add_title_box(s, '制度落地六步法', top=emu(0.6))
steps = [
    ('建制度', '有制度吗？', CLR['blue']), ('明职责', '知道谁负责吗？', CLR['blue']),
    ('讲培训', '制度有人教过吗？', CLR['amber']), ('建台账', '做了有记录吗？', CLR['amber']),
    ('强检查', '检查不是在演戏吗？', CLR['accent']), ('严问责', '出问题了有人扛吗？', CLR['accent']),
]
sw_step, gap = emu(1.8), emu(0.15)
for i, (action, question, clr) in enumerate(steps):
    x = emu(0.7) + i*(sw_step+gap)
    add_label(s, action, emu(1.5), x, sw_step, emu(0.55), bg_color=clr, font_size=16)
    tb = add_text_box(s, [f'← {question}'], top=emu(2.2), left=x-emu(0.2), width=sw_step+emu(0.4), height=emu(0.4), font_size=11, color=CLR['accent'])
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if i < 5:
        tb = add_text_box(s, ['→'], top=emu(1.6), left=x+sw_step, width=emu(0.3), height=emu(0.4), font_size=18, color=CLR['gray'])
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_text_box(s, [
    ('核心结论', True, CLR['title']),
    '打通最后一公里的关键，不在于制度有多少，而在于执行有多实',
    '六个反问句 = 请各位回去做一次部门内的自我体检',
], top=emu(3.0), font_size=16, line_spacing=1.4)

# ═══════════════════════════════════
# SLIDE 21: KEY INSIGHTS
# ═══════════════════════════════════
print('Slide 21: Key Insights')
s = add_slide(1); add_decor(s)
add_title_box(s, '关键启示', top=emu(0.6))
add_text_box(s, [
    ('三句话', True, CLR['title']),
    '❶ 内控不是墙上的制度，是做出来的',
    '❷ 台账不是纸上的记录，是出了事能查到的',
    '❸ 审批不是流程走过场，是责任落实到人',
    ('四位一体框架', True, CLR['title']),
    '1号文 → 用系统减少"人"的随意性',
    '2号文 → 用框架消除"管"的盲区',
    '15号文 → 用执行填补"制"与"行"之间的鸿沟',
    '46号令 → 用问责守住该守的底线',
], top=emu(1.5), font_size=16, line_spacing=1.35)

# ═══════════════════════════════════
# SLIDE 22: SELF-CHECK
# ═══════════════════════════════════
print('Slide 22: Self-Check')
s = add_slide(1); add_decor(s)
add_title_box(s, '合规自查要点（课后作业）', top=emu(0.6))
checks = [
    ('对照1号文', '财务系统打通了吗？预警自动了吗？', CLR['blue']),
    ('对照2号文', '穿透到底了吗？资金可追溯了吗？', CLR['blue']),
    ('对照15号文', '内控制度全覆盖了吗？整改闭环了吗？', CLR['amber']),
    ('对照46号令', '红线全员知晓了吗？容错边界明确了吗？', CLR['accent']),
]
y = emu(1.5)
for label, question, clr in checks:
    add_label(s, label, y, emu(0.8), emu(2.2), emu(0.4), bg_color=clr, font_size=13)
    add_text_box(s, [question], top=y+emu(0.05), left=emu(3.2), width=emu(8.5), height=emu(0.4), font_size=16)
    y += emu(0.8)
add_text_box(s, [
    ('核心原则八字诀', True, CLR['title']),
    '制度有 · 执行实  |  台账有 · 信息全  |  流程有 · 留痕清  |  审批严 · 追责准',
], top=emu(4.6), font_size=18, line_spacing=1.3)

# ═══════════════════════════════════
# SLIDE 23: TRANSITION
# ═══════════════════════════════════
print('Slide 23: Transition')
s = add_slide(1); add_decor(s)
add_title_box(s, '结语', top=emu(0.6))
add_text_box(s, [
    ('送大家三句话', True, CLR['title']),
    ('🔍 每一次审计都是一次体检', True, CLR['dark']),
    '体检不是为难你，是帮你发现隐患',
    ('🪞 每一个问题都是一面镜子', True, CLR['dark']),
    '照出来的不是你一个人的问题，是整个管理体系的问题',
    ('📈 每一项整改都是一次升级', True, CLR['dark']),
    '把问题改到位了，你的管理水平就上了一个台阶',
], top=emu(1.5), font_size=18, line_spacing=1.4)

# ═══════════════════════════════════
# SLIDE 24: THANK YOU
# ═══════════════════════════════════
print('Slide 24: Thank You')
s = add_slide(1); add_decor(s)
add_title_box(s, '谢谢大家', top=emu(2.0), font_size=48, color=CLR['title'],
              align=PP_ALIGN.CENTER, width=SW-emu(2.0), left=emu(1.0))
tb = add_text_box(s, [
    '"提升全员审计风险意识"——核心词不是"审计"，是"全员"',
    '审计团队来查是一年一次，但风险发生是每时每刻',
    '真正防住风险的人，是你们每一个业务岗位上的每一个人',
    '',
    '四川融策会计师事务所  |  2026年6月',
    '欢迎会后交流提问',
], top=emu(3.5), font_size=16, color=CLR['dark'], left=emu(1.5), width=SW-emu(3.0))
for p in tb.text_frame.paragraphs:
    p.alignment = PP_ALIGN.CENTER

# ── Save ──
print(f'\nSaving to {OUT}...')
prs.save(OUT)
print(f'Saved to {OUT}')

# Copy to desktop
try:
    shutil.copy2(OUT, DESK)
    print(f'Copied to {DESK}')
except Exception as e:
    print(f'Copy to desktop failed: {e}')
    print(f'File is at: {OUT}')
