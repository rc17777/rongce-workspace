import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os, zipfile, io

tmpl = r'C:\Users\scrccpa\Desktop\物资管理制度培训-2025.1.20.pptx'
outdir = r'D:\openclaw-workspace\scripts\ppt_assets'
os.makedirs(outdir, exist_ok=True)

prs = Presentation(tmpl)

# Extract images from specific slides
with zipfile.ZipFile(tmpl, 'r') as z:
    # Map slide->image relationships
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                print(f'Slide {si+1}: {shape.name} pos=({shape.left/914400:.1f},{shape.top/914400:.1f}) size=({shape.width/914400:.1f}x{shape.height/914400:.1f})')
                # Get image blob
                try:
                    image = shape.image
                    ext = image.content_type.split('/')[-1]
                    fname = f'slide{si+1}_{shape.name}.{ext}'
                    with open(os.path.join(outdir, fname), 'wb') as f:
                        f.write(image.blob)
                    print(f'  -> saved {fname} ({len(image.blob)} bytes)')
                except Exception as e:
                    print(f'  -> ERROR: {e}')

print('\nDone extracting images.')
