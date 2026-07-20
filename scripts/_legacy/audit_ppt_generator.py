#!/usr/bin/env python3
"""
审计PPT生成器 — 审计报告一键转PPT

两条路径:
  A. 本地模式: python-pptx 直接生成（立即可用）
  B. ppt-master模式: 生成prompt文件（需要ppt-master环境）

用法:
  python scripts/audit_ppt_generator.py report.md                          # 本地生成
  python scripts/audit_ppt_generator.py report.md --mode prompt            # 生成ppt-master prompt
  python scripts/audit_ppt_generator.py report.md --style professional     # 指定风格
  python scripts/audit_ppt_generator.py --demo                             # 演示模式
"""
from __future__ import annotations

import sys
import re
import textwrap
import argparse
import json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "output" / "ppt"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ============================================================
#  融策品牌配色
# ============================================================

BRAND = {
    "primary": RGBColor(0x1A, 0x3A, 0x6E),     # 深蓝
    "accent": RGBColor(0xC0, 0x39, 0x2B),       # 融策红
    "dark": RGBColor(0x2C, 0x3E, 0x50),          # 深灰蓝
    "light": RGBColor(0xF5, 0xF7, 0xFA),         # 浅灰蓝
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "text": RGBColor(0x33, 0x33, 0x33),
    "muted": RGBColor(0x7F, 0x8C, 0x8D),
}

PRESET_STYLES = {
    "professional": {  # 政府审计风格
        "name": "专业稳健",
        "cover_color": RGBColor(0x1A, 0x3A, 0x6E),
        "section_color": RGBColor(0x2C, 0x3E, 0x50),
        "accent_color": RGBColor(0xC0, 0x39, 0x2B),
        "bg_color": RGBColor(0xF5, 0xF7, 0xFA),
    },
    "modern": {  # 现代简洁风格
        "name": "现代简洁",
        "cover_color": RGBColor(0x0D, 0x47, 0xA1),
        "section_color": RGBColor(0x15, 0x65, 0xC0),
        "accent_color": RGBColor(0xFF, 0x6F, 0x00),
        "bg_color": RGBColor(0xFA, 0xFA, 0xFA),
    },
    "elegant": {  # 沉稳大气风格
        "name": "沉稳大气",
        "cover_color": RGBColor(0x1B, 0x1B, 0x1B),
        "section_color": RGBColor(0x37, 0x47, 0x4F),
        "accent_color": RGBColor(0xB7, 0x1C, 0x1C),
        "bg_color": RGBColor(0xEF, 0xEB, 0xE9),
    },
}


# ============================================================
#  内容解析
# ============================================================

def parse_markdown(text: str) -> list[dict]:
    """解析Markdown为幻灯片结构"""
    slides = []
    lines = text.strip().split('\n')

    current_slide = None
    current_content = []

    h1_pattern = re.compile(r'^#\s+(.+)')
    h2_pattern = re.compile(r'^##\s+(.+)')

    for line in lines:
        stripped = line.strip()
        if not stripped:
            if current_content:
                current_content.append('')
            continue

        h1 = h1_pattern.match(stripped)
        h2 = h2_pattern.match(stripped)

        if h1 or h2:
            if current_slide:
                current_slide["body"] = '\n'.join(current_content).strip()
                slides.append(current_slide)

            title = h1.group(1) if h1 else h2.group(1)
            level = 1 if h1 else 2
            current_slide = {"title": title, "level": level, "body": "", "bullets": []}
            current_content = []
        else:
            if current_slide is None:
                current_slide = {"title": "审计报告", "level": 1, "body": "", "bullets": []}
                current_content = []

            # 检测列表项
            if re.match(r'^[\d\-\*•>]\s', stripped) or stripped.startswith('> '):
                clean = re.sub(r'^[\d\-\*•>]+\s*', '', stripped)
                current_slide.setdefault("bullets", []).append(clean)
            else:
                current_content.append(stripped)

    if current_slide:
        current_slide["body"] = '\n'.join(current_content).strip()
        slides.append(current_slide)

    return slides


# ============================================================
#  PPT 生成引擎
# ============================================================

class AuditPPTGenerator:
    """审计报告PPT生成器"""

    def __init__(self, style: str = "professional", title: str = "审计报告"):
        self.style = PRESET_STYLES.get(style, PRESET_STYLES["professional"])
        self.title = title
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate(self, slides: list[dict], output_path: str | Path | None = None) -> Path:
        """生成PPT"""
        if output_path is None:
            output_path = OUTPUT_DIR / f"审计报告_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
        output_path = Path(output_path)

        # 封面
        self._cover_slide(slides[0] if slides else {"title": self.title})

        # 目录
        sections = [s for s in slides if s["level"] == 1]
        if sections:
            self._toc_slide(sections)

        # 内容页
        for slide in slides:
            if slide["level"] == 1:
                self._section_slide(slide)
            else:
                self._content_slide(slide)

        # 结尾页
        self._ending_slide()

        self.prs.save(str(output_path))
        return output_path

    def _add_bg(self, slide, color=None):
        """设置幻灯片背景色"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color or self.style["bg_color"]

    def _add_rect(self, slide, left, top, width, height, color, alpha=None):
        """添加矩形色块"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_textbox(self, slide, left, top, width, height, text, font_size=18,
                     color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or self.style["cover_color"]
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_bullet_list(self, slide, left, top, width, height, items,
                         font_size=16, color=None, spacing=Pt(8)):
        """添加项目符号列表"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color or BRAND["text"]
            p.font.name = '微软雅黑'
            p.space_after = spacing
            p.level = 0

        return txBox

    def _add_footer(self, slide, text=None):
        """添加页脚"""
        self._add_textbox(
            slide, Cm(1), Cm(17.5), Cm(20), Cm(1),
            text or "四川融策会计师事务所 | 四川融策工程咨询公司",
            font_size=9, color=BRAND["muted"], alignment=PP_ALIGN.CENTER
        )

    def _cover_slide(self, slide_data):
        """封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._add_bg(slide, self.style["cover_color"])

        # 装饰线
        self._add_rect(slide, Cm(3), Cm(7), Cm(1.5), Pt(4), self.style["accent_color"])

        # 报告类型标签
        self._add_textbox(slide, Cm(3), Cm(6), Cm(18), Cm(1.5),
                          "审 计 报 告",
                          font_size=16, color=RGBColor(0xFF, 0xFF, 0xFF),
                          bold=False, font_name='微软雅黑')

        # 标题
        title_text = slide_data.get("title", self.title)
        self._add_textbox(slide, Cm(3), Cm(8), Cm(22), Cm(3),
                          title_text,
                          font_size=36, color=RGBColor(0xFF, 0xFF, 0xFF),
                          bold=True, font_name='微软雅黑')

        # 公司信息
        self._add_textbox(slide, Cm(3), Cm(13), Cm(20), Cm(2),
                          "四川融策会计师事务所\n四川融策工程咨询公司",
                          font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD),
                          font_name='微软雅黑')

        # 日期
        self._add_textbox(slide, Cm(3), Cm(15.5), Cm(10), Cm(1),
                          datetime.now().strftime('%Y年%m月'),
                          font_size=12, color=RGBColor(0x99, 0xAA, 0xBB),
                          font_name='微软雅黑')

    def _toc_slide(self, sections):
        """目录页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_rect(slide, Cm(2), Cm(2.5), Pt(6), Cm(1.5), self.style["accent_color"])
        self._add_textbox(slide, Cm(2.8), Cm(2.5), Cm(15), Cm(2),
                          "目  录", font_size=32, bold=True,
                          color=self.style["cover_color"])

        items = [s["title"] for s in sections[:6]]
        y_start = 5.5
        for i, item in enumerate(items):
            num_color = self.style["accent_color"] if i < 3 else BRAND["muted"]
            self._add_textbox(slide, Cm(4), Cm(y_start + i * 1.8), Cm(2), Cm(1.5),
                              f"0{i+1}", font_size=36, bold=True,
                              color=num_color)
            self._add_textbox(slide, Cm(6.5), Cm(y_start + i * 1.8 + 0.2), Cm(18), Cm(1.5),
                              item, font_size=20, color=BRAND["text"])

        self._add_footer(slide)

    def _section_slide(self, slide_data):
        """章节标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, self.style["section_color"])

        self._add_rect(slide, Cm(3), Cm(7.5), Cm(3), Pt(4), self.style["accent_color"])
        self._add_textbox(slide, Cm(3), Cm(8.2), Cm(22), Cm(3),
                          slide_data["title"],
                          font_size=32, color=RGBColor(0xFF, 0xFF, 0xFF),
                          bold=True)

        self._add_footer(slide)

    def _content_slide(self, slide_data):
        """内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        # 顶栏
        self._add_rect(slide, Cm(0), Cm(0), Cm(33.867), Cm(1.8), self.style["cover_color"])
        self._add_textbox(slide, Cm(2), Cm(0.3), Cm(25), Cm(1.2),
                          slide_data["title"],
                          font_size=18, color=BRAND["white"], bold=True)

        # 页脚
        self._add_footer(slide)

        bullets = slide_data.get("bullets", [])
        body = slide_data.get("body", "")

        if bullets:
            self._add_bullet_list(slide, Cm(2), Cm(2.5), Cm(28), Cm(14),
                                  bullets, font_size=16)
        elif body:
            # 分段显示
            paragraphs = [p for p in body.split('\n') if p.strip()]
            if paragraphs:
                txBox = slide.shapes.add_textbox(Cm(2), Cm(2.5), Cm(28), Cm(14))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, para in enumerate(paragraphs[:20]):  # 限制20段
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = para[:200]
                    p.font.size = Pt(14)
                    p.font.color.rgb = BRAND["text"]
                    p.font.name = '微软雅黑'
                    p.space_after = Pt(6)

    def _ending_slide(self, slide_data=None):
        """结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, self.style["cover_color"])

        self._add_textbox(slide, Cm(3), Cm(6), Cm(20), Cm(2),
                          "感谢聆听", font_size=40, bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF),
                          alignment=PP_ALIGN.CENTER)

        self._add_textbox(slide, Cm(3), Cm(9), Cm(20), Cm(3),
                          "四川融策会计师事务所\n四川融策工程咨询公司\n\n地址：四川省成都市\n电话：028-XXXXXXXX",
                          font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                          alignment=PP_ALIGN.CENTER)


# ============================================================
#  ppt-master Prompt 生成器
# ============================================================

def generate_ppt_master_prompt(md_path: str, style: str = "professional",
                               output_pptx: str = "审计报告.pptx") -> str:
    """生成ppt-master兼容的prompt"""
    md_content = Path(md_path).read_text(encoding="utf-8")

    style_desc = {
        "professional": "政府机关专业风格：深蓝色主色调，白色背景，布局干净严谨，字体端正，适合财政部门汇报",
        "modern": "现代简洁风格：蓝色系搭配橙色点缀，大面积留白，信息密度克制",
        "elegant": "沉稳大气风格：深色封面，暖灰背景，红色点缀，适合高层汇报",
    }

    prompt = f"""请根据以下Markdown内容生成一个专业的审计汇报PPT。

【风格要求】
{style_desc.get(style, style_desc['professional'])}

【内容要求】
1. 封面页：包含报告标题、委托单位、编制单位（四川融策会计师事务所/四川融策工程咨询公司）
2. 目录页：列出所有一级章节
3. 每个一级标题作为章节分隔页
4. 内容页：每个二级标题一页，正文要点清晰
5. 结尾页：感谢聆听+公司信息

【特别要求】
- 使用16:9宽屏比例
- 文字大小不低于14pt
- 图表区域自动生成占位图（注明"请插入相关图表"）
- 页脚标注"四川融策会计师事务所"

【Markdown内容】
{md_content}

请现在开始生成，输出.pptx文件。
"""
    return prompt


# ============================================================
#  CLI
# ============================================================

DEMO_MD = """# 绩效评价报告

## 核心发现摘要
> XX市养老服务体系建设专项资金绩效评价总体得分为82.5分，等级为"良"
> 资金拨付及时率达95%，但项目产出质量存在短板
> 建议优化资金分配机制，强化过程监控

# 一、项目概述

## 1.1 项目背景
XX市2025年度养老服务体系建设专项资金主要用于社区养老服务中心建设、居家养老服务补贴、养老机构运营补助三大方向。截至2025年12月31日，专项资金已拨付3,050万元，拨付率95.3%。

## 1.2 评价目标与方法
本次绩效评价采用"资料审查+实地走访+问卷调查+数据分析"相结合的方式。覆盖全市6个区县、24个社区养老服务中心、12家养老机构。综合得分82.5分（满分100分），等级为"良"。

# 二、资金管理

## 2.1 预算执行
项目立项依据充分，绩效目标设置较为合理。但存在部分区县资金分配未充分考虑老年人口分布的实际需求，资源配置与需求存在错位。

## 2.2 资金使用合规性
- XX区民政局将养老专项资金85万元用于单位日常公用经费，不符合专款专用规定
- 3个社区养老服务中心建设项目未按合同约定及时拨付工程款，拖欠施工单位款项合计120万元
- 部分居家养老服务记录不完整，抽查发现约15%的服务工单缺少服务对象签字确认

# 三、问题发现

## 3.1 主要问题
- 资金被挤占挪用：涉及金额85万元
- 项目进度滞后：2个社区养老服务中心建设延期6个月以上
- 资产闲置：养老机构床位利用率仅62%
- 农村地区养老服务覆盖不足

## 3.2 风险等级分布
- 高风险：5项（涉及资金占用、项目延期）
- 中风险：6项（管理规范性问题）
- 低风险：4项（制度完善性）

# 四、改进建议

## 4.1 优化资金分配
- 建立"因素法+项目法"相结合的资金分配模型
- 将老年人口数量、老龄化程度、经济发展水平等因素纳入分配权重

## 4.2 强化过程监控
- 建立专项资金使用"红黄绿"预警机制
- 推广居家养老服务信息化管理平台
- 加强民政与卫健部门协作，推进医养结合
"""


def main():
    parser = argparse.ArgumentParser(description="审计PPT生成器")
    parser.add_argument("input", nargs="?", help="输入Markdown文件路径")
    parser.add_argument("--mode", choices=["local", "prompt"], default="local",
                        help="生成模式: local(本地PPT) / prompt(ppt-master提示词)")
    parser.add_argument("--style", choices=list(PRESET_STYLES.keys()), default="professional",
                        help="PPT风格")
    parser.add_argument("-o", "--output", help="输出文件路径")
    parser.add_argument("--demo", action="store_true", help="演示模式")
    args = parser.parse_args()

    if args.demo:
        # 生成临时md文件
        demo_path = OUTPUT_DIR / f"demo_audit_{datetime.now().strftime('%H%M')}.md"
        demo_path.write_text(DEMO_MD, encoding="utf-8")
        args.input = str(demo_path)
        print(f"[DEMO] 演示素材: {demo_path.name}")

    if not args.input:
        parser.print_help()
        return

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"[ERROR] 文件不存在: {args.input}")
        return

    md_content = input_path.read_text(encoding="utf-8")

    if args.mode == "prompt":
        prompt = generate_ppt_master_prompt(str(input_path), args.style)
        prompt_path = args.output or OUTPUT_DIR / f"ppt_master_prompt_{datetime.now().strftime('%Y%m%d_%H%M')}.txt"
        Path(prompt_path).parent.mkdir(parents=True, exist_ok=True)
        Path(prompt_path).write_text(prompt, encoding="utf-8")
        print(f"[OK] ppt-master提示词已生成: {prompt_path}")
        print(f"\n预览:\n{prompt[:500]}...")
        return

    # 本地生成模式
    slides = parse_markdown(md_content)
    print(f"[*] 解析到 {len(slides)} 页幻灯片")

    if not slides:
        print("[ERROR] 未能解析出幻灯片内容")
        return

    title = slides[0].get("title", "审计报告")
    generator = AuditPPTGenerator(style=args.style, title=title)
    output_path = generator.generate(slides, args.output)
    print(f"[OK] PPT已生成: {output_path}")
    print(f"     共 {len(generator.prs.slides)} 页幻灯片")
    print(f"     风格: {PRESET_STYLES[args.style]['name']}")


if __name__ == "__main__":
    main()
