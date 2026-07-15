#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
融策会计师事务所 · 业务能力展示 PPTX v2
政府风格 + 顶级商务风格
配色: 深蓝#0A1F3F / 青绿#1A5C6E / 铜金#C5955C / 暖灰#F5F2EC
字体: 标题微软雅黑 / 正文宋体
"""

import sys, io
sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ─── Colors ───
DEEP_BLUE = RGBColor(0x0A, 0x1F, 0x3F)
TEAL = RGBColor(0x1A, 0x5C, 0x6E)
GOLD = RGBColor(0xC5, 0x95, 0x5C)
WARM_GRAY = RGBColor(0xF5, 0xF2, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
GOLD_LIGHT = RGBColor(0xE8, 0xCF, 0xA0)

SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)
MARGIN_L = Cm(2.8)
CONTENT_W = Cm(28.267)

FONT_T = '微软雅黑'
FONT_B = '宋体'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, x, y, w, h, fc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background()
    if fc: s.fill.solid(); s.fill.fore_color.rgb = fc
    else: s.fill.background()
    return s

def tb(slide, x, y, w, h, text, fn=FONT_B, fs=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, ls=1.2):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tx.word_wrap = True
    tf = tx.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.name = fn; p.font.size = Pt(fs); p.font.bold = bold
    p.font.color.rgb = color; p.alignment = align
    p.space_after = Pt(0); p.space_before = Pt(0)
    pPr = p._p.get_or_add_pPr()
    ln = pPr.makeelement(qn('a:lnSpc'), {})
    sp = ln.makeelement(qn('a:spcPct'), {'val': str(int(ls * 100000))})
    ln.append(sp); pPr.append(ln)
    return tx

def ml_tb(slide, x, y, w, h, lines, fn=FONT_B, fs=14, color=DARK_TEXT, align=PP_ALIGN.LEFT, ls=1.5):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tx.word_wrap = True; tf = tx.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.name = fn; p.font.size = Pt(fs)
        p.font.color.rgb = color; p.alignment = align
        p.space_after = Pt(4); p.space_before = Pt(0)
        pPr = p._p.get_or_add_pPr()
        ln = pPr.makeelement(qn('a:lnSpc'), {})
        sp = ln.makeelement(qn('a:spcPct'), {'val': str(int(ls * 100000))})
        ln.append(sp); pPr.append(ln)
    return tx

def top_bar(slide): rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(0.15), GOLD)
def bottom_bar(slide): rect(slide, Cm(0), Cm(18.6), SLIDE_W, Cm(0.45), DEEP_BLUE)
def gold_line(slide, x, y, w): rect(slide, x, y, w, Cm(0.08), GOLD)
def pn(slide, n): tb(slide, Cm(30.5), Cm(18.65), Cm(3), Cm(0.4), f'{n}/10', fs=9, color=WHITE, align=PP_ALIGN.RIGHT)

# ═══ P1: Cover ═══
s = blank(); set_bg(s, DEEP_BLUE)
rect(s, Cm(0), Cm(0), SLIDE_W, Cm(0.4), GOLD)
rect(s, Cm(0), Cm(18.65), SLIDE_W, Cm(0.4), GOLD)
rect(s, Cm(2.5), Cm(4), Cm(0.06), Cm(11), GOLD)
tb(s, Cm(4), Cm(4.5), Cm(26), Cm(3), '四川融策会计师事务所', FONT_T, 40, True, WHITE, PP_ALIGN.LEFT)
tb(s, Cm(4), Cm(7.5), Cm(26), Cm(2), 'SICHUAN RONGCE CPA', FONT_B, 14, False, GOLD_LIGHT, PP_ALIGN.LEFT)
tb(s, Cm(4), Cm(10), Cm(26), Cm(2), '政府审计  ·  绩效评价  ·  工程咨询  ·  财政评审', FONT_B, 16, False, WHITE, PP_ALIGN.LEFT)
gold_line(s, Cm(4), Cm(12.5), Cm(8))
tb(s, Cm(4), Cm(13.2), Cm(26), Cm(1.5), '专业 · 诚信 · 高效  |  值得信赖的审计服务伙伴', FONT_B, 12, False, GOLD_LIGHT, PP_ALIGN.LEFT)
tb(s, Cm(4), Cm(14.5), Cm(10), Cm(1.5), '2026', FONT_T, 20, True, GOLD, PP_ALIGN.LEFT)

# ═══ P2: Company Overview ═══
s = blank(); set_bg(s, WHITE); top_bar(s); bottom_bar(s); pn(s, 2)
tb(s, MARGIN_L, Cm(1.2), CONTENT_W, Cm(1.5), '公司概览', FONT_T, 28, True, DEEP_BLUE)
gold_line(s, MARGIN_L, Cm(2.8), Cm(5))
tb(s, MARGIN_L, Cm(3.5), CONTENT_W, Cm(3), '四川融策会计师事务所主营政府审计业务，涵盖绩效评价、资产清查、专项债申报、监督检查等财政职能；四川融策工程咨询公司主营预算编制、财政评审、全过程工程咨询与工程结算。', FONT_B, 14, False, DARK_TEXT, PP_ALIGN.LEFT, 1.6)
for i, (n, l, d) in enumerate([('300+','服务客户','覆盖政府与企事业单位'),('50+','专业团队','注册会计师+工程师双资质'),('800+','服务项目','12年行业经验积累')]):
    cx = Cm(2.8) + Cm(9.4) * i
    rect(s, cx, Cm(7.5), Cm(8.2), Cm(5), WARM_GRAY)
    tb(s, cx, Cm(8), Cm(8.2), Cm(2.5), n, FONT_T, 36, True, GOLD, PP_ALIGN.CENTER)
    tb(s, cx, Cm(10.5), Cm(8.2), Cm(0.8), l, FONT_B, 14, True, DEEP_BLUE, PP_ALIGN.CENTER)
    tb(s, cx, Cm(11.3), Cm(8.2), Cm(0.8), d, FONT_B, 10, False, GRAY_TEXT, PP_ALIGN.CENTER)
tb(s, MARGIN_L, Cm(13.5), CONTENT_W, Cm(1), '双业务线协同 · 会计师事务所 + 工程咨询公司', FONT_B, 11, False, GRAY_TEXT, PP_ALIGN.CENTER)

# ═══ P3: Core Business ═══
s = blank(); set_bg(s, WHITE); top_bar(s); bottom_bar(s); pn(s, 3)
tb(s, MARGIN_L, Cm(1.2), CONTENT_W, Cm(1.5), '核心业务板块', FONT_T, 28, True, DEEP_BLUE)
gold_line(s, MARGIN_L, Cm(2.8), Cm(5))
biz = [('经济责任审计','任中·离任·自然资源'),('专项资金审计','社保·营养餐等'),('预算执行审计','部门预算·财政收支'),('招投标审计','串标围标检测'),('预算绩效管理','目标·监控·评价'),('工程全过程咨询','预算·评审·结算'),('财政评审与结算','投资评审·决算')]
for i, (nm, dc) in enumerate(biz):
    cx = MARGIN_L + Cm(9.4) * (i if i < 4 else i - 4)
    cy = Cm(4) if i < 4 else Cm(8.5)
    rect(s, cx, cy, Cm(8.2), Cm(3.8), WARM_GRAY)
    rect(s, cx, cy, Cm(0.2), Cm(3.8), DEEP_BLUE)
    tb(s, cx + Cm(0.8), cy + Cm(0.5), Cm(7), Cm(1.2), nm, FONT_T, 15, True, DEEP_BLUE)
    tb(s, cx + Cm(0.8), cy + Cm(2), Cm(7), Cm(1.2), dc, FONT_B, 11, False, GRAY_TEXT)
tb(s, MARGIN_L, Cm(13.5), CONTENT_W, Cm(1), '审计 + 工程咨询 · 双轮驱动  |  12大业务线全覆盖', FONT_B, 11, False, GRAY_TEXT, PP_ALIGN.CENTER)

# ═══ P4: Audit Services ═══
s = blank(); set_bg(s, WHITE); top_bar(s); bottom_bar(s); pn(s, 4)
tb(s, MARGIN_L, Cm(1.2), CONTENT_W, Cm(1.5), '政府审计业务线', FONT_T, 28, True, DEEP_BLUE)
gold_line(s, MARGIN_L, Cm(2.8), Cm(5))
audit = [('绩效评价','事前评估→事中监控→事后评价\n覆盖财政支出、项目支出、政策评价全链条'),('资产清查','固定资产盘点·往来款清理\n专项资金清理·账实核对\n全流程资产管理工作'),('专项债申报','项目规划·方案编制\n评审对接·资金使用监管\n全过程专项债咨询服务'),('监督检查','财政监督检查·会计信息质量检查\n内部控制评价·专项资金检查\n助力规范财政管理')]
for i, (nm, dc) in enumerate(audit):
    cx = MARGIN_L + Cm(7.3) * i
    rect(s, cx, Cm(4), Cm(6.5), Cm(9), WARM_GRAY)
    rect(s, cx, Cm(4), Cm(6.5), Cm(0.15), DEEP_BLUE)
    tb(s, cx + Cm(0.5), Cm(4.8), Cm(5.5), Cm(1.2), nm, FONT_T, 16, True, DEEP_BLUE, PP_ALIGN.CENTER)
    gold_line(s, cx + Cm(1.5), Cm(6.3), Cm(3.5))
    tb(s, cx + Cm(0.5), Cm(7), Cm(5.5), Cm(5.5), dc, FONT_B, 11, False, DARK_TEXT, PP_ALIGN.CENTER, 1.5)

# ═══ P5: Cases ═══
s = blank(); set_bg(s, WHITE); top_bar(s); bottom_bar(s); pn(s, 5)
tb(s, MARGIN_L, Cm(1.2), CONTENT_W, Cm(1.5), '典型案例', FONT_T, 28, True, DEEP_BLUE)
gold_line(s, MARGIN_L, Cm(2.8), Cm(5))
cases = [('某市财政局','绩效评价项目','50+项目·资金3亿元',['完成财政支出绩效评价50余项','涉及教育、社保、农业等领域','提出整改建议200余条']),('某交通局','资产清查项目','2.5亿资产·20+项目',['全面清查交通系统固定资产','盘活闲置资产5000余万元','建立资产管理制度体系']),('某县审计局','经济责任审计','15+项目·整改8000万',['完成经责审计15项','发现违规资金8000余万元','推动出台3项管理制度'])]
for i, (cl, tg, st, dt) in enumerate(cases):
    cx = MARGIN_L + Cm(9.8) * i
    rect(s, cx, Cm(4), Cm(8.8), Cm(8), WARM_GRAY)
    tb(s, cx + Cm(0.6), Cm(4.5), Cm(7.6), Cm(1), cl, FONT_T, 16, True, DEEP_BLUE)
    rect(s, cx + Cm(0.6), Cm(5.6), Cm(3.5), Cm(0.7), GOLD)
    tb(s, cx + Cm(0.6), Cm(5.6), Cm(3.5), Cm(0.7), tg, FONT_B, 9, True, DEEP_BLUE, PP_ALIGN.CENTER)
    tb(s, cx + Cm(0.6), Cm(6.6), Cm(7.6), Cm(0.8), st, FONT_B, 11, True, TEAL)
    ml_tb(s, cx + Cm(0.6), Cm(7.6), Cm(7.6), Cm(4), dt, FONT_B, 10, DARK_TEXT, PP_ALIGN.LEFT, 1.5)

# ═══ P6: Key Data ═══
s = blank(); set_bg(s, DEEP_BLUE)
rect(s, Cm(0), Cm(0), SLIDE_W, Cm(0.4), GOLD)
rect(s, Cm(0), Cm(18.65), SLIDE_W, Cm(0.4), GOLD)
tb(s, MARGIN_L, Cm(1.5), CONTENT_W, Cm(1.5), '关键数据', FONT_T, 28, True, WHITE)
gold_line(s, MARGIN_L, Cm(3.2), Cm(5))
tb(s, MARGIN_L, Cm(4.5), CONTENT_W, Cm(4), '800+', FONT_T, 72, True, GOLD, PP_ALIGN.CENTER)
tb(s, MARGIN_L, Cm(8.5), CONTENT_W, Cm(1), '累计服务项目', FONT_B, 16, False, GOLD_LIGHT, PP_ALIGN.CENTER)
gold_line(s, Cm(10), Cm(10.5), Cm(14))
for i, (n, l) in enumerate([('92%','客户续约率'),('50人+','专业团队'),('30个','覆盖市县')]):
    cx = MARGIN_L + Cm(9.5) * i
    tb(s, cx, Cm(11), Cm(8), Cm(2), n, FONT_T, 36, True, WHITE, PP_ALIGN.CENTER)
    tb(s, cx, Cm(13.5), Cm(8), Cm(1), l, FONT_B, 13, False, GOLD_LIGHT, PP_ALIGN.CENTER)
tb(s, Cm(30.5), Cm(18.2), Cm(3), Cm(0.4), '6/10', FONT_B, 9, False, GOLD_LIGHT, PP_ALIGN.RIGHT)

# ═══ P7: Quality Assurance ═══
s = blank(); set_bg(s, WHITE); top_bar(s); bottom_bar(s); pn(s, 7)
tb(s, MARGIN_L, Cm(1.2), CONTENT_W, Cm(1.5), '质量保障体系', FONT_T, 28, True, DEEP_BLUE)
gold_line(s, MARGIN_L, Cm(2.8), Cm(5))
qa = [('三级复核制','项目组→部门→质控部\n三级交叉复核，层层把关\n确保报告零差错'),('AI辅助复核','15维度智能检查系统\n自动检测错别字/金额/日期\n法规引用一致性校验'),('法规数据库','13,000+审计法规实时更新\n覆盖审计法/会计法/预算法等\n智能化法规匹配检索'),('客户满意度','98%客户好评率\n12年持续服务经验\n定期回访持续改进')]
for i, (nm, dc) in enumerate(qa):
    cx = MARGIN_L + Cm(7.3) * i
    rect(s, cx, Cm(4), Cm(6.5), Cm(6.5), WARM_GRAY)
    rect(s, cx + Cm(0.5), Cm(4.5), Cm(0.15), Cm(2.5), GOLD)
    tb(s, cx + Cm(1), Cm(4.5), Cm(5), Cm(1.2), nm, FONT_T, 14, True, DEEP_BLUE)
    tb(s, cx + Cm(0.8), Cm(6.2), Cm(5), Cm(4), dc, FONT_B, 10, False, DARK_TEXT, PP_ALIGN.LEFT, 1.5)
tb(s, MARGIN_L, Cm(11.5), CONTENT_W, Cm(1), '从项目立项到报告出具的全流程质量管控，确保每一份报告经得起检验', FONT_B, 11, False, GRAY_TEXT, PP_ALIGN.CENTER)

# ═══ P8: Mission ═══
s = blank(); set_bg(s, DEEP_BLUE)
rect(s, Cm(0), Cm(0), SLIDE_W, Cm(0.4), GOLD)
rect(s, Cm(0), Cm(18.65), SLIDE_W, Cm(0.4), GOLD)
tb(s, MARGIN_L, Cm(5), CONTENT_W, Cm(3), '以专业审计服务\n助力政府治理现代化', FONT_T, 36, True, WHITE, PP_ALIGN.CENTER, 1.4)
gold_line(s, Cm(14), Cm(9.5), Cm(6))
for i, v in enumerate(['专业立身','诚信为本','高效服务']):
    cx = MARGIN_L + Cm(9) * i
    rect(s, cx + Cm(2.5), Cm(10.5), Cm(3.5), Cm(0.7), GOLD)
    tb(s, cx + Cm(2.5), Cm(10.5), Cm(3.5), Cm(0.7), v, FONT_T, 14, True, DEEP_BLUE, PP_ALIGN.CENTER)
tb(s, Cm(30.5), Cm(18.2), Cm(3), Cm(0.4), '8/10', FONT_B, 9, False, GOLD_LIGHT, PP_ALIGN.RIGHT)

# ═══ P9: Engineering Consulting ═══
s = blank(); set_bg(s, WHITE); top_bar(s); bottom_bar(s); pn(s, 9)
tb(s, MARGIN_L, Cm(1.2), CONTENT_W, Cm(1.5), '全过程工程咨询', FONT_T, 28, True, DEEP_BLUE)
gold_line(s, MARGIN_L, Cm(2.8), Cm(5))
tb(s, MARGIN_L, Cm(3.5), CONTENT_W, Cm(2), '从预算编制到工程结算，覆盖项目建设全生命周期，以专业能力为政府投资项目保驾护航。', FONT_B, 14, False, DARK_TEXT, PP_ALIGN.LEFT, 1.5)
stages = [('预算编制','投资估算\n设计概算\n施工图预算'),('财政评审','预算评审\n招标控制价\n评审报告'),('全过程咨询','跟踪审计\n进度款审核\n变更管理'),('工程结算','结算审核\n竣工决算\n决算报告')]
for i, (nm, dc) in enumerate(stages):
    cx = MARGIN_L + Cm(7.3) * i
    rect(s, cx, Cm(6.5), Cm(6.5), Cm(5.5), WARM_GRAY)
    rect(s, cx, Cm(6.5), Cm(6.5), Cm(0.15), DEEP_BLUE)
    tb(s, cx + Cm(0.5), Cm(7), Cm(5.5), Cm(1.2), nm, FONT_T, 16, True, DEEP_BLUE, PP_ALIGN.CENTER)
    gold_line(s, cx + Cm(1.5), Cm(8.5), Cm(3.5))
    tb(s, cx + Cm(0.5), Cm(9.2), Cm(5.5), Cm(2.5), dc, FONT_B, 11, False, DARK_TEXT, PP_ALIGN.CENTER, 1.5)
# 关键数据
tb(s, MARGIN_L, Cm(13), Cm(9), Cm(1.5), '500+ 完成项目', FONT_T, 18, True, DEEP_BLUE, PP_ALIGN.CENTER)
tb(s, Cm(12.5), Cm(13), Cm(9), Cm(1.5), '10亿+ 评审金额', FONT_T, 18, True, DEEP_BLUE, PP_ALIGN.CENTER)
tb(s, Cm(22), Cm(13), Cm(9), Cm(1.5), '1.2亿 节约资金', FONT_T, 18, True, DEEP_BLUE, PP_ALIGN.CENTER)

# ═══ P10: Contact ═══
s = blank(); set_bg(s, DEEP_BLUE)
rect(s, Cm(0), Cm(0), SLIDE_W, Cm(0.4), GOLD)
rect(s, Cm(0), Cm(18.65), SLIDE_W, Cm(0.4), GOLD)
tb(s, MARGIN_L, Cm(3), CONTENT_W, Cm(2), '联系我们', FONT_T, 36, True, WHITE, PP_ALIGN.CENTER)
gold_line(s, Cm(14), Cm(5.5), Cm(6))
tb(s, MARGIN_L, Cm(6.5), CONTENT_W, Cm(1.5), '四川融策 · 与您同行', FONT_B, 16, False, GOLD_LIGHT, PP_ALIGN.CENTER)
tb(s, MARGIN_L, Cm(8.5), CONTENT_W, Cm(2), '无论您是政府部门还是企事业单位，融策都将以专业、诚信、高效的服务，为您提供最优质的审计与工程咨询解决方案。', FONT_B, 13, False, GOLD_LIGHT, PP_ALIGN.CENTER, 1.5)
info = [('公司地址','四川省成都市高新区'),('联系电话','028-XXXXXXX'),('电子邮箱','rongce@rongcecpa.com'),('官方网站','www.rongcecpa.com')]
for i, (l, v) in enumerate(info):
    cx = MARGIN_L + Cm(7) * i
    tb(s, cx, Cm(12), Cm(6.5), Cm(1), l, FONT_B, 12, True, GOLD, PP_ALIGN.CENTER)
    tb(s, cx, Cm(13), Cm(6.5), Cm(1), v, FONT_B, 12, False, WHITE, PP_ALIGN.CENTER)
tb(s, Cm(30.5), Cm(18.2), Cm(3), Cm(0.4), '10/10', FONT_B, 9, False, GOLD_LIGHT, PP_ALIGN.RIGHT)

# ═══ Save ═══
out = r'C:\Users\scrccpa\.openclaw\workspace\output\rongce-capability-20260715\rongce-capability-v2.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'OK: {out}')
print(f'Slides: {len(prs.slides)}')