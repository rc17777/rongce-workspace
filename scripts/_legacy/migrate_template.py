"""
模板迁移脚本 v3：将审计风险培训内容套入物资管理培训模板
"""

import os, copy, io, time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = r"C:\Users\scrccpa\Desktop\物资管理制度培训-2025.1.20.pptx"
CONTENT  = r"C:\Users\scrccpa\Desktop\轨道培训\四川轨道公司审计风险培训-v2.pptx"

tpl = Presentation(TEMPLATE)
src = Presentation(CONTENT)
print("Template: %d slides, Content: %d slides" % (len(tpl.slides), len(src.slides)))

# ── 查找layout ──
def find_layout(prs, name):
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name == name:
            return i
    return 0

# ── 创建输出，清空 ──
out = Presentation(TEMPLATE)
while len(out.slides) > 0:
    rId = out.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    out.part.drop_rel(rId)
    out.slides._sldIdLst.remove(out.slides._sldIdLst[0])

L_TITLE   = find_layout(out, '标题幻灯片')
L_SECTION = find_layout(out, '节标题')
L_CONTENT = find_layout(out, '标题和内容')
print("Layouts - title:%d section:%d content:%d" % (L_TITLE, L_SECTION, L_CONTENT))

# ── 颜色常量 ──
DARK  = RGBColor(0x33, 0x33, 0x33)
GRAY  = RGBColor(0x66, 0x66, 0x66)
BLUE  = RGBColor(0x0B, 0x75, 0xB4)

# ── 工具函数 ──
def add_textbox(slide, text, x, y, w, h, font_name='微软雅黑', font_size=28,
                bold=True, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(6)
    return txBox

def add_bg_images(slide):
    """从模板第3页复制背景装饰图（底部条+Logo）"""
    proto = tpl.slides[2]
    for shape in proto.shapes:
        if shape.shape_type == 13:  # Picture
            try:
                image_blob = shape.image.blob
                pic = slide.shapes.add_picture(
                    io.BytesIO(image_blob), shape.left, shape.top,
                    shape.width, shape.height)
                if shape.rotation != 0:
                    pic.rotation = shape.rotation
            except Exception as e:
                print("  bg warn: %s" % e)

def extract_texts(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)
    return texts

# ── 页类型映射 ──
PAGE_MAP = {
    1: 'title', 2: 'agenda', 3: 'section',
    4: 'content', 5: 'content', 6: 'content', 7: 'content',
    8: 'content', 9: 'content', 10: 'content',
    11: 'section', 12: 'content', 13: 'content', 14: 'content',
    15: 'content', 16: 'content', 17: 'content',
    18: 'section', 19: 'content', 20: 'content',
    21: 'content', 22: 'content', 23: 'content', 24: 'end',
}

# ── 逐页重建 ──
for src_idx in range(len(src.slides)):
    page_num = src_idx + 1
    ptype = PAGE_MAP.get(page_num, 'content')
    src_slide = src.slides[src_idx]
    
    # 选layout
    if ptype in ('title', 'end'):
        layout_idx = L_TITLE
    elif ptype == 'section':
        layout_idx = L_SECTION
    else:
        layout_idx = L_CONTENT
    new_slide = out.slides.add_slide(out.slide_layouts[layout_idx])
    
    # 背景图（封面和结束页不加）
    if ptype not in ('title', 'end'):
        add_bg_images(new_slide)
    
    # 内容
    texts = extract_texts(src_slide)
    if not texts:
        continue
    
    title = texts[0]
    body = texts[1:]
    
    if ptype == 'title':
        add_textbox(new_slide, "提升公司全员\n审计风险意识",
                   Cm(4), Cm(5), Cm(26), Cm(5),
                   '楷体', 40, True, BLUE, PP_ALIGN.CENTER)
        add_textbox(new_slide, "四川轨道公司专题培训",
                   Cm(4), Cm(11), Cm(26), Cm(2),
                   '楷体', 22, False, GRAY, PP_ALIGN.CENTER)
        add_textbox(new_slide, "四川融策会计师事务所    2026年6月",
                   Cm(4), Cm(13.5), Cm(26), Cm(1.5),
                   '楷体', 16, False, GRAY, PP_ALIGN.CENTER)
    
    elif ptype == 'section':
        add_textbox(new_slide, title,
                   Cm(3), Cm(6), Cm(28), Cm(5),
                   '方正小标宋简体', 48, True, BLUE, PP_ALIGN.CENTER)
        if body:
            add_textbox(new_slide, body[0],
                       Cm(4), Cm(11.5), Cm(26), Cm(3),
                       '微软雅黑', 20, False, GRAY, PP_ALIGN.CENTER)
    
    elif ptype == 'agenda':
        add_textbox(new_slide, title,
                   Cm(2), Cm(1.5), Cm(30), Cm(2),
                   '微软雅黑', 28, True, DARK, PP_ALIGN.CENTER)
        body_text = '\n'.join(body[:12])
        add_textbox(new_slide, body_text,
                   Cm(3), Cm(4), Cm(28), Cm(13),
                   '微软雅黑', 16, False, DARK, PP_ALIGN.LEFT)
    
    elif ptype == 'end':
        add_textbox(new_slide, "谢谢！",
                   Cm(4), Cm(6), Cm(26), Cm(4),
                   '楷体', 48, True, BLUE, PP_ALIGN.CENTER)
        add_textbox(new_slide, "欢迎会后交流提问",
                   Cm(4), Cm(11), Cm(26), Cm(2),
                   '楷体', 20, False, GRAY, PP_ALIGN.CENTER)
    
    else:  # content
        add_textbox(new_slide, title,
                   Cm(1.5), Cm(1.3), Cm(31), Cm(2.5),
                   '微软雅黑', 26, True, DARK, PP_ALIGN.CENTER)
        if body:
            body_text = '\n'.join(body)
            add_textbox(new_slide, body_text,
                       Cm(2), Cm(4), Cm(30), Cm(13.5),
                       '微软雅黑', 14, False, DARK, PP_ALIGN.LEFT)

# ── 保存 ──
ts = time.strftime("%H%M%S")
out_file = r"D:\openclaw-workspace\output\sc-v3-" + ts + ".pptx"
out.save(out_file)
print("\nDone: " + out_file)
print("Slides: " + str(len(out.slides)))
