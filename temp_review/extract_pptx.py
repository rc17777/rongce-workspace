# -*- coding: utf-8 -*-
"""提取验收报告pptx文本"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

path = r'C:\Users\scrccpa\Desktop\新建文件夹\1.深度行四川站合同审核资料\2、信通院\制造业数字化转型促进中心 深度行 （四川站） 验收报告-信通院.pptx'
prs = Presentation(path)
out = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)
        if shape.has_table:
            for row in shape.table.rows:
                texts.append(' | '.join(c.text.strip() for c in row.cells))
    if texts:
        out.append(f'--- slide {i} ---')
        out.extend(texts)
text = '\n'.join(out)
with open(r'C:\Users\scrccpa\.openclaw\workspace\temp_review\acceptance.txt', 'w', encoding='utf-8') as f:
    f.write(text)
print(f'共{len(prs.slides)}页，提取{len(text)}字符')
