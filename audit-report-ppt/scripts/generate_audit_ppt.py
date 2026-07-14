#!/usr/bin/env python3
"""
审计报告PPT生成器
自动调用python-pptx生成汇报级PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
sys.stdout.reconfigure(encoding='utf-8')

class AuditReportPPTGenerator:
    """审计报告PPT生成器"""
    
    COLORS = {
        'primary': RGBColor(0, 51, 102),
        'secondary': RGBColor(0, 102, 204),
        'accent': RGBColor(255, 102, 0),
        'warning': RGBColor(204, 0, 0),
        'success': RGBColor(0, 128, 0),
        'text': RGBColor(51, 51, 51),
        'text_light': RGBColor(102, 102, 102),
        'background': RGBColor(245, 247, 250),
    }
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def add_text(self, slide, text, x, y, w, h, size=16, bold=False, color=None, align=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        if color: p.font.color.rgb = color
        if align: p.alignment = align
        return tf
    
    def add_cover_slide(self, title, subtitle, presenter=""):
        """T01 封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_text(slide, title, 0.5, 2.5, 12, 1.5, 44, True, self.COLORS['primary'], PP_ALIGN.CENTER)
        self.add_text(slide, subtitle, 0.5, 4.2, 12, 0.8, 20, False, self.COLORS['text_light'], PP_ALIGN.CENTER)
        if presenter:
            self.add_text(slide, presenter, 0.5, 5.5, 12, 0.8, 16, False, self.COLORS['text'], PP_ALIGN.CENTER)
        return slide
    
    def add_title_slide(self, title):
        """T03 章节过渡页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_text(slide, title, 2, 3, 9, 1.5, 36, True, self.COLORS['primary'], PP_ALIGN.CENTER)
        return slide
    
    def add_data_slide(self, title, data_points):
        """T10 数字突出页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_text(slide, title, 0.5, 0.3, 12, 0.8, 32, True, self.COLORS['primary'])
        for i, (val, lbl, color) in enumerate(data_points[:3]):
            x = 1 + i * 4
            self.add_text(slide, val, x, 2, 3.5, 1.2, 56, True, color)
            self.add_text(slide, lbl, x, 3.3, 3.5, 0.8, 16, False, self.COLORS['text'])
        return slide
    
    def add_comparison_slide(self, title, before_items, after_items):
        """T06 对比页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_text(slide, title, 0.5, 0.3, 12, 0.8, 32, True, self.COLORS['primary'])
        tf_l = self.add_text(slide, "整改前", 0.5, 1.5, 5.5, 0.6, 24, True, self.COLORS['warning'])
        for item in before_items:
            p = tf_l.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.COLORS['text']
        tf_r = self.add_text(slide, "整改后", 7, 1.5, 5.5, 0.6, 24, True, self.COLORS['success'])
        for item in after_items:
            p = tf_r.add_paragraph()
            p.text = f"✅ {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.COLORS['text']
        return slide
    
    def add_problem_slide(self, title, problems):
        """T14 问题清单页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_text(slide, title, 0.5, 0.3, 12, 0.8, 32, True, self.COLORS['primary'])
        tf = self.add_text(slide, "", 0.5, 1.5, 12, 5.5, 14, False, self.COLORS['text'])
        for i, prob in enumerate(problems[:6], 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = f"{i}. {prob}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.COLORS['text']
        return slide
    
    def add_rectification_slide(self, title, rectifications):
        """整改建议页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_text(slide, title, 0.5, 0.3, 12, 0.8, 32, True, self.COLORS['primary'])
        tf = self.add_text(slide, "", 0.5, 1.5, 12, 5.5, 14, False, self.COLORS['text'])
        for rec in rectifications[:4]:
            priority = rec.get('priority', 'P2')
            color = {'P0': self.COLORS['warning'], 'P1': self.COLORS['accent']}.get(priority, self.COLORS['text'])
            p = tf.paragraphs[0] if priority == rectifications[0].get('priority', 'P2') else tf.add_paragraph()
            p.text = f"[{priority}] {rec['problem']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = color
            p2 = tf.add_paragraph()
            p2.text = f"  → {rec['solution']}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = self.COLORS['text']
            p3 = tf.add_paragraph()
            p3.text = f"  {rec.get('unit', '')} | {rec.get('deadline', '')}"
            p3.font.size = Pt(12)
            p3.font.color.rgb = self.COLORS['text_light']
        return slide
    
    def save(self, path):
        self.prs.save(path)
        print(f"PPT saved: {path}")
        return path

# 示例
if __name__ == '__main__':
    gen = AuditReportPPTGenerator()
    
    gen.add_cover_slide(
        title="2025年度某区财政收支审计报告",
        subtitle="某区审计局 | 审计期间: 2025年1-12月",
        presenter="XXX | 2026年6月3日"
    )
    
    gen.add_title_slide("一、审计发现")
    
    gen.add_data_slide("三个核心数字", [
        ("856.3万", "违规金额", gen.COLORS['warning']),
        ("45%", "制度执行率", gen.COLORS['accent']),
        ("68%", "问题整改率", gen.COLORS['success']),
    ])
    
    gen.add_comparison_slide(
        "整改前后对比",
        ["制度执行率: 45%(2024年)", "预算调整率: 28%(2024年)", "违规发放: XXX万元(2024年)"],
        ["制度执行率: 82%(2025年, +37%)", "预算调整率: 12%(2025年, -16%)", "违规发放: 已清退XXX万元(2025年)"]
    )
    
    gen.add_problem_slide("主要问题清单", [
        "违规发放津补贴XXX万元",
        "预算编制不精细，调整率高达28%",
        "固定资产管理不规范，账实不符",
        "政府采购程序不完整",
    ])
    
    gen.add_rectification_slide("审计建议", [
        {'priority': 'P0', 'problem': '违规发放津补贴', 'solution': '清退违规金额XXX万元', 'unit': '某局', 'deadline': '2026年3月31日前'},
        {'priority': 'P1', 'problem': '预算编制不精细', 'solution': '建立预算编制模板', 'unit': '财政局', 'deadline': '2026年6月30日前'},
    ])
    
    gen.save(r"C:\Users\Admin\审计报告_演示.pptx")
